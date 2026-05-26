from contextlib import asynccontextmanager
from fastapi import FastAPI, File, UploadFile, HTTPException, Request, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, JSONResponse
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded
from loguru import logger
from typing import List, Optional
import io
import os
import re
import math
import time
import uuid
import asyncio
import secrets
import zipfile
import tempfile
import functools
import httpx
import sys
from urllib.parse import quote as url_quote
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor

import database as db
import payment as pay
import cache as _cache
from img2pdf_improved import (
    do_img2pdf_single, do_imgs2pdf_multi, validate_params as img2pdf_validate, is_image_bytes,
)

# ─── Config ──────────────────────────────────────────────────────────────────

BOT_TOKEN     = os.environ.get("BOT_TOKEN", "")
APP_URL       = os.environ.get("APP_URL", "https://nematovbegz0d.github.io/studentsTools/EduBot.html?v=8")
# Telegram webhook secret token — Telegram'ga setWebhook qilganda biz uzatamiz, har
# kelgan so'rovda X-Telegram-Bot-Api-Secret-Token header bilan tasdiqlanadi.
# Agar bo'sh bo'lsa, BOT_TOKEN dan hash hosil qilamiz (deterministik, hech kim taxmin qilolmaydi).
WEBHOOK_SECRET = os.environ.get("WEBHOOK_SECRET", "")
if not WEBHOOK_SECRET and BOT_TOKEN:
    import hashlib as _hl
    WEBHOOK_SECRET = _hl.sha256(f"webhook:{BOT_TOKEN}".encode()).hexdigest()[:48]

MAX_FILE_MB = 20
MAX_FILE_BYTES = MAX_FILE_MB * 1024 * 1024
_start_time = time.time()

FONT_REGULAR = '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'
FONT_BOLD    = '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf'

# ─── Sentry (optional) ───────────────────────────────────────────────────────
# Set SENTRY_DSN env var to enable error monitoring. Without it, code is a no-op.
_SENTRY_DSN = os.environ.get("SENTRY_DSN", "")
if _SENTRY_DSN:
    try:
        import sentry_sdk
        from sentry_sdk.integrations.fastapi import FastApiIntegration
        from sentry_sdk.integrations.starlette import StarletteIntegration
        sentry_sdk.init(
            dsn=_SENTRY_DSN,
            traces_sample_rate=float(os.environ.get("SENTRY_TRACE_RATE", "0.1")),
            profiles_sample_rate=0.0,
            environment=os.environ.get("RAILWAY_ENVIRONMENT", "production"),
            release=os.environ.get("RAILWAY_GIT_COMMIT_SHA", "")[:7] or "dev",
            integrations=[StarletteIntegration(), FastApiIntegration()],
            send_default_pii=False,
        )
    except Exception as _sentry_err:
        # Don't fail boot if sentry is misconfigured
        pass

# ─── Logging ─────────────────────────────────────────────────────────────────

logger.remove()
logger.add(sys.stdout, format="{time:HH:mm:ss} | {level:<7} | {message}", level="INFO")

# ─── Rate limiter ─────────────────────────────────────────────────────────────

limiter = Limiter(key_func=get_remote_address, default_limits=["60/minute"])

# ─── Thread pool for CPU-intensive conversions ────────────────────────────────
# pdf2docx blocks the event loop — run it in a thread pool

_converter_pool = ThreadPoolExecutor(max_workers=2, thread_name_prefix="pdf2docx")


def _do_pdf2docx(pdf_path: str, docx_path: str) -> None:
    """CPU-bound, blocking — runs in _converter_pool thread."""
    from pdf2docx import Converter
    cv = Converter(pdf_path)
    try:
        cv.convert(docx_path, multi_processing=False, start=0)
    finally:
        cv.close()


def _do_pdf2docx_libreoffice(pdf_path: str, docx_path: str) -> None:
    """LibreOffice writer_pdf_import filter → DOCX. Best quality, requires libreoffice in PATH."""
    import subprocess, shutil
    outdir = os.path.dirname(docx_path)
    lo_home = tempfile.mkdtemp(prefix="lo_")
    env = os.environ.copy()
    env["HOME"] = lo_home
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--norestore",
             "--convert-to", "docx", "--infilter=writer_pdf_import",
             "--outdir", outdir, pdf_path],
            env=env, capture_output=True, timeout=60,
        )
    finally:
        shutil.rmtree(lo_home, ignore_errors=True)
    if result.returncode != 0:
        stderr = (result.stderr or b"").decode(errors="replace")[:300]
        raise RuntimeError(f"libreoffice exit {result.returncode}: {stderr}")
    expected = os.path.join(outdir, os.path.splitext(os.path.basename(pdf_path))[0] + ".docx")
    if not os.path.exists(expected) or os.path.getsize(expected) < 100:
        raise RuntimeError("LibreOffice DOCX chiqarilmadi yoki bo'sh")
    if os.path.abspath(expected) != os.path.abspath(docx_path):
        shutil.move(expected, docx_path)


def _do_pdf2docx_pymupdf(pdf_path: str, docx_path: str) -> None:
    """PyMuPDF block-level extraction → python-docx. Preserves headings, bold, italic, colors."""
    import fitz
    from docx import Document
    from docx.shared import Pt, RGBColor

    pdf = fitz.open(pdf_path)
    doc_out = Document()

    # Median font size — faqat birinchi 3 sahifadan namuna (OOM oldini olish)
    sample_sizes = []
    for page in pdf[:3]:
        for block in page.get_text("dict")["blocks"]:
            if block.get("type") == 0:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        sz = span.get("size", 0)
                        if sz > 0 and span.get("text", "").strip():
                            sample_sizes.append(sz)
                            if len(sample_sizes) >= 500:
                                break
                    if len(sample_sizes) >= 500:
                        break
            if len(sample_sizes) >= 500:
                break
        if len(sample_sizes) >= 500:
            break
    body_size = sorted(sample_sizes)[len(sample_sizes) // 2] if sample_sizes else 11.0

    for page_num, page in enumerate(pdf):
        if page_num > 0:
            doc_out.add_paragraph()

        for block in page.get_text("dict")["blocks"]:
            if block.get("type") != 0:
                continue
            all_spans = [
                s for ln in block.get("lines", [])
                for s in ln.get("spans", [])
                if s.get("text", "").strip()
            ]
            if not all_spans:
                continue

            avg_size = sum(s["size"] for s in all_spans) / len(all_spans)
            block_text = " ".join(s["text"].strip() for s in all_spans)

            if avg_size >= body_size * 1.6:
                doc_out.add_heading(block_text, level=1)
            elif avg_size >= body_size * 1.35:
                doc_out.add_heading(block_text, level=2)
            elif avg_size >= body_size * 1.15:
                doc_out.add_heading(block_text, level=3)
            else:
                p = doc_out.add_paragraph()
                for span in all_spans:
                    text = span.get("text", "")
                    if not text:
                        continue
                    run = p.add_run(text)
                    flags = span.get("flags", 0)
                    run.bold = bool(flags & 16)
                    run.italic = bool(flags & 2)
                    sz = span.get("size", body_size)
                    if 4 <= sz <= 72:
                        run.font.size = Pt(round(sz))
                    color = span.get("color", 0)
                    if color:
                        r, g, b = (color >> 16) & 0xFF, (color >> 8) & 0xFF, color & 0xFF
                        if (r, g, b) != (0, 0, 0):
                            try:
                                run.font.color.rgb = RGBColor(r, g, b)
                            except Exception:
                                pass

    pdf.close()
    doc_out.save(docx_path)


def _do_pdf2docx_ocr(pdf_path: str, docx_path: str) -> None:
    """Scanned PDF fallback: Tesseract OCR → python-docx (Tesseract already in Docker)."""
    import fitz
    import pytesseract
    from PIL import Image as PILImage
    from docx import Document

    pdf = fitz.open(pdf_path)
    doc_out = Document()
    doc_out.add_heading("OCR natijasi", level=1)

    for page_num, page in enumerate(pdf):
        if page_num > 0:
            doc_out.add_page_break()

        mat = fitz.Matrix(2.0, 2.0)  # 2x zoom → OCR aniqligini oshiradi
        pix = page.get_pixmap(matrix=mat)
        img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)

        try:
            text = pytesseract.image_to_string(img, lang="uzb+rus+eng", timeout=30)
        except Exception:
            text = pytesseract.image_to_string(img, lang="eng", timeout=30)

        for line in (text or "").split("\n"):
            line = line.strip()
            if line:
                doc_out.add_paragraph(line)

    pdf.close()
    doc_out.save(docx_path)


# ✅ YANGILANDI: docling (IBM) — AI bilan jadvallarni aniq taniydi, birlamchi usul
def _do_pdf2docx_docling(pdf_path: str, docx_path: str) -> None:
    from docling.document_converter import DocumentConverter
    converter = DocumentConverter()
    result = converter.convert(pdf_path)
    doc = result.document
    # docling 2.x API
    saved = False
    try:
        docx_bytes = doc.export_to_docx()
        with open(docx_path, "wb") as f:
            f.write(docx_bytes)
        saved = True
    except AttributeError:
        pass
    if not saved:
        try:
            from docling.backend.docx_backend import DocxDocumentBackend
            with open(docx_path, "wb") as f:
                DocxDocumentBackend().write(doc, f)
            saved = True
        except Exception:
            pass
    if not saved:
        raise RuntimeError("docling DOCX saqlash usuli topilmadi")


def _do_pdf2docx_best(pdf_path: str, docx_path: str, is_scanned: bool) -> str:
    """Tries best available method in order. Returns method name on success."""
    if is_scanned:
        chain = [
            ("ocr-tesseract", lambda: _do_pdf2docx_ocr(pdf_path, docx_path)),
            ("pdf2docx",      lambda: _do_pdf2docx(pdf_path, docx_path)),
        ]
    else:
        # ✅ YANGILANDI: DOCLING_DISABLED=1 bo'lsa docling o'tkazib yuboriladi
        chain = []
        if os.environ.get("DOCLING_DISABLED", "").lower() not in ("1", "true", "yes"):
            chain.append(("docling", lambda: _do_pdf2docx_docling(pdf_path, docx_path)))
        chain.extend([
            ("libreoffice",   lambda: _do_pdf2docx_libreoffice(pdf_path, docx_path)),
            ("pdf2docx",      lambda: _do_pdf2docx(pdf_path, docx_path)),
            ("pymupdf",       lambda: _do_pdf2docx_pymupdf(pdf_path, docx_path)),
        ])

    last_err = None
    for name, fn in chain:
        try:
            fn()
            if os.path.exists(docx_path) and os.path.getsize(docx_path) >= 2000:
                return name
        except Exception as e:
            last_err = e
            logger.warning(f"pdf2docx [{name}] muvaffaqiyatsiz: {type(e).__name__}: {str(e)[:120]}")
        try:
            if os.path.exists(docx_path):
                os.unlink(docx_path)
        except Exception:
            pass

    raise RuntimeError(
        f"Barcha usullar muvaffaqiyatsiz: "
        f"{type(last_err).__name__ if last_err else 'Unknown'}: {str(last_err)[:100]}"
    )


# ─── Rembg session ───────────────────────────────────────────────────────────

_rembg_session = None

def get_rembg_session():
    global _rembg_session
    if _rembg_session is None:
        from rembg import new_session
        model = os.environ.get("REMBG_MODEL", "birefnet-general")
        try:
            _rembg_session = new_session(model)
            logger.info(f"rembg session yaratildi ({model})")
        except Exception:
            _rembg_session = new_session("u2netp")
            logger.info("rembg session yaratildi (u2netp fallback)")
    return _rembg_session

# ─── Lifespan ─────────────────────────────────────────────────────────────────

@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("EduBot Backend ishga tushmoqda...")
    logger.info("ADMIN_TOKEN: set" if ADMIN_TOKEN else "ADMIN_TOKEN: NOT SET — Railway Variables ga qoshing")

    # ── Ma'lumotlar bazasini ishga tushirish ──
    try:
        await db.init_db()
        db_type = "PostgreSQL" if db._USE_PG else "SQLite"
        logger.info(f"{db_type} DB ishga tushdi")
    except Exception as e:
        logger.error(f"DB init xatosi: {e}")

    domain = os.environ.get("RAILWAY_PUBLIC_DOMAIN", "")
    if BOT_TOKEN and domain:
        webhook_url = f"https://{domain}/webhook"
        try:
            async with httpx.AsyncClient(timeout=10) as client:
                # ✅ secret_token Telegram'dan har webhook so'rov bilan X-Telegram-Bot-Api-Secret-Token
                # header sifatida qaytadi — boshqa hech kim webhook'ni soxtalashtirolmaydi
                payload = {"url": webhook_url, "drop_pending_updates": True}
                if WEBHOOK_SECRET:
                    payload["secret_token"] = WEBHOOK_SECRET
                r = await client.post(
                    f"https://api.telegram.org/bot{BOT_TOKEN}/setWebhook",
                    json=payload,
                )
                if r.json().get("ok"):
                    logger.info(f"Webhook: {webhook_url} (secret_token: {'✓' if WEBHOOK_SECRET else 'NO'})")
                else:
                    logger.error(f"Webhook xatosi: {r.json()}")
        except Exception as e:
            logger.error(f"Webhook ulanishda xato: {e}")
    else:
        logger.warning("BOT_TOKEN yoki RAILWAY_PUBLIC_DOMAIN yo'q")

    # ── Rembg modelini background'da yuklash ─────────────────────────
    # Startup ni bloklamasligi uchun background task sifatida ishlatiladi.
    # isnet-general-use modeli katta (170MB) — yuklanishi vaqt olishi mumkin.
    async def _preload_rembg():
        try:
            loop = asyncio.get_event_loop()
            await asyncio.wait_for(
                loop.run_in_executor(None, get_rembg_session),
                timeout=120.0,
            )
            logger.info("rembg model preloaded ✓")
        except asyncio.TimeoutError:
            logger.warning("rembg preload timeout — birinchi so'rovda yuklanadi")
        except Exception as e:
            logger.warning(f"rembg preload muvaffaqiyatsiz — birinchi so'rovda yuklanadi: {e}")

    asyncio.create_task(_preload_rembg())

    yield
    logger.info("EduBot Backend to'xtatilmoqda...")
    global _rembg_session
    _rembg_session = None

# ─── App ─────────────────────────────────────────────────────────────────────

app = FastAPI(title="EduBot Backend", version="2.0.0", lifespan=lifespan)
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)

# ✅ CORS — aniq originlar (CORS_ORIGINS env: vergul bilan ajratilgan ro'yxat)
# Default: GitHub Pages + Telegram Web (telegram.org va t.me)
_default_origins = (
    "https://nematovbegz0d.github.io,"
    "https://web.telegram.org,"
    "https://telegram.org,"
    "https://t.me,"
    # Lokal dev (VS Code Live Server, Vite, CRA, va h.k.)
    "http://127.0.0.1:5500,"
    "http://localhost:5500,"
    "http://127.0.0.1:3000,"
    "http://localhost:3000,"
    "http://127.0.0.1:5173,"
    "http://localhost:5173,"
    "http://127.0.0.1:8000,"
    "http://localhost:8000"
)
_cors_raw = os.environ.get("CORS_ORIGINS", _default_origins)
_cors_origins = [o.strip() for o in _cors_raw.split(",") if o.strip()]

app.add_middleware(
    CORSMiddleware,
    allow_origins=_cors_origins,
    allow_credentials=False,
    allow_methods=["GET", "POST", "DELETE", "OPTIONS"],
    allow_headers=["Content-Type", "Authorization", "X-User-Id"],
    expose_headers=[
        "X-Info", "X-Page-Count", "X-Saved-Percent", "X-Warning",
        "X-Password",  # pdflock yaratgan parol — HTTPS orqali xavfsiz, faqat user ko'radi
        "Content-Disposition",
    ],
)
logger.info(f"CORS allowed origins: {_cors_origins}")

# ─── Helpers ──────────────────────────────────────────────────────────────────

def check_size(data: bytes, endpoint: str = ""):
    if len(data) > MAX_FILE_BYTES:
        mb = round(len(data) / 1024 / 1024, 1)
        raise HTTPException(status_code=413, detail=f"Fayl {mb}MB. Maksimal: {MAX_FILE_MB}MB")

def safe_header(value) -> str:
    """HTTP/1.1 headers must be Latin-1 encodable. Strip emojis and other
    non-Latin-1 chars (✅ • → ⚠️ 📎 ў ҳ ғ қ ...) — these would otherwise raise
    UnicodeEncodeError when Starlette encodes the response."""
    if not isinstance(value, str):
        value = str(value)
    return "".join(ch for ch in value if ord(ch) < 256).strip() or "ok"

async def tg_send(chat_id: int, text: str, reply_markup: Optional[dict] = None):
    payload = {"chat_id": chat_id, "text": text, "parse_mode": "HTML"}
    if reply_markup:
        payload["reply_markup"] = reply_markup
    async with httpx.AsyncClient(timeout=10) as client:
        await client.post(f"https://api.telegram.org/bot{BOT_TOKEN}/sendMessage", json=payload)

def pil_fonts(sizes: list):
    from PIL import ImageFont
    try:
        return [ImageFont.truetype(FONT_BOLD if i == 0 else FONT_REGULAR, s) for i, s in enumerate(sizes)]
    except Exception:
        return [ImageFont.load_default() for _ in sizes]

def text_center_x(draw, text, font, W):
    try:
        bbox = draw.textbbox((0, 0), text, font=font)
        return (W - (bbox[2] - bbox[0])) // 2
    except Exception:
        return W // 4

# ─── Health ───────────────────────────────────────────────────────────────────

@app.get("/")
@app.get("/health")
def health():
    ocr_ok = False
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        ocr_ok = True
    except Exception:
        pass
    return {
        "status": "ok", "version": "2.0.0",
        "uptime": int(time.time() - _start_time),
        "ocr": ocr_ok, "rembg_ready": _rembg_session is not None,
        "max_file_mb": MAX_FILE_MB,
    }

# ─── Telegram Webhook ─────────────────────────────────────────────────────────

@app.post("/webhook")
async def webhook(request: Request):
    if not BOT_TOKEN:
        return {"ok": False}
    # ✅ Telegram secret_token tekshiruv — spoofed webhook so'rovlarni rad etadi
    if WEBHOOK_SECRET:
        got = request.headers.get("X-Telegram-Bot-Api-Secret-Token", "")
        if not secrets.compare_digest(got, WEBHOOK_SECRET):
            logger.warning(f"Webhook secret mismatch from {get_remote_address(request)}")
            raise HTTPException(status_code=403, detail="Forbidden")
    try:
        update = await request.json()
    except Exception:
        return {"ok": False}

    message = update.get("message") or update.get("edited_message")
    if not message:
        return {"ok": True}

    chat_id    = message.get("chat", {}).get("id")
    text       = message.get("text", "")
    first_name = message.get("from", {}).get("first_name", "Do'st")
    logger.info(f"Webhook: chat={chat_id} text={text[:30]!r}")

    # Foydalanuvchini DBga saqlash
    from_user = message.get("from", {})
    try:
        await db.upsert_user(
            user_id=from_user.get("id"),
            username=from_user.get("username"),
            first_name=from_user.get("first_name"),
        )
    except Exception as e:
        logger.warning(f"upsert_user xatosi: {e}")

    if text.startswith("/start"):
        await tg_send(
            chat_id,
            f"Salom, <b>{first_name}</b>! 👋\n\n"
            f"📚 <b>EduBot</b> — talabalar uchun 28+ ta bepul xizmat:\n\n"
            f"• 📄 PDF birlashtirish, bo'lish, himoyalash\n"
            f"• 🖼 Rasm → PDF, siqish, fon olib tashlash\n"
            f"• 🌐 Tarjima, Wikipedia, kitob qidirish\n"
            f"• 📊 Formula, grafik, statistika\n"
            f"• 🎓 Sertifikat, jadval, QR kod\n\n"
            f"Quyidagi tugmani bosib ilovani oching 👇",
            reply_markup={"inline_keyboard": [[{"text": "📱 Ilovani ochish", "web_app": {"url": APP_URL}}]]},
        )
    else:
        await tg_send(
            chat_id,
            "📱 Barcha xizmatlar ilovada mavjud:",
            reply_markup={"inline_keyboard": [[{"text": "📱 Ilovani ochish", "web_app": {"url": APP_URL}}]]},
        )
    return {"ok": True}

# ─── Send file to Telegram ────────────────────────────────────────────────────

@app.post("/api/send-file")
@limiter.limit("20/minute")
async def send_file(
    request: Request,
    user_id: str = Form(...),
    filename: str = Form(...),
    file: UploadFile = File(...),
):
    if not BOT_TOKEN:
        raise HTTPException(status_code=503, detail="Bot sozlanmagan")
    data = await file.read()
    check_size(data, "/api/send-file")
    t0 = time.time()
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.post(
            f"https://api.telegram.org/bot{BOT_TOKEN}/sendDocument",
            data={"chat_id": user_id, "caption": f"📎 {filename}\n\n🤖 EduBot"},
            files={"document": (filename, data, file.content_type or "application/octet-stream")},
        )
    result = r.json()
    if not result.get("ok"):
        raise HTTPException(status_code=500, detail=result.get("description", "Xatolik"))
    logger.info(f"sendDocument: user={user_id} {filename} {time.time()-t0:.1f}s")
    return {"ok": True}

# ─── PDF: Merge ───────────────────────────────────────────────────────────────

_MERGEPDF_MAX_FILES = 30


# ✅ YANGILANDI: pikepdf (qpdf asosida) — buzilgan PDFlarni tuzatadi, linearizatsiya, MPL-2.0
def _do_mergepdf(data_list: list) -> tuple:
    import pikepdf

    out_pdf = pikepdf.new()
    page_count = 0
    for data in data_list:
        if data[:4] != b"%PDF":
            raise ValueError("Faqat PDF fayllar qabul qilinadi")
        try:
            src = pikepdf.open(io.BytesIO(data))
        except pikepdf.PasswordError:
            raise ValueError("Himoyalangan PDF birlashtirish uchun ochib bo'lmadi")
        out_pdf.pages.extend(src.pages)
        page_count += len(src.pages)
        src.close()
    buf = io.BytesIO()
    out_pdf.save(buf, linearize=True)
    out_pdf.close()
    info = f"{len(data_list)} fayl, {page_count} sahifa"
    return buf.getvalue(), info, page_count


@app.post("/api/mergepdf")
@limiter.limit("10/minute")
async def merge_pdf(request: Request, files: List[UploadFile] = File(...)):
    try:
        if len(files) > _MERGEPDF_MAX_FILES:
            raise HTTPException(status_code=400, detail=f"Maksimal {_MERGEPDF_MAX_FILES} fayl")
        data_list = []
        total = 0
        for f in files:
            data = await f.read()
            total += len(data)
            if total > MAX_FILE_BYTES * 3:
                raise HTTPException(status_code=413, detail="Fayllar jami hajmi juda katta")
            data_list.append(data)
        loop = asyncio.get_event_loop()
        try:
            out_bytes, info, _ = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _do_mergepdf, data_list),
                timeout=60.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="PDF birlashtirish vaqti tugadi")
        logger.info(f"mergepdf: {info}")
        return Response(
            content=out_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=merged.pdf",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"mergepdf xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"mergepdf: {type(e).__name__}: {str(e)[:160]}")

# ─── PDF: Split ───────────────────────────────────────────────────────────────

_SPLITPDF_MAX_PAGES = 200


def _parse_page_ranges(pages_str: str, total: int) -> list:
    """Parse "1-5,8,10-15" → [0,1,2,3,4,7,9,10,11,12,13,14] (0-indexed).
    None or empty → all pages."""
    if not pages_str or not pages_str.strip():
        return list(range(total))
    result = set()
    for part in pages_str.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            start, end = int(a), int(b)
            if start > end:
                raise ValueError(f"Noto'g'ri diapazon: {a}-{b} (bosh > oxir)")
            result.update(range(start - 1, end))
        else:
            result.add(int(part) - 1)
    return sorted(x for x in result if 0 <= x < total)


# ✅ YANGILANDI: pikepdf — har bir sahifani mustaqil PDF sifatida ajratadi, linearize=True
def _do_splitpdf(data: bytes, page_indices: list) -> tuple:
    import pikepdf

    if data[:4] != b"%PDF":
        raise ValueError("PDF fayl emas")
    try:
        doc = pikepdf.open(io.BytesIO(data))
    except pikepdf.PasswordError:
        raise ValueError("PDF parol bilan himoyalangan")
    n = len(doc.pages)
    if n == 0:
        doc.close()
        raise ValueError("PDF sahifasiz")
    if len(page_indices) > _SPLITPDF_MAX_PAGES:
        doc.close()
        raise ValueError(f"Maksimal {_SPLITPDF_MAX_PAGES} sahifa ajratish mumkin")

    zf_buf = io.BytesIO()
    with zipfile.ZipFile(zf_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        pad = len(str(max(page_indices) + 1))
        for i in page_indices:
            out = pikepdf.new()
            out.pages.append(doc.pages[i])
            pb = io.BytesIO()
            out.save(pb, linearize=True)
            out.close()
            zf.writestr(f"page_{str(i + 1).zfill(pad)}.pdf", pb.getvalue())
    doc.close()
    return zf_buf.getvalue(), len(page_indices)


@app.post("/api/splitpdf")
@limiter.limit("10/minute")
async def split_pdf(
    request: Request,
    file: UploadFile = File(...),
    pages: Optional[str] = Form(None),
):
    try:
        data = await file.read()
        check_size(data, "/api/splitpdf")
        import pikepdf as _pike_check
        try:
            _pdf_check = _pike_check.open(io.BytesIO(data))
        except _pike_check.PasswordError:
            raise HTTPException(status_code=422, detail="PDF parol bilan himoyalangan.")
        except Exception:
            raise HTTPException(status_code=422, detail="Bu fayl PDF emas.")
        total_n = len(_pdf_check.pages)
        _pdf_check.close()
        try:
            indices = _parse_page_ranges(pages, total_n)
        except Exception:
            raise HTTPException(status_code=400, detail="Sahifa oralig'i noto'g'ri. Format: 1-5,8,10-15")
        if not indices:
            raise HTTPException(status_code=400, detail="Sahifalar ro'yxati bo'sh")
        loop = asyncio.get_event_loop()
        try:
            zip_bytes, count = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _do_splitpdf, data, indices),
                timeout=45.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="PDF ajratish vaqti tugadi")
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e))
        info = f"{count} sahifa → {count} fayl"
        logger.info(f"splitpdf: {info}")
        return Response(
            content=zip_bytes,
            media_type="application/zip",
            headers={
                "Content-Disposition": "attachment; filename=pages.zip",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"splitpdf xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"splitpdf: {type(e).__name__}: {str(e)[:160]}")


# ─── PDF: Page selection ──────────────────────────────────────────────────────

@app.post("/api/pdfpages")
@limiter.limit("10/minute")
async def pdf_select_pages(
    request: Request,
    file: UploadFile = File(...),
    pages: str = Form(...),
):
    """Extract a subset of pages into a new PDF. pages = '1,3,5-10,15'"""
    try:
        data = await file.read()
        check_size(data, "/api/pdfpages")
        if data[:4] != b"%PDF":
            raise HTTPException(status_code=422, detail="Bu fayl PDF emas.")
        # ✅ YANGILANDI: pikepdf — sahifalarni tanlash va linearize qilib saqlash
        import pikepdf as _pike
        try:
            _pdf_tmp = _pike.open(io.BytesIO(data))
        except _pike.PasswordError:
            raise HTTPException(status_code=422, detail="PDF parol bilan himoyalangan.")
        except Exception:
            raise HTTPException(status_code=422, detail="Bu fayl PDF emas.")
        total_n = len(_pdf_tmp.pages)
        _pdf_tmp.close()
        try:
            indices = _parse_page_ranges(pages, total_n)
        except Exception:
            raise HTTPException(status_code=400, detail="Sahifa oralig'i noto'g'ri. Format: 1-5,8,10")
        if not indices:
            raise HTTPException(status_code=400, detail="Sahifalar ro'yxati bo'sh")

        def _extract():
            src = _pike.open(io.BytesIO(data))
            out = _pike.new()
            for i in indices:
                out.pages.append(src.pages[i])
            buf = io.BytesIO()
            out.save(buf, linearize=True)
            src.close()
            return buf.getvalue()

        loop = asyncio.get_event_loop()
        try:
            out_bytes = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _extract),
                timeout=30.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="Sahifa ajratish vaqti tugadi")
        info = f"{len(indices)} sahifa tanlandi (jami {total_n} dan)"
        return Response(
            content=out_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=selected.pdf",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"pdfpages xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"pdfpages: {type(e).__name__}: {str(e)[:160]}")

# ─── PDF: Text extraction ─────────────────────────────────────────────────────

_PDFTEXT_MAX_PAGES = 50

def _do_pdftext(data: bytes) -> str:
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    if doc.is_encrypted:
        doc.close()
        raise ValueError("PDF parol bilan himoyalangan.")
    n = doc.page_count
    if n == 0:
        doc.close()
        return "PDF bo'sh (0 sahifa)."

    render_n = min(n, _PDFTEXT_MAX_PAGES)
    parts = []
    for i in range(render_n):
        # sort=True: reads spans in visual reading order (top→bottom, left→right)
        # Handles multi-column layouts and tables far better than pypdf
        page_text = doc[i].get_text("text", sort=True).strip()
        if page_text:
            if n > 1:
                parts.append(f"── Sahifa {i + 1} ──")
            parts.append(page_text)
    doc.close()

    result = re.sub(r'\n{3,}', '\n\n', "\n\n".join(parts)).strip()

    if not result:
        # pdfplumber fallback — jadvalli PDFlar uchun yaxshiroq
        try:
            import pdfplumber
            pb_parts = []
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                for i, page in enumerate(pdf.pages[:render_n]):
                    text = page.extract_text() or ""
                    tables = page.extract_tables() or []
                    for table in tables:
                        text += "\n" + "\n".join(
                            " | ".join(str(c or "").strip() for c in row)
                            for row in table if any(c for c in row)
                        )
                    if text.strip():
                        if n > 1:
                            pb_parts.append(f"── Sahifa {i + 1} ──")
                        pb_parts.append(text.strip())
            result = re.sub(r'\n{3,}', '\n\n', "\n\n".join(pb_parts)).strip()
        except Exception:
            pass

    if not result:
        return ("❌ Matn topilmadi. Bu skanerlangan PDF bo'lishi mumkin — "
                "OCR xizmatidan foydalaning.")
    if n > _PDFTEXT_MAX_PAGES:
        result += (f"\n\n⚠️ Faqat birinchi {_PDFTEXT_MAX_PAGES} sahifa ko'rsatildi "
                   f"(PDF jami {n} sahifali).")
    return result

@app.post("/api/pdftext")
@limiter.limit("20/minute")
async def pdf_text(request: Request, file: UploadFile = File(...)):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/pdftext")
        if data[:4] != b'%PDF':
            raise HTTPException(status_code=422, detail="Bu fayl PDF emas.")
        loop = asyncio.get_event_loop()
        try:
            result = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _do_pdftext, data),
                timeout=30.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Matn ajratish 30 soniyadan oshdi. Kichikroq PDF tanlang.")
        logger.info(f"pdftext: {len(result)} belgi, {time.time()-t0:.1f}s")
        return JSONResponse({"result": result})
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"pdftext xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"pdftext: {type(e).__name__}: {str(e)[:160]}")

# ─── PDF: Lock ────────────────────────────────────────────────────────────────

# ✅ YANGILANDI: pikepdf AES-256 shifrlash — qpdf asosida, Encryption() API
def _do_pdflock(data: bytes, user_pwd: str, owner_pwd: str,
                allow_print: bool, allow_copy: bool, allow_modify: bool) -> tuple:
    import pikepdf

    try:
        doc = pikepdf.open(io.BytesIO(data))
    except pikepdf.PasswordError:
        raise ValueError("PDF allaqachon parol bilan himoyalangan.")
    n = len(doc.pages)
    if n == 0:
        doc.close()
        raise ValueError("PDF bo'sh (0 sahifa).")
    if n > 200:
        doc.close()
        raise ValueError(f"PDF {n} sahifa. Maksimal 200 sahifa qabul qilinadi.")

    enc = pikepdf.Encryption(
        owner=owner_pwd,
        user=user_pwd,
        aes=True,
        allow=pikepdf.Permissions(
            print_lowres=allow_print,
            print_highres=allow_print,
            extract=allow_copy,
            modify_other=allow_modify,
            modify_annotation=allow_modify,
        ),
    )
    buf = io.BytesIO()
    doc.save(buf, encryption=enc, linearize=True)
    doc.close()

    perms = []
    if allow_print:  perms.append("chop")
    if allow_copy:   perms.append("nusxa")
    if allow_modify: perms.append("tahrir")
    perm_str = ", ".join(perms) if perms else "faqat o'qish"
    info = f"{n} sahifa • AES-256 • {perm_str}"
    return buf.getvalue(), info


@app.post("/api/pdflock")
@limiter.limit("15/minute")
async def lock_pdf(
    request: Request,
    file: UploadFile = File(...),
    password: str = Form(""),
    allow_print: bool = Form(True),
    allow_copy: bool = Form(False),
    allow_modify: bool = Form(False),
):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/pdflock")
        if data[:4] != b'%PDF':
            raise HTTPException(status_code=422, detail="Bu fayl PDF emas.")
        user_pwd  = password.strip()[:64] if password.strip() else secrets.token_urlsafe(9)[:12]
        owner_pwd = secrets.token_urlsafe(16)
        loop = asyncio.get_event_loop()
        try:
            out_bytes, info = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(_do_pdflock, data, user_pwd, owner_pwd,
                                      allow_print, allow_copy, allow_modify),
                ),
                timeout=45.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Parol qo'yish 45 soniyadan oshdi. Kichikroq fayl tanlang.")
        logger.info(f"pdflock: {len(data)//1024}KB, {info}, {time.time()-t0:.1f}s")
        return Response(
            content=out_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=locked.pdf",
                "X-Password": safe_header(user_pwd),
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"pdflock xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"pdflock: {type(e).__name__}: {str(e)[:160]}")

# ─── PDF: Watermark ───────────────────────────────────────────────────────────

# ✅ YANGILANDI: DejaVu fontni bir marta global cache qilish
_wm_font_cached: str = ""

def _get_wm_font() -> str:
    """Register DejaVu-Bold once and cache the result."""
    global _wm_font_cached
    if _wm_font_cached:
        return _wm_font_cached
    import os as _os
    _DEJAVU = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
    if _os.path.exists(_DEJAVU):
        try:
            from reportlab.pdfbase import pdfmetrics as _pm
            from reportlab.pdfbase.ttfonts import TTFont as _TTF
            _pm.registerFont(_TTF("DejaVuSans-Bold", _DEJAVU))
            _wm_font_cached = "DejaVuSans-Bold"
            return _wm_font_cached
        except Exception:
            pass
    _wm_font_cached = "Helvetica-Bold"
    return _wm_font_cached


def _make_wm_page(pw: float, ph: float, text: str,
                   opacity: float = 0.22, angle: int = 42, repeat: bool = True):
    """Create a ReportLab watermark overlay for one (pw×ph) page size."""
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.colors import Color
    from pypdf import PdfReader as _PR

    font_size = max(20, min(56, int(pw * 0.07)))
    spacing = max(pw, ph) * 0.38

    wm_buf = io.BytesIO()
    c = rl_canvas.Canvas(wm_buf, pagesize=(pw, ph))
    c.saveState()
    c.translate(pw / 2, ph / 2)
    c.rotate(angle)
    c.setFillColor(Color(0.5, 0.5, 0.5, alpha=max(0.05, min(0.9, opacity))))
    c.setFont(_get_wm_font(), font_size)
    offsets = (-spacing, 0, spacing) if repeat else (0,)
    for offset in offsets:
        c.drawCentredString(0, offset, text)
    c.restoreState()
    c.save()
    return _PR(io.BytesIO(wm_buf.getvalue())).pages[0]

def _do_watermark(data: bytes, wm_text: str,
                  opacity: float = 0.22, angle: int = 42, repeat: bool = True) -> tuple:
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(io.BytesIO(data))
    if reader.is_encrypted:
        raise ValueError("Himoyalangan PDF ga watermark qo'yib bo'lmaydi.")
    n = len(reader.pages)
    if n == 0:
        raise ValueError("PDF bo'sh (0 sahifa).")
    if n > 200:
        raise ValueError(f"PDF {n} sahifa. Maksimal 200 sahifa qabul qilinadi.")

    text = wm_text.strip()[:60] or "EduBot"

    writer = PdfWriter()
    wm_cache: dict = {}

    for page in reader.pages:
        pw = float(page.mediabox.width)
        ph = float(page.mediabox.height)
        key = (round(pw), round(ph))
        if key not in wm_cache:
            wm_cache[key] = _make_wm_page(pw, ph, text, opacity, angle, repeat)
        page.merge_page(wm_cache[key])
        writer.add_page(page)

    out = io.BytesIO()
    writer.write(out)
    info = f"{n} sahifa • \"{text}\" watermark • {int(opacity*100)}% shaffof"
    return out.getvalue(), info


@app.post("/api/watermark")
@limiter.limit("15/minute")
async def watermark_pdf(
    request: Request,
    file: UploadFile = File(...),
    text: str = Form(""),
    opacity: float = Form(0.22),
    angle: int = Form(42),
    repeat: bool = Form(True),
):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/watermark")
        if data[:4] != b'%PDF':
            raise HTTPException(status_code=422, detail="Bu fayl PDF emas.")
        opacity = max(0.05, min(0.9, opacity))
        angle   = angle % 360
        loop = asyncio.get_event_loop()
        try:
            out_bytes, info = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(_do_watermark, data, text, opacity, angle, repeat),
                ),
                timeout=45.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Watermark qo'yish 45 soniyadan oshdi. Kichikroq fayl tanlang.")
        logger.info(f"watermark: {len(data)//1024}KB, {info}, {time.time()-t0:.1f}s")
        return Response(
            content=out_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=watermarked.pdf",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"watermark xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"watermark: {type(e).__name__}: {str(e)[:160]}")

# ─── PDF: To image ────────────────────────────────────────────────────────────

_PDF2IMG_MAX_PAGES = 25   # more than enough for presentations & reports

_PDF2IMG_VALID_DPI = {72, 96, 150, 300}
_PDF2IMG_VALID_FMT = {"png", "jpeg", "webp"}


def _do_pdf2img(data: bytes, dpi: int = 150, fmt: str = "png", quality: int = 85) -> tuple:
    """PDF → single image or ZIP of images. Single-page → raw image; multi-page → ZIP."""
    import fitz

    dpi     = dpi if dpi in _PDF2IMG_VALID_DPI else 150
    fmt     = fmt if fmt in _PDF2IMG_VALID_FMT else "png"
    quality = max(50, min(95, quality))
    zoom    = dpi / 72
    mat     = fitz.Matrix(zoom, zoom)
    ext     = "jpg" if fmt == "jpeg" else fmt

    doc = fitz.open(stream=data, filetype="pdf")
    if doc.is_encrypted:
        doc.close()
        raise ValueError("PDF parol bilan himoyalangan. Avval parolini oching.")
    total_pages = doc.page_count
    if total_pages == 0:
        doc.close()
        raise ValueError("PDF bo'sh (0 sahifa).")

    render_n  = min(total_pages, _PDF2IMG_MAX_PAGES)
    truncated = total_pages > _PDF2IMG_MAX_PAGES

    # Single page → return raw image (better UX — direct download, no ZIP)
    if render_n == 1:
        pix       = doc[0].get_pixmap(matrix=mat, alpha=False)
        img_bytes = pix.tobytes(fmt, quality=quality) if fmt != "png" else pix.tobytes("png")
        doc.close()
        media = f"image/{fmt}"
        info  = f"✅ 1 sahifa · {dpi} DPI · {fmt.upper()}"
        return img_bytes, info, media, f"page_1.{ext}"

    # Multi-page → ZIP
    padding = len(str(render_n))
    zf_buf  = io.BytesIO()
    total_size = 0
    zip_mode = zipfile.ZIP_STORED if fmt == "png" else zipfile.ZIP_DEFLATED
    with zipfile.ZipFile(zf_buf, "w", zip_mode) as zf:
        for i in range(render_n):
            pix       = doc[i].get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes(fmt, quality=quality) if fmt != "png" else pix.tobytes("png")
            total_size += len(img_bytes)
            zf.writestr(f"page_{str(i + 1).zfill(padding)}.{ext}", img_bytes)
            del pix
    doc.close()

    avg_kb     = total_size // 1024 // render_n if render_n else 0
    info_parts = [f"✅ {render_n} sahifa · {dpi} DPI · {fmt.upper()} · ~{avg_kb} KB/sahifa"]
    if truncated:
        info_parts.append(f"⚠️ Faqat {_PDF2IMG_MAX_PAGES} sahifa (PDF jami {total_pages})")
    return zf_buf.getvalue(), " · ".join(info_parts), "application/zip", "pages.zip"


@app.post("/api/pdf2img")
@limiter.limit("10/minute")
async def pdf_to_img(
    request: Request,
    file: UploadFile = File(...),
    dpi: int = Form(150),
    fmt: str = Form("png"),
    quality: int = Form(85),
):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/pdf2img")

        loop = asyncio.get_event_loop()
        try:
            content, info, media_type, filename = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(_do_pdf2img, data, dpi, fmt, quality),
                ),
                timeout=90.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail=f"Konversiya 90 soniyadan oshdi — maksimal {_PDF2IMG_MAX_PAGES} sahifa.")

        logger.info(f"pdf2img: {len(data)//1024}KB → {len(content)//1024}KB, {time.time()-t0:.1f}s")
        return Response(
            content=content,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"pdf2img xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"pdf2img: {type(e).__name__}: {str(e)[:160]}")

# ─── PDF: Compress ────────────────────────────────────────────────────────────

_COMPRESSPDF_LEVELS = {
    "screen":   (52,  96),   # minimal sifat, maksimal siqish
    "ebook":    (65, 110),   # o'rtacha, ekran uchun
    "printer":  (82, 150),   # yuqori sifat
    "prepress": (92, 200),   # maxsus bosib chiqarish
}


def _do_compresspdf(data: bytes, level: str = "ebook") -> tuple:
    import fitz, subprocess, tempfile

    orig_size = len(data)

    # ── Try Ghostscript first (best for image-heavy PDFs) ────────────
    _GS_SETTINGS = {
        "screen":   "/screen",    # 72 DPI, smallest
        "ebook":    "/ebook",     # 150 DPI, good for screens
        "printer":  "/printer",   # 300 DPI, print quality
        "prepress": "/prepress",  # 300 DPI, colour managed
    }
    gs_level = _GS_SETTINGS.get(level, "/ebook")
    # ✅ Tuzatildi: try/finally bilan vaqtinchalik fayllar har doim tozalanadi,
    #    fayl handle context manager bilan yopiladi, fitz.open() ham close qilinadi
    import os as _os
    fin_path = fout_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as fin:
            fin.write(data)
            fin_path = fin.name
        fout_path = fin_path.replace(".pdf", "_gs_out.pdf")
        proc = subprocess.run(
            ["gs", "-dBATCH", "-dNOPAUSE", "-dQUIET", "-sDEVICE=pdfwrite",
             f"-dPDFSETTINGS={gs_level}", "-dCompatibilityLevel=1.6",
             f"-sOutputFile={fout_path}", fin_path],
            timeout=60, capture_output=True,
        )
        if proc.returncode == 0 and _os.path.exists(fout_path):
            with open(fout_path, "rb") as _gs_f:
                gs_bytes = _gs_f.read()
            if len(gs_bytes) < orig_size:
                saved = max(0, round((1 - len(gs_bytes) / orig_size) * 100))
                _gs_doc = fitz.open(stream=gs_bytes, filetype="pdf")
                try:
                    n_gs = _gs_doc.page_count
                finally:
                    _gs_doc.close()
                return (gs_bytes,
                        f"{n_gs} sahifa • {orig_size // 1024} KB → {len(gs_bytes) // 1024} KB"
                        f" • {saved}% tejaldi (Ghostscript {level})",
                        saved)
    except (FileNotFoundError, Exception):
        pass  # Ghostscript not installed — fall back to PyMuPDF
    finally:
        for _p in (fin_path, fout_path):
            if _p and _os.path.exists(_p):
                try: _os.unlink(_p)
                except Exception: pass

    doc = fitz.open(stream=data, filetype="pdf")
    n = doc.page_count

    if n == 0:
        doc.close()
        raise ValueError("PDF bo'sh (0 sahifa).")
    if n > 100:
        doc.close()
        raise ValueError(f"PDF {n} sahifa. Maksimal 100 sahifa qabul qilinadi.")

    # ── Detect content type (text vs scanned/image) ─────────────────
    sample_chars = sum(len(doc[i].get_text()) for i in range(min(3, n)))
    is_text_pdf = sample_chars > 150

    # ── Path A: text PDF — structure-level compression ───────────────
    try:
        doc.scrub()
    except Exception:
        pass
    buf_a = io.BytesIO()
    doc.save(buf_a, garbage=4, deflate=True, clean=True)
    doc.close()
    result_a = buf_a.getvalue()

    # ── Path B: scanned/image PDF — re-render at lower DPI ───────────
    # Only used when PDF has no extractable text.
    result_b = None
    if not is_text_pdf:
        # Level parametri ustunlik qiladi
        if level in _COMPRESSPDF_LEVELS:
            quality, dpi = _COMPRESSPDF_LEVELS[level]
        else:
            orig_mb = orig_size / (1024 * 1024)
            if orig_mb > 5:
                quality, dpi = 52, 96
            elif orig_mb > 2:
                quality, dpi = 62, 110
            else:
                quality, dpi = 72, 130
        zoom = dpi / 72
        mat  = fitz.Matrix(zoom, zoom)

        doc2    = fitz.open(stream=data, filetype="pdf")
        out_doc = fitz.open()
        for page in doc2:
            pix       = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("jpeg", quality=quality)
            del pix
            new_page = out_doc.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(new_page.rect, stream=img_bytes)
        doc2.close()
        buf_b    = io.BytesIO()
        out_doc.save(buf_b, garbage=4, deflate=True)
        result_b = buf_b.getvalue()

    # ── Choose best result ───────────────────────────────────────────
    if result_b is not None and len(result_b) < len(result_a):
        out = result_b
        method = "rasterize"
    else:
        out = result_a
        method = "deflate"

    # Never return something larger than the original
    if len(out) >= orig_size:
        out    = data
        method = "original"

    saved = max(0, round((1 - len(out) / orig_size) * 100))
    info  = (f"{n} sahifa • {orig_size // 1024} KB → {len(out) // 1024} KB"
             f" • {saved}% tejaldi"
             + (" (matn saqlanadi)" if method == "deflate" else ""))
    return out, info, saved

@app.post("/api/compresspdf")
@limiter.limit("10/minute")
async def compress_pdf(
    request: Request,
    file: UploadFile = File(...),
    level: str = Form("ebook"),   # screen | ebook | printer | prepress
):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/compresspdf")
        if data[:4] != b'%PDF':
            raise HTTPException(status_code=422, detail="Bu fayl PDF emas.")
        if level not in _COMPRESSPDF_LEVELS:
            level = "ebook"
        loop = asyncio.get_event_loop()
        fn = functools.partial(_do_compresspdf, data, level)
        try:
            out_bytes, info, saved = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, fn),
                timeout=60.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Siqish 60 soniyadan oshdi. Kichikroq fayl yoki kamroq sahifa tanlang.")
        logger.info(f"compresspdf: {info}, {time.time()-t0:.1f}s")
        return Response(
            content=out_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=compressed.pdf",
                "X-Saved-Percent": str(saved),
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"compresspdf xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500,
            detail=f"compresspdf: {type(e).__name__}: {str(e)[:160]}")

# ─── PDF → DOCX ───────────────────────────────────────────────────────────────

@app.post("/api/pdf2docx")
@limiter.limit("10/minute")
async def pdf_to_docx(request: Request, file: UploadFile = File(...)):
    t0 = time.time()
    uid = uuid.uuid4().hex
    pdf_path = f"/tmp/edubot_{uid}.pdf"
    docx_path = f"/tmp/edubot_{uid}.docx"
    try:
        data = await file.read()
        check_size(data, "/api/pdf2docx")

        # ── 1. Sahifalar sonini tekshirish ───────────────────────────
        from pypdf import PdfReader
        try:
            reader = PdfReader(io.BytesIO(data))
            page_count = len(reader.pages)
        except Exception:
            raise HTTPException(status_code=422,
                detail="PDF fayl ochib bo'lmadi. Fayl buzilgan yoki parol bilan himoyalangan.")

        if page_count == 0:
            raise HTTPException(status_code=422, detail="PDF bo'sh (0 sahifa).")

        MAX_PAGES = 50
        if page_count > MAX_PAGES:
            raise HTTPException(status_code=400,
                detail=f"PDF {page_count} sahifali. Maksimal {MAX_PAGES} sahifa. "
                       f"PDF'ni bo'lib (split) qayta urinib ko'ring.")

        # ── 2. Skanerlangan PDF aniqlaish (matn yo'qligi) ────────────
        sample = min(3, page_count)
        total_chars = sum(
            len((reader.pages[i].extract_text() or "").strip())
            for i in range(sample)
        )
        is_scanned = total_chars < 20

        # ── 3. Vaqtinchalik faylga yozish ────────────────────────────
        with open(pdf_path, "wb") as f:
            f.write(data)

        # ── 4. Thread pool'da konversiya (event loop bloklanmaydi) ───
        loop = asyncio.get_event_loop()
        try:
            conv_method = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _do_pdf2docx_best, pdf_path, docx_path, is_scanned),
                timeout=120.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail=f"Konversiya {page_count} sahifa uchun 120 soniyadan oshdi. "
                       f"Faylni kichikroq qismlarga bo'lib qayta urinib ko'ring.")

        # ── 5. Natijani o'qish ────────────────────────────────────────
        if not os.path.exists(docx_path):
            raise HTTPException(status_code=500, detail="Konversiya muvaffaqiyatsiz (fayl yaratilmadi).")

        with open(docx_path, "rb") as f:
            out = f.read()

        if len(out) < 2000:
            raise HTTPException(status_code=422,
                detail="Konversiya natijasi bo'sh. PDF tuzilishi qo'llab-quvvatlanmaydi.")

        # ── 6. Info xabar ─────────────────────────────────────────────
        info_parts = [f"{page_count} sahifa aylantirildi ({conv_method})"]
        if is_scanned:
            info_parts.append("OCR orqali")
        info = " - ".join(info_parts)

        elapsed = time.time() - t0
        logger.info(f"pdf2docx [{conv_method}]: {page_count}p {'scan' if is_scanned else 'text'}, "
                    f"{len(data)//1024}KB → {len(out)//1024}KB, {elapsed:.1f}s")

        return Response(
            content=out,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={
                "Content-Disposition": "attachment; filename=converted.docx",
                "X-Info": safe_header(info),
                "X-Page-Count": str(page_count),
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"pdf2docx xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500,
            detail=f"pdf2docx: {type(e).__name__}: {str(e)[:160]}")
    finally:
        for path in (pdf_path, docx_path):
            try:
                if os.path.exists(path):
                    os.unlink(path)
            except Exception:
                pass

# ─── DOCX → PDF ───────────────────────────────────────────────────────────────

def _rl_fonts():
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    try:
        pdfmetrics.registerFont(TTFont('DV', FONT_REGULAR))
        pdfmetrics.registerFont(TTFont('DV-Bold', FONT_BOLD))
        return 'DV', 'DV-Bold'
    except Exception:
        return 'Helvetica', 'Helvetica-Bold'


def _do_docx2pdf(data: bytes, fn_regular: str, fn_bold: str) -> bytes:
    """
    CPU-bound DOCX→PDF conversion using python-docx + ReportLab.
    Preserves: headings (H1–H4), bold/italic/underline/color/font-size,
               paragraph alignment, tables (borders + alternating rows),
               embedded images, bulleted & numbered lists.
    Runs in _converter_pool (non-blocking).
    """
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph as DocxParagraph
    from docx.table import Table as DocxTable
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        Image as RLImage, PageBreak,
    )
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    from PIL import Image as PILImage

    doc = Document(io.BytesIO(data))
    PAGE_W = A4[0] - 4 * cm  # usable text width

    ALIGN_MAP = {
        WD_ALIGN_PARAGRAPH.LEFT:    TA_LEFT,
        WD_ALIGN_PARAGRAPH.CENTER:  TA_CENTER,
        WD_ALIGN_PARAGRAPH.RIGHT:   TA_RIGHT,
        WD_ALIGN_PARAGRAPH.JUSTIFY: TA_JUSTIFY,
        None: TA_LEFT,
    }

    # ── Base styles ────────────────────────────────────────────────
    S = {
        'h1':  ParagraphStyle('h1',  fontName=fn_bold,    fontSize=20, leading=26, spaceAfter=8,  spaceBefore=18),
        'h2':  ParagraphStyle('h2',  fontName=fn_bold,    fontSize=16, leading=21, spaceAfter=6,  spaceBefore=14),
        'h3':  ParagraphStyle('h3',  fontName=fn_bold,    fontSize=13, leading=17, spaceAfter=5,  spaceBefore=10),
        'h4':  ParagraphStyle('h4',  fontName=fn_bold,    fontSize=12, leading=15, spaceAfter=4,  spaceBefore=8),
        'body':ParagraphStyle('body',fontName=fn_regular, fontSize=11, leading=16, spaceAfter=6),
        'lb':  ParagraphStyle('lb',  fontName=fn_regular, fontSize=11, leading=16, spaceAfter=3,  leftIndent=18),
        'ln':  ParagraphStyle('ln',  fontName=fn_regular, fontSize=11, leading=16, spaceAfter=3,  leftIndent=18),
        'tc':  ParagraphStyle('tc',  fontName=fn_regular, fontSize=9,  leading=13),
        'tch': ParagraphStyle('tch', fontName=fn_bold,    fontSize=9,  leading=13),
    }

    _style_ctr = [0]

    def derived(base, alignment, extra_indent=0):
        _style_ctr[0] += 1
        return ParagraphStyle(
            f'_d{_style_ctr[0]}',
            parent=base,
            alignment=alignment,
            leftIndent=base.leftIndent + extra_indent,
        )

    def esc(t: str) -> str:
        return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    def runs_to_markup(para) -> str:
        """Convert paragraph runs to ReportLab XML markup."""
        if not para.runs:
            return esc(para.text or "")
        parts = []
        for run in para.runs:
            if not run.text:
                continue
            t = esc(run.text)
            # Font color
            try:
                clr = run.font.color
                if clr and clr.type and clr.rgb:
                    t = f'<font color="#{clr.rgb}">{t}</font>'
            except Exception:
                pass
            # Font size (only if meaningfully different from default 11pt)
            try:
                if run.font.size and run.font.size.pt:
                    fs = round(run.font.size.pt)
                    if fs > 0 and abs(fs - 11) > 1:
                        t = f'<font size="{fs}">{t}</font>'
            except Exception:
                pass
            if run.bold:      t = f'<b>{t}</b>'
            if run.italic:    t = f'<i>{t}</i>'
            if run.underline: t = f'<u>{t}</u>'
            parts.append(t)
        result = "".join(parts)
        return result if result.strip() else esc(para.text or "")

    def heading_level(element, para) -> int:
        """Return heading level 1-4, or 0 for non-heading."""
        style_name = (para.style.name or "").lower()
        for n in (1, 2, 3, 4, 5, 6):
            if f"heading {n}" in style_name or f"заголовок {n}" in style_name:
                return min(n, 4)
        # Fallback: check outline level in XML
        try:
            pPr = element.find(qn("w:pPr"))
            if pPr is not None:
                ol = pPr.find(qn("w:outlineLvl"))
                if ol is not None:
                    lvl = int(ol.get(qn("w:val"), 9))
                    if lvl <= 3:
                        return lvl + 1
        except Exception:
            pass
        return 0

    # ── Extract images from relationships ──────────────────────────
    images: dict[str, bytes] = {}
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                images[rel.rId] = rel.target_part.blob
    except Exception:
        pass

    # ── Track numbered list counters per numId ─────────────────────
    list_counters: dict[str, int] = {}

    # ── Build story ────────────────────────────────────────────────
    story = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        # ── Paragraph ─────────────────────────────────────────────
        if tag == "p":
            try:
                para = DocxParagraph(element, doc)
                try:
                    style_name = (para.style.name or "").lower()
                except Exception:
                    style_name = ""

                # Images inside this paragraph
                for blip in element.findall(".//" + qn("a:blip")):
                    rId = blip.get(qn("r:embed")) or blip.get(qn("r:link"))
                    if rId and rId in images:
                        try:
                            buf = io.BytesIO(images[rId])
                            pil = PILImage.open(buf)
                            iw, ih = pil.size
                            if iw > PAGE_W:
                                ratio = PAGE_W / iw
                                iw, ih = PAGE_W, ih * ratio
                            out_buf = io.BytesIO()
                            pil.convert("RGB").save(out_buf, format="PNG")
                            out_buf.seek(0)
                            story.append(RLImage(out_buf, width=iw, height=ih))
                            story.append(Spacer(1, 6))
                        except Exception:
                            pass

                markup = runs_to_markup(para)
                if not markup.strip():
                    story.append(Spacer(1, 3))
                    continue

                try:
                    align = ALIGN_MAP.get(para.alignment, TA_LEFT)
                except Exception:
                    align = TA_LEFT
                hlvl = heading_level(element, para)

                if hlvl == 1:
                    st = derived(S["h1"], align)
                elif hlvl == 2:
                    st = derived(S["h2"], align)
                elif hlvl == 3:
                    st = derived(S["h3"], align)
                elif hlvl >= 4:
                    st = derived(S["h4"], align)
                elif "list bullet" in style_name or "list paragraph" in style_name:
                    markup = f"• {markup}"
                    st = derived(S["lb"], align)
                elif "list number" in style_name:
                    try:
                        numId = element.find(".//" + qn("w:numId"))
                        key = numId.get(qn("w:val"), "0") if numId is not None else "0"
                    except Exception:
                        key = "0"
                    list_counters[key] = list_counters.get(key, 0) + 1
                    markup = f"{list_counters[key]}. {markup}"
                    st = derived(S["ln"], align)
                else:
                    st = derived(S["body"], align)

                try:
                    story.append(Paragraph(markup, st))
                except Exception:
                    try:
                        story.append(Paragraph(esc(para.text or ""), S["body"]))
                    except Exception:
                        pass
            except Exception:
                pass  # skip problematic paragraph elements

        # ── Table ─────────────────────────────────────────────────
        elif tag == "tbl":
            try:
                table = DocxTable(element, doc)
                rows_data = []
                for r_idx, row in enumerate(table.rows):
                    try:
                        row_cells = []
                        is_header = r_idx == 0
                        for cell in row.cells:
                            try:
                                cell_para_text = " ".join(p.text for p in cell.paragraphs).strip()
                            except Exception:
                                cell_para_text = ""
                            row_cells.append(Paragraph(esc(cell_para_text),
                                                       S["tch"] if is_header else S["tc"]))
                        rows_data.append(row_cells)
                    except Exception:
                        pass

                if rows_data:
                    col_n = max(len(r) for r in rows_data)
                    col_w = PAGE_W / col_n if col_n else PAGE_W
                    tbl = Table(rows_data, colWidths=[col_w] * col_n, repeatRows=1)
                    tbl.setStyle(TableStyle([
                        ("BACKGROUND",    (0, 0), (-1,  0), colors.HexColor("#e8e8e8")),
                        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, colors.HexColor("#f7f7f7")]),
                        ("GRID",          (0, 0), (-1, -1), 0.5, colors.HexColor("#bbbbbb")),
                        ("TOPPADDING",    (0, 0), (-1, -1), 4),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                        ("LEFTPADDING",   (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING",  (0, 0), (-1, -1), 6),
                        ("VALIGN",        (0, 0), (-1, -1), "TOP"),
                        ("WORDWRAP",      (0, 0), (-1, -1), "WORD"),
                    ]))
                    story.append(tbl)
                    story.append(Spacer(1, 10))
            except Exception:
                pass  # skip problematic table elements

        # ── Page break ────────────────────────────────────────────
        elif tag == "sectPr":
            pass  # section properties — skip

    if not story:
        story.append(Paragraph("(Bo'sh hujjat)", S["body"]))

    # ── Build PDF ─────────────────────────────────────────────────
    out = io.BytesIO()
    pdf_doc = SimpleDocTemplate(
        out, pagesize=A4,
        leftMargin=2 * cm, rightMargin=2 * cm,
        topMargin=2.5 * cm, bottomMargin=2.5 * cm,
        title=doc.core_properties.title or "EduBot Document",
        author="EduBot",
    )
    try:
        pdf_doc.build(story)
    except Exception:
        # Fallback: rebuild with only safe text/spacer elements
        safe = [item for item in story
                if isinstance(item, (Paragraph, Spacer, PageBreak))]
        if not safe:
            safe = [Paragraph("(Hujjatni aylantirish qisman bajarildi)", S["body"])]
        out2 = io.BytesIO()
        pdf_doc2 = SimpleDocTemplate(
            out2, pagesize=A4,
            leftMargin=2 * cm, rightMargin=2 * cm,
            topMargin=2.5 * cm, bottomMargin=2.5 * cm,
            title="EduBot Document", author="EduBot",
        )
        pdf_doc2.build(safe)
        return out2.getvalue()
    return out.getvalue()


def _do_docx2pdf_libreoffice(data: bytes) -> bytes:
    """LibreOffice headless → highest quality. HOME isolated per call to avoid concurrency conflicts."""
    import subprocess
    with tempfile.TemporaryDirectory() as tmpdir:
        docx_path = os.path.join(tmpdir, "input.docx")
        pdf_path  = os.path.join(tmpdir, "input.pdf")
        with open(docx_path, "wb") as f:
            f.write(data)
        env = os.environ.copy()
        env["HOME"] = tmpdir  # unique profile dir — prevents concurrent LibreOffice conflicts
        result = subprocess.run(
            ["libreoffice", "--headless", "--norestore",
             "--convert-to", "pdf", "--outdir", tmpdir, docx_path],
            env=env, capture_output=True, timeout=50,
        )
        if result.returncode != 0:
            stderr = (result.stderr or b"").decode(errors="replace")[:300]
            raise RuntimeError(f"libreoffice exit {result.returncode}: {stderr}")
        if not os.path.exists(pdf_path) or os.path.getsize(pdf_path) < 100:
            raise RuntimeError("libreoffice PDF chiqarilmadi yoki bo'sh")
        with open(pdf_path, "rb") as f:
            return f.read()


def _do_docx2pdf_mammoth(data: bytes) -> bytes:
    """mammoth → HTML → WeasyPrint PDF. Mid-quality, good Unicode/Cyrillic support."""
    import mammoth
    from weasyprint import HTML
    result = mammoth.convert_to_html(io.BytesIO(data))
    html_body = result.value or ""
    if not html_body.strip():
        raise ValueError("mammoth: bo'sh HTML natija")
    html = (
        "<!DOCTYPE html><html><head>"
        "<meta charset='utf-8'>"
        "<style>body{font-family:DejaVu Sans,sans-serif;font-size:11pt;margin:2cm;}"
        "h1{font-size:18pt}h2{font-size:15pt}h3{font-size:13pt}"
        "table{border-collapse:collapse;width:100%}"
        "td,th{border:1px solid #bbb;padding:4px 6px}"
        "</style></head><body>"
        + html_body
        + "</body></html>"
    )
    pdf_bytes = HTML(string=html).write_pdf()
    if not pdf_bytes or len(pdf_bytes) < 100:
        raise RuntimeError("weasyprint: bo'sh PDF natija")
    return pdf_bytes


def _do_docx2pdf_best(data: bytes, fn_regular: str, fn_bold: str) -> tuple:
    """Tries LibreOffice → mammoth+WeasyPrint → ReportLab. Returns (pdf_bytes, method_name)."""
    attempts = [
        ("libreoffice", lambda: _do_docx2pdf_libreoffice(data)),
        ("mammoth",     lambda: _do_docx2pdf_mammoth(data)),
        ("reportlab",   lambda: _do_docx2pdf(data, fn_regular, fn_bold)),
    ]
    last_err = None
    for name, fn in attempts:
        try:
            result = fn()
            if result and len(result) > 500:
                return result, name
        except Exception as e:
            last_err = e
            logger.warning(f"docx2pdf [{name}] muvaffaqiyatsiz: {type(e).__name__}: {str(e)[:120]}")
    raise RuntimeError(
        f"Barcha 3 usul (LibreOffice, mammoth, ReportLab) muvaffaqiyatsiz: {type(last_err).__name__}: {str(last_err)[:100]}"
    )


@app.post("/api/docx2pdf")
@limiter.limit("10/minute")
async def docx_to_pdf(request: Request, file: UploadFile = File(...)):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/docx2pdf")
        fn, fn_bold = _rl_fonts()

        # ── 1. Fayl validatsiyasi ─────────────────────────────────
        fname = (file.filename or "").lower()
        if fname.endswith(".doc") and not fname.endswith(".docx"):
            raise HTTPException(status_code=400,
                detail="Eski .doc format qo'llab-quvvatlanmaydi. "
                       "Faylni Microsoft Word'da ochib, 'Word hujjati (.docx)' formatida saqlang.")

        # ── 2. Hujjat tarkibini oldindan tekshirish ───────────────
        try:
            from docx import Document as _DocCheck
            _doc_check = _DocCheck(io.BytesIO(data))
            para_count  = sum(1 for p in _doc_check.paragraphs if p.text.strip())
            table_count = len(_doc_check.tables)
            del _doc_check
        except Exception as e:
            raise HTTPException(status_code=422,
                detail="DOCX fayl ochib bo'lmadi. Fayl buzilgan yoki parol bilan himoyalangan.")

        if para_count == 0 and table_count == 0:
            raise HTTPException(status_code=422, detail="Hujjat bo'sh (matn yoki jadval topilmadi).")

        # ── 3. Thread pool'da konversiya (90s timeout) ────────────
        loop = asyncio.get_event_loop()
        try:
            pdf_bytes, conv_method = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _do_docx2pdf_best, data, fn, fn_bold),
                timeout=90.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Konversiya 90 soniyadan oshdi. Hujjat juda katta yoki murakkab.")

        if len(pdf_bytes) < 500:
            raise HTTPException(status_code=500, detail="Konversiya natijasi bo'sh. Qayta urinib ko'ring.")

        # ── 4. Info xabar ─────────────────────────────────────────
        info_parts = [f"{para_count} paragraf"]
        if table_count:
            info_parts.append(f"{table_count} jadval")
        info_parts.append(f"aylantiriLdi ({conv_method})")
        info = " - ".join(info_parts)

        logger.info(f"docx2pdf [{conv_method}]: {para_count}p+{table_count}t, "
                    f"{len(data)//1024}KB → {len(pdf_bytes)//1024}KB, {time.time()-t0:.1f}s")

        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=document.pdf",
                "X-Info": safe_header(info),
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"docx2pdf xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500,
            detail=f"docx2pdf: {type(e).__name__}: {str(e)[:160]}")


# ─── DOCX Edit ────────────────────────────────────────────────────────────────

def _do_docxedit(data: bytes, find: str, replace: str, case_sensitive: bool) -> tuple:
    """Find-and-replace text in a DOCX file (runs + tables)."""
    from docx import Document
    import copy

    doc = Document(io.BytesIO(data))
    count = 0

    def _replace_in_para(para):
        nonlocal count
        if not para.runs:
            return
        import re as _re

        if case_sensitive:
            # Per-run first (preserves bold/italic/size of each run)
            hit = 0
            for r in para.runs:
                if find in r.text:
                    hit += r.text.count(find)
                    r.text = r.text.replace(find, replace)
            if hit:
                count += hit
                return
            # Cross-run fallback
            full = "".join(r.text for r in para.runs)
            if find not in full:
                return
            count += full.count(find)
            new_full = full.replace(find, replace)
        else:
            pattern = _re.compile(_re.escape(find), _re.IGNORECASE)
            # Per-run first
            hit = 0
            for r in para.runs:
                ms = pattern.findall(r.text)
                if ms:
                    hit += len(ms)
                    r.text = pattern.sub(replace, r.text)
            if hit:
                count += hit
                return
            # Cross-run fallback
            full = "".join(r.text for r in para.runs)
            if not pattern.search(full):
                return
            count += len(pattern.findall(full))
            new_full = pattern.sub(replace, full)

        # Cross-run: put result in first run, clear the rest (first run's
        # formatting is kept; formatting of other runs is unavoidably lost)
        para.runs[0].text = new_full
        for r in para.runs[1:]:
            r.text = ""

    for para in doc.paragraphs:
        _replace_in_para(para)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _replace_in_para(para)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue(), count


@app.post("/api/docxedit")
@limiter.limit("15/minute")
async def docx_edit(
    request: Request,
    file: UploadFile = File(...),
    find: str = Form(...),
    replace: str = Form(...),
    case_sensitive: bool = Form(False),
):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/docxedit")

        if not find:
            raise HTTPException(status_code=400, detail="'find' bo'sh bo'lmasligi kerak")
        if len(find) > 500:
            raise HTTPException(status_code=400, detail="'find' 500 belgidan oshmasligi kerak")
        if len(replace) > 500:
            raise HTTPException(status_code=400, detail="'replace' 500 belgidan oshmasligi kerak")

        loop = asyncio.get_event_loop()
        fn = functools.partial(_do_docxedit, data, find, replace, case_sensitive)
        out_bytes, count = await asyncio.wait_for(
            loop.run_in_executor(_converter_pool, fn),
            timeout=30.0,
        )

        orig_name = os.path.splitext(file.filename or "document")[0]
        out_name = f"{orig_name}_edited.docx"
        enc_name = url_quote(out_name)

        logger.info(f"docxedit: {count} almashtirish, {time.time()-t0:.1f}s")
        return Response(
            content=out_bytes,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{enc_name}",
                "X-Info": safe_header(f"{count} ta so'z almashtirildi"),
            },
        )
    except HTTPException:
        raise
    except asyncio.TimeoutError:
        raise HTTPException(status_code=504, detail="Vaqt tugadi. Kichikroq fayl yuboring.")
    except Exception as e:
        logger.error(f"docxedit xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500,
            detail=f"docxedit: {type(e).__name__}: {str(e)[:160]}")


# ─── Image → PDF ──────────────────────────────────────────────────────────────

@app.post("/api/img2pdf")
@limiter.limit("20/minute")
async def img_to_pdf(
    request: Request,
    file: UploadFile = File(...),
    page_size: str   = Form("a4"),    # a4 | a3 | a5 | letter | legal | original
    margin_mm: float = Form(10.0),    # 0–30 mm
    fit_mode: str    = Form("fit"),   # fit | fill | center
    dpi: int         = Form(150),     # 72 | 96 | 150 | 300
):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/img2pdf")

        if not is_image_bytes(data):
            raise HTTPException(status_code=422,
                detail="Rasm fayli emas. JPG, PNG, WebP, TIFF, BMP, GIF yuklang.")

        page_size, margin_mm, fit_mode, dpi = img2pdf_validate(page_size, margin_mm, fit_mode, dpi)
        fname = file.filename or "image"

        loop = asyncio.get_event_loop()
        try:
            pdf_bytes, info = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(do_img2pdf_single, data, page_size, margin_mm, fit_mode, dpi, fname),
                ),
                timeout=30.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408, detail="Konversiya 30 soniyadan oshdi.")
        except ValueError as e:
            raise HTTPException(status_code=422, detail=str(e))

        logger.info(f"img2pdf: {len(data)//1024}KB → {len(pdf_bytes)//1024}KB [{page_size} {fit_mode}] {time.time()-t0:.1f}s")
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=image.pdf",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"img2pdf xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"img2pdf: {type(e).__name__}: {str(e)[:160]}")

# ─── Images → PDF ─────────────────────────────────────────────────────────────

@app.post("/api/imgs2pdf")
@limiter.limit("10/minute")
async def imgs_to_pdf(
    request: Request,
    files: List[UploadFile] = File(...),
    page_size: str   = Form("a4"),
    margin_mm: float = Form(10.0),
    fit_mode: str    = Form("fit"),
    dpi: int         = Form(150),
):
    t0 = time.time()
    try:
        if not files:
            raise HTTPException(status_code=400, detail="Rasm yuklanmadi.")
        if len(files) > 50:
            raise HTTPException(status_code=400,
                detail=f"Juda ko'p rasm ({len(files)} ta). Maksimal 50 ta.")

        all_data, all_names, total = [], [], 0
        for f in files:
            d = await f.read()
            total += len(d)
            if total > MAX_FILE_BYTES * 3:
                raise HTTPException(status_code=413,
                    detail="Rasmlarning umumiy hajmi juda katta (max 60 MB).")
            if not is_image_bytes(d):
                continue  # rasm bo'lmagan fayllar o'tkazib yuboriladi
            all_data.append(d)
            all_names.append(f.filename or f"image_{len(all_data)}")

        if not all_data:
            raise HTTPException(status_code=422, detail="Yuklangan fayllar ichida rasm topilmadi.")

        page_size, margin_mm, fit_mode, dpi = img2pdf_validate(page_size, margin_mm, fit_mode, dpi)

        loop = asyncio.get_event_loop()
        try:
            pdf_bytes, info = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(do_imgs2pdf_multi, all_data, all_names, page_size, margin_mm, fit_mode, dpi),
                ),
                timeout=90.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408, detail="Konversiya 90 soniyadan oshdi.")
        except ValueError as e:
            raise HTTPException(status_code=422, detail=str(e))

        logger.info(f"imgs2pdf: {len(all_data)} rasm {total//1024}KB → {len(pdf_bytes)//1024}KB [{page_size}] {time.time()-t0:.1f}s")
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=images.pdf",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"imgs2pdf xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"imgs2pdf: {type(e).__name__}: {str(e)[:160]}")

# ─── 3x4 Passport foto ────────────────────────────────────────────────────────

def _do_photo3x4(data: bytes, bg: str = "white", sheet: bool = True) -> tuple[bytes, str]:
    """
    Rasm → 3×4 sm passport foto (300 DPI).
    1. OpenCV Haar Cascade orqali yuz aniqlanadi.
    2. Yuz markaziga qarab crop qilinadi (foto balandligining ~70% yuz).
    3. 354×472 px ga resize (3×4 sm @ 300 DPI).
    4. sheet=True → A4 da 2×3 = 6 ta foto chiqariladi.
    """
    import cv2
    import numpy as np
    from PIL import Image, ImageOps, ImageDraw

    W, H = 354, 472  # 3×4 sm @ 300 DPI

    BG_COLORS = {"white": (255, 255, 255), "blue": (100, 149, 237), "gray": (200, 200, 200)}
    bg_rgb = BG_COLORS.get(bg, (255, 255, 255))

    img_pil = Image.open(io.BytesIO(data))
    img_pil = ImageOps.exif_transpose(img_pil)
    if img_pil.mode != "RGB":
        img_pil = img_pil.convert("RGB")

    iw, ih = img_pil.size

    # OpenCV yuz aniqlash
    img_cv = cv2.cvtColor(np.array(img_pil), cv2.COLOR_RGB2GRAY)
    cascade = cv2.CascadeClassifier(cv2.data.haarcascades + "haarcascade_frontalface_default.xml")
    faces = cascade.detectMultiScale(img_cv, scaleFactor=1.1, minNeighbors=5, minSize=(60, 60))

    if len(faces) > 0:
        # Eng katta yuzni ol
        fx, fy, fw, fh = sorted(faces, key=lambda f: f[2] * f[3], reverse=True)[0]

        # Foto balandligi: yuz → fotoning ~65% ni egallaydi
        photo_h = int(fh / 0.65)
        photo_w = int(photo_h * W / H)

        # Yuz fotoning tepasidan ~25% pastda bo'lsin
        cx = fx + fw // 2
        cy = fy + fh // 2
        x0 = cx - photo_w // 2
        y0 = cy - int(photo_h * 0.30)

        x0 = max(0, min(x0, iw - photo_w))
        y0 = max(0, min(y0, ih - photo_h))
        crop = img_pil.crop((x0, y0, x0 + photo_w, y0 + photo_h))
        detected = True
    else:
        # Yuz topilmadi — rasmni markazdan portrait nisbatda crop qilish
        target_ratio = W / H
        if iw / ih > target_ratio:
            crop_w = int(ih * target_ratio)
            crop = img_pil.crop(((iw - crop_w) // 2, 0, (iw + crop_w) // 2, ih))
        else:
            crop_h = int(iw / target_ratio)
            crop = img_pil.crop((0, 0, iw, crop_h))
        detected = False

    photo = crop.resize((W, H), Image.LANCZOS)

    if not sheet:
        buf = io.BytesIO()
        photo.save(buf, format="JPEG", quality=95, dpi=(300, 300))
        msg = "✅ 3×4 sm · 300 DPI" + ("" if detected else " · Yuz topilmadi, markaz crop")
        return buf.getvalue(), msg

    # A4 varaqdagi 6 ta foto (2 ustun × 3 qator) @ 300 DPI
    A4W, A4H = 2480, 3508
    MARGIN, GAP = 100, 50

    canvas = Image.new("RGB", (A4W, A4H), (255, 255, 255))
    draw = ImageDraw.Draw(canvas)

    for row in range(3):
        for col in range(2):
            px = MARGIN + col * (W + GAP)
            py = MARGIN + row * (H + GAP)
            canvas.paste(photo, (px, py))
            # Kesish chizig'i (yupqa kulrang ramka)
            draw.rectangle([px - 1, py - 1, px + W, py + H], outline=(180, 180, 180), width=1)

    buf = io.BytesIO()
    canvas.save(buf, format="JPEG", quality=95, dpi=(300, 300))
    msg = "✅ A4 · 6 ta 3×4 sm foto · 300 DPI" + ("" if detected else " · Yuz topilmadi")
    return buf.getvalue(), msg


@app.post("/api/photo3x4")
@limiter.limit("15/minute")
async def photo3x4(
    request: Request,
    file: UploadFile  = File(...),
    bg: str           = Form("white"),   # white | blue | gray
    sheet: str        = Form("true"),    # "true" → A4 da 6 ta, "false" → bitta
):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/photo3x4")

        if not is_image_bytes(data):
            raise HTTPException(status_code=422,
                detail="Rasm fayli emas. JPG yoki PNG yuklang.")

        bg      = bg if bg in ("white", "blue", "gray") else "white"
        do_sheet = sheet.lower() not in ("false", "0", "no")

        loop = asyncio.get_event_loop()
        try:
            jpg_bytes, info = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(_do_photo3x4, data, bg, do_sheet),
                ),
                timeout=30.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408, detail="Konversiya 30 soniyadan oshdi.")

        fname = "photo3x4_sheet.jpg" if do_sheet else "photo3x4.jpg"
        logger.info(f"photo3x4: {len(data)//1024}KB → {len(jpg_bytes)//1024}KB [{bg} sheet={do_sheet}] {time.time()-t0:.1f}s")
        return Response(
            content=jpg_bytes,
            media_type="image/jpeg",
            headers={
                "Content-Disposition": f"attachment; filename={fname}",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"photo3x4 xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"photo3x4: {type(e).__name__}: {str(e)[:160]}")


# ─── CV (Resume) Generator ────────────────────────────────────────────────────

_CV_TEMPLATES = {
    "modern":  {"accent": (79, 70, 229),  "bg": (248, 247, 255), "text": (15, 15, 30)},
    "classic": {"accent": (30, 64, 175),  "bg": (255, 255, 255), "text": (17, 24, 39)},
    "minimal": {"accent": (5, 150, 105),  "bg": (255, 255, 255), "text": (17, 24, 39)},
    "dark":    {"accent": (167, 139, 250), "bg": (15, 15, 30),   "text": (240, 240, 255)},
}


def _do_cv(
    name: str,
    title: str,
    email: str,
    phone: str,
    location: str,
    summary: str,
    skills: list[str],
    education: list[dict],   # [{degree, school, year}]
    experience: list[dict],  # [{position, company, period, desc}]
    languages: list[str],
    template: str = "modern",
) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import Color, HexColor, white, black
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle
    )
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import os as _os

    tmpl = _CV_TEMPLATES.get(template, _CV_TEMPLATES["modern"])
    accent  = Color(tmpl["accent"][0]/255, tmpl["accent"][1]/255, tmpl["accent"][2]/255)
    bg_col  = Color(tmpl["bg"][0]/255,     tmpl["bg"][1]/255,     tmpl["bg"][2]/255)
    txt_col = Color(tmpl["text"][0]/255,   tmpl["text"][1]/255,   tmpl["text"][2]/255)

    # Register DejaVu for Unicode (Cyrillic + Latin)
    _dv_dir = "/usr/share/fonts/truetype/dejavu"
    _font_reg  = "DejaVuSans"
    _font_bold = "DejaVuSans-Bold"
    try:
        if _os.path.exists(f"{_dv_dir}/DejaVuSans.ttf"):
            pdfmetrics.registerFont(TTFont(_font_reg,  f"{_dv_dir}/DejaVuSans.ttf"))
            pdfmetrics.registerFont(TTFont(_font_bold, f"{_dv_dir}/DejaVuSans-Bold.ttf"))
    except Exception:
        _font_reg = _font_bold = "Helvetica"

    W, H = A4
    ML, MR, MT, MB = 18*mm, 18*mm, 18*mm, 18*mm

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=ML, rightMargin=MR, topMargin=MT, bottomMargin=MB,
    )

    def _style(name_s, **kw):
        base = kw.pop("parent", None)
        s = ParagraphStyle(name_s, fontName=_font_reg, textColor=txt_col, **kw)
        return s

    s_name    = _style("cvname",  fontName=_font_bold, fontSize=26, leading=30, textColor=accent)
    s_title   = _style("cvtitle", fontName=_font_reg,  fontSize=13, leading=16, textColor=txt_col, spaceAfter=2)
    s_contact = _style("cvcont",  fontName=_font_reg,  fontSize=9,  leading=13, textColor=txt_col)
    s_section = _style("cvsec",   fontName=_font_bold, fontSize=11, leading=14, textColor=accent,
                        spaceBefore=10, spaceAfter=3)
    s_body    = _style("cvbody",  fontName=_font_reg,  fontSize=9.5, leading=14, textColor=txt_col)
    s_bold    = _style("cvbold",  fontName=_font_bold, fontSize=9.5, leading=14, textColor=txt_col)
    s_small   = _style("cvsmall", fontName=_font_reg,  fontSize=8.5, leading=12, textColor=txt_col)
    s_skill   = _style("cvskill", fontName=_font_reg,  fontSize=9,   leading=13, textColor=txt_col)

    def hr():
        return HRFlowable(width="100%", thickness=0.8, color=accent, spaceAfter=4, spaceBefore=2)

    story = []

    # ── Header ────────────────────────────────────────────────────────
    story.append(Paragraph(name or "Ism Familiya", s_name))
    if title:
        story.append(Paragraph(title, s_title))

    contacts = [x for x in [email, phone, location] if x]
    if contacts:
        story.append(Paragraph("  •  ".join(contacts), s_contact))
    story.append(Spacer(1, 4*mm))

    # ── Summary ───────────────────────────────────────────────────────
    if summary:
        story.append(hr())
        story.append(Paragraph("Haqida / О себе", s_section))
        story.append(Paragraph(summary, s_body))
        story.append(Spacer(1, 2*mm))

    # ── Experience ────────────────────────────────────────────────────
    if experience:
        story.append(hr())
        story.append(Paragraph("Ish tajribasi / Опыт работы", s_section))
        for exp in experience:
            pos     = exp.get("position", "")
            company = exp.get("company", "")
            period  = exp.get("period", "")
            desc    = exp.get("desc", "")
            header  = f"<b>{pos}</b>" + (f" — {company}" if company else "")
            story.append(Paragraph(header, s_bold))
            if period:
                story.append(Paragraph(period, s_small))
            if desc:
                story.append(Paragraph(desc, s_body))
            story.append(Spacer(1, 2*mm))

    # ── Education ─────────────────────────────────────────────────────
    if education:
        story.append(hr())
        story.append(Paragraph("Ta'lim / Образование", s_section))
        for edu in education:
            degree = edu.get("degree", "")
            school = edu.get("school", "")
            year   = edu.get("year", "")
            header = f"<b>{degree}</b>" + (f" — {school}" if school else "")
            story.append(Paragraph(header, s_bold))
            if year:
                story.append(Paragraph(year, s_small))
            story.append(Spacer(1, 2*mm))

    # ── Skills ────────────────────────────────────────────────────────
    if skills:
        story.append(hr())
        story.append(Paragraph("Ko'nikmalar / Навыки", s_section))
        # Chips in 3 columns
        rows = [skills[i:i+3] for i in range(0, len(skills), 3)]
        tdata = []
        for row in rows:
            while len(row) < 3:
                row.append("")
            tdata.append([Paragraph(f"• {s}", s_skill) for s in row])
        if tdata:
            col_w = (W - ML - MR) / 3
            tbl = Table(tdata, colWidths=[col_w]*3)
            tbl.setStyle(TableStyle([
                ("VALIGN", (0,0), (-1,-1), "TOP"),
                ("LEFTPADDING",  (0,0), (-1,-1), 0),
                ("RIGHTPADDING", (0,0), (-1,-1), 4),
                ("BOTTOMPADDING",(0,0), (-1,-1), 3),
            ]))
            story.append(tbl)
        story.append(Spacer(1, 2*mm))

    # ── Languages ─────────────────────────────────────────────────────
    if languages:
        story.append(hr())
        story.append(Paragraph("Tillar / Языки", s_section))
        story.append(Paragraph("  •  ".join(languages), s_body))

    doc.build(story)
    return buf.getvalue()


@app.post("/api/cv")
@limiter.limit("10/minute")
async def make_cv(request: Request):
    t0 = time.time()
    try:
        body = await request.json()

        name     = (body.get("name") or "").strip()[:80]
        title    = (body.get("title") or "").strip()[:100]
        email    = (body.get("email") or "").strip()[:80]
        phone    = (body.get("phone") or "").strip()[:30]
        location = (body.get("location") or "").strip()[:80]
        summary  = (body.get("summary") or "").strip()[:800]
        template = (body.get("template") or "modern").strip().lower()
        if template not in _CV_TEMPLATES:
            template = "modern"

        raw_skills = body.get("skills") or []
        if isinstance(raw_skills, str):
            raw_skills = [s.strip() for s in raw_skills.split(",") if s.strip()]
        skills = [str(s).strip()[:40] for s in raw_skills[:20]]

        raw_langs = body.get("languages") or []
        if isinstance(raw_langs, str):
            raw_langs = [s.strip() for s in raw_langs.split(",") if s.strip()]
        languages = [str(s).strip()[:40] for s in raw_langs[:10]]

        education = []
        for edu in (body.get("education") or [])[:8]:
            if isinstance(edu, dict):
                education.append({
                    "degree": str(edu.get("degree") or "")[:80],
                    "school": str(edu.get("school") or "")[:80],
                    "year":   str(edu.get("year")   or "")[:20],
                })

        experience = []
        for exp in (body.get("experience") or [])[:8]:
            if isinstance(exp, dict):
                experience.append({
                    "position": str(exp.get("position") or "")[:80],
                    "company":  str(exp.get("company")  or "")[:80],
                    "period":   str(exp.get("period")   or "")[:40],
                    "desc":     str(exp.get("desc")     or "")[:400],
                })

        if not name:
            raise HTTPException(status_code=400, detail="Ism (name) majburiy")

        loop = asyncio.get_event_loop()
        pdf_bytes = await asyncio.wait_for(
            loop.run_in_executor(
                _converter_pool,
                functools.partial(
                    _do_cv, name, title, email, phone, location,
                    summary, skills, education, experience, languages, template,
                ),
            ),
            timeout=30.0,
        )
        fname = f"cv_{name.replace(' ', '_')[:30]}.pdf"
        logger.info(f"cv: {name!r} {template} {len(pdf_bytes)//1024}KB {time.time()-t0:.1f}s")
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{fname}"',
                "X-Info": safe_header(f"{template} · {len(experience)} tajriba · {len(education)} ta'lim"),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"cv xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"cv: {type(e).__name__}: {str(e)[:160]}")


# ─── Excel → PDF ──────────────────────────────────────────────────────────────

def _do_xlsx2pdf(data: bytes) -> tuple:
    """
    Excel → PDF with smart column widths, auto landscape, date/number formatting.
    Pass 1 — scan all sheets → determine orientation.
    Pass 2 — build full document.
    Runs in _converter_pool (non-blocking).
    """
    from openpyxl import load_workbook
    from reportlab.platypus import (
        SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak,
    )
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from datetime import datetime as _dt, date as _date, time as _time

    fn, fn_bold = _rl_fonts()
    MAX_ROWS    = 500
    MAX_SHEETS  = 10
    MAX_COLS    = 20
    MARGIN      = 1.2 * cm

    def fmt_val(v) -> str:
        if v is None:
            return ''
        if isinstance(v, (_dt,)):
            return v.strftime('%d.%m.%Y %H:%M') if v.hour or v.minute else v.strftime('%d.%m.%Y')
        if isinstance(v, _date):
            return v.strftime('%d.%m.%Y')
        if isinstance(v, _time):
            return v.strftime('%H:%M')
        if isinstance(v, bool):
            return 'Ha' if v else "Yo'q"
        if isinstance(v, float):
            # Integer-valued float (e.g. 1.0, 2.0)
            if abs(v) < 1e14 and v == int(v):
                return str(int(v))
            return f'{v:g}'
        return str(v)

    def esc(t: str) -> str:
        return t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

    def calc_col_widths(all_rows, col_n: int, avail_w: float) -> list:
        """Proportional widths based on max content length, clamped min/max."""
        max_lens = [1] * col_n
        for row in all_rows[:100]:
            for j in range(col_n):
                v = row[j] if j < len(row) else None
                max_lens[j] = max(max_lens[j], min(len(fmt_val(v)), 40))
        total = sum(max_lens) or col_n
        MIN_W, MAX_W = 1.0 * cm, 9.0 * cm
        widths = [max(MIN_W, min(MAX_W, avail_w * l / total)) for l in max_lens]
        # Normalise so columns fill the page exactly
        scale = avail_w / sum(widths)
        return [w * scale for w in widths]

    # ── Pass 1: quick scan — find max non-empty columns across all sheets ──
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    max_col_n = 0
    for name in wb.sheetnames[:MAX_SHEETS]:
        ws = wb[name]
        for row in ws.iter_rows(values_only=True, max_row=5, max_col=MAX_COLS):
            non_empty = sum(1 for v in row if v is not None)
            max_col_n = max(max_col_n, non_empty)
    wb.close()

    # Landscape if any sheet has > 7 meaningful columns
    if max_col_n > 7:
        page_w, page_h = A4[1], A4[0]
        orientation = "Landscape"
    else:
        page_w, page_h = A4
        orientation = "Portrait"
    avail_w = page_w - 2 * MARGIN

    # ── Pass 2: full build (read_only=False to access cell formatting) ──────
    wb = load_workbook(io.BytesIO(data), read_only=False, data_only=True)

    def _cell_bg(cell) -> str | None:
        try:
            fill = cell.fill
            if fill and fill.fill_type == "solid":
                rgb = fill.fgColor.rgb  # "FFRRGGBB"
                if rgb and rgb not in ("00000000", "FF000000", "FFFFFFFF"):
                    return f"#{rgb[2:]}"
        except Exception:
            pass
        return None

    title_st = ParagraphStyle('xt', fontName=fn_bold, fontSize=12,
                               spaceAfter=6, spaceBefore=4,
                               textColor=colors.HexColor('#1a3a5c'))
    note_st  = ParagraphStyle('xn', fontName=fn, fontSize=8,
                               textColor=colors.HexColor('#888888'), spaceAfter=4)

    story = []
    total_rows   = 0
    total_sheets = 0
    any_truncated = False

    for name in wb.sheetnames[:MAX_SHEETS]:
        ws = wb[name]
        # values_only=False — to access cell formatting (fill, font)
        raw_cells = list(ws.iter_rows(values_only=False,
                                      max_row=MAX_ROWS + 1, max_col=MAX_COLS))
        if not raw_cells:
            continue

        truncated = len(raw_cells) > MAX_ROWS
        raw_cells = raw_cells[:MAX_ROWS]
        if truncated:
            any_truncated = True

        # Extract plain values for column width calc
        raw = [[cell.value for cell in row] for row in raw_cells]

        # Trim trailing all-None columns
        col_n = MAX_COLS
        while col_n > 1:
            if all((r[col_n - 1] if col_n - 1 < len(r) else None) is None
                   for r in raw[:30]):
                col_n -= 1
            else:
                break
        col_n = max(1, col_n)

        # Trim trailing all-None rows
        while raw and all(v is None for v in raw[-1][:col_n]):
            raw.pop()
            raw_cells.pop()
        if not raw:
            continue

        # Font size adapts to column count
        if   col_n <= 4:  fs = 10
        elif col_n <= 7:  fs = 9
        elif col_n <= 11: fs = 8
        else:             fs = 7

        tc_st  = ParagraphStyle(f'tc_{name}',  fontName=fn,      fontSize=fs, leading=fs + 3)
        tch_st = ParagraphStyle(f'tch_{name}', fontName=fn_bold, fontSize=fs, leading=fs + 3)

        col_widths = calc_col_widths(raw, col_n, avail_w)

        # Build table rows + collect cell background colors
        table_data = []
        cell_bgs = []  # list of (row, col, hex_color)
        for r_idx, (row_cells, row_vals) in enumerate(zip(raw_cells, raw)):
            is_hdr = r_idx == 0
            cells = []
            for j in range(col_n):
                cell_obj = row_cells[j] if j < len(row_cells) else None
                v = row_vals[j] if j < len(row_vals) else None
                txt = esc(fmt_val(v))
                is_bold = is_hdr or bool(
                    cell_obj and cell_obj.font and cell_obj.font.bold
                )
                cells.append(Paragraph(txt, tch_st if is_bold else tc_st))
                if cell_obj and not is_hdr:
                    bg = _cell_bg(cell_obj)
                    if bg:
                        cell_bgs.append((r_idx, j, bg))
            table_data.append(cells)

        # Section header
        if story:
            story.append(PageBreak())
        story.append(Paragraph(name, title_st))
        if truncated:
            story.append(Paragraph(
                f'⚠️ Faqat birinchi {MAX_ROWS} qator ko\'rsatildi', note_st))

        tbl = Table(table_data, colWidths=col_widths, repeatRows=1)
        style_cmds = [
            # Header row
            ('BACKGROUND',    (0, 0), (-1,  0), colors.HexColor('#cfe2f3')),
            ('LINEBELOW',     (0, 0), (-1,  0), 1.2, colors.HexColor('#2e7cbf')),
            # Alternating data rows
            ('ROWBACKGROUNDS',(0, 1), (-1, -1),
             [colors.white, colors.HexColor('#f0f6fb')]),
            # Grid
            ('GRID',          (0, 0), (-1, -1), 0.35, colors.HexColor('#b0c8e0')),
            # Padding
            ('TOPPADDING',    (0, 0), (-1, -1), 3),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
            ('LEFTPADDING',   (0, 0), (-1, -1), 4),
            ('RIGHTPADDING',  (0, 0), (-1, -1), 4),
            ('VALIGN',        (0, 0), (-1, -1), 'TOP'),
            ('WORDWRAP',      (0, 0), (-1, -1), 'WORD'),
        ]
        # Apply per-cell background colors from Excel
        for r_idx, c_idx, hex_color in cell_bgs:
            style_cmds.append(
                ('BACKGROUND', (c_idx, r_idx), (c_idx, r_idx), colors.HexColor(hex_color))
            )
        tbl.setStyle(TableStyle(style_cmds))
        story.append(tbl)
        total_rows   += len(raw)
        total_sheets += 1

    wb.close()

    if not story:
        raise ValueError("Excel fayl bo'sh yoki o'qib bo'lmadi.")

    out = io.BytesIO()
    doc = SimpleDocTemplate(
        out, pagesize=(page_w, page_h),
        leftMargin=MARGIN, rightMargin=MARGIN,
        topMargin=MARGIN + 0.3 * cm, bottomMargin=MARGIN,
        title="EduBot — Excel to PDF",
    )
    doc.build(story)

    sheets_word = f"{total_sheets} varaq" if total_sheets > 1 else "1 varaq"
    info = f"✅ {sheets_word} · {total_rows} qator · {orientation}"
    if any_truncated:
        info += f" · (max {MAX_ROWS} qator/varaq)"
    return out.getvalue(), info


@app.post("/api/xlsx2pdf")
@limiter.limit("10/minute")
async def xlsx_to_pdf(request: Request, file: UploadFile = File(...)):
    t0 = time.time()
    try:
        data = await file.read()
        check_size(data, "/api/xlsx2pdf")
        # xlsx is a ZIP-based format; check magic bytes PK
        if len(data) < 4 or data[:2] != b'PK':
            raise HTTPException(status_code=422,
                detail="Bu fayl Excel (.xlsx) emas. .xlsx yoki .xls fayl yuklang.")

        loop = asyncio.get_event_loop()
        try:
            pdf_bytes, info = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _do_xlsx2pdf, data),
                timeout=60.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Konversiya 60 soniyadan oshdi. Excel fayl juda katta.")

        logger.info(f"xlsx2pdf: {len(data)//1024}KB → {len(pdf_bytes)//1024}KB, {time.time()-t0:.1f}s")
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=spreadsheet.pdf",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"xlsx2pdf xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"xlsx2pdf: {type(e).__name__}: {str(e)[:160]}")

# ─── PPTX Compress (ZIP approach) ────────────────────────────────────────────

_PPTX_IMG_EXTS = {"jpg", "jpeg", "png", "bmp", "webp", "gif", "tiff"}
_PPTX_MAX_DIM  = 1920


# ✅ YANGILANDI: python-pptx — rasm metadata va munosabatlarini saqlab siqish
def _do_compresspptx(data: bytes) -> tuple:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from PIL import Image

    orig_size = len(data)
    compressed = 0

    prs = Presentation(io.BytesIO(data))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                img_blob = shape.image.blob
                img_ext = shape.image.ext.lower()
                img = Image.open(io.BytesIO(img_blob))
                if max(img.width, img.height) > _PPTX_MAX_DIM:
                    ratio = _PPTX_MAX_DIM / max(img.width, img.height)
                    img = img.resize(
                        (int(img.width * ratio), int(img.height * ratio)),
                        Image.LANCZOS,
                    )
                if len(img_blob) > 2 * 1024 * 1024:
                    quality = 65
                elif len(img_blob) > 512 * 1024:
                    quality = 72
                else:
                    quality = 82
                buf = io.BytesIO()
                if img.mode == "RGBA" or img_ext in ("png", "gif"):
                    img.save(buf, format="PNG", optimize=True)
                else:
                    img.convert("RGB").save(buf, format="JPEG", quality=quality, optimize=True)
                candidate = buf.getvalue()
                if len(candidate) < len(img_blob):
                    blip = shape.element.blipFill.blip
                    rId = blip.get(qn("r:embed"))
                    image_part = shape.part.related_parts[rId]
                    image_part._blob = candidate
                    compressed += 1
            except Exception:
                continue

    out_buf = io.BytesIO()
    prs.save(out_buf)
    out_bytes = out_buf.getvalue()
    if len(out_bytes) >= orig_size:
        out_bytes = data
    saved = max(0, round((1 - len(out_bytes) / orig_size) * 100))
    info = f"{compressed} ta rasm siqildi, {saved}% kichiklashdi"
    return out_bytes, info, saved


@app.post("/api/compresspptx")
@limiter.limit("10/minute")
async def compress_pptx(request: Request, file: UploadFile = File(...)):
    try:
        data = await file.read()
        check_size(data, "/api/compresspptx")
        loop = asyncio.get_event_loop()
        try:
            out_bytes, info, saved = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _do_compresspptx, data),
                timeout=45.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="PPTX compression timed out")
        logger.info(f"compresspptx: {info}")
        return Response(
            content=out_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": "attachment; filename=compressed.pptx",
                "X-Info": safe_header(info),
                "X-Saved-Percent": str(saved),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"compresspptx xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"compresspptx: {type(e).__name__}: {str(e)[:160]}")

# ─── Image compress ───────────────────────────────────────────────────────────

_IMGCOMPRESS_MAX_DIM = 4000
_IMGCOMPRESS_ALLOWED = {b"\xff\xd8\xff", b"\x89PNG", b"RIFF", b"WEBP"}  # JPEG, PNG, WebP


def _do_imgcompress(data: bytes, output_format: str = "jpeg", user_quality: int = 0) -> tuple:
    orig_size = len(data)

    # ✅ YANGILANDI: pyvips birlamchi (10x tez, 4x kam RAM, streaming pipeline)
    try:
        import pyvips
        vimg = pyvips.Image.new_from_buffer(data, "")
        w, h = vimg.width, vimg.height
        if max(w, h) > _IMGCOMPRESS_MAX_DIM:
            scale = _IMGCOMPRESS_MAX_DIM / max(w, h)
            vimg = vimg.resize(scale)
        if user_quality:
            quality = max(30, min(95, user_quality))
        elif orig_size > 3 * 1024 * 1024:
            quality = 65
        elif orig_size > 1024 * 1024:
            quality = 72
        else:
            quality = 82
        has_alpha = vimg.hasalpha()
        fmt = output_format
        if has_alpha and fmt == "jpeg":
            fmt = "webp"
        if fmt == "webp":
            out_bytes = vimg.webpsave_buffer(Q=quality)
            media_type, ext = "image/webp", "webp"
        elif fmt == "png":
            out_bytes = vimg.pngsave_buffer()
            media_type, ext = "image/png", "png"
        else:
            if has_alpha:
                vimg = vimg.flatten()
            out_bytes = vimg.jpegsave_buffer(Q=quality, optimize_coding=True)
            media_type, ext = "image/jpeg", "jpg"
        if len(out_bytes) >= orig_size:
            out_bytes = data
            quality = 0
            if data[:3] == b"\xff\xd8\xff":
                media_type, ext = "image/jpeg", "jpg"
            elif data[:8] == b"\x89PNG\r\n\x1a\n":
                media_type, ext = "image/png", "png"
            elif data[:4] == b"RIFF" and data[8:12] == b"WEBP":
                media_type, ext = "image/webp", "webp"
            else:
                media_type, ext = "image/jpeg", "jpg"
        saved = max(0, round((1 - len(out_bytes) / orig_size) * 100))
        info = f"{w}x{h} -> {vimg.width}x{vimg.height}, {fmt}, saved {saved}%"
        return out_bytes, info, saved, media_type, ext
    except ImportError:
        pass
    except Exception as _vips_err:
        logger.debug(f"pyvips xato, Pillow fallback: {_vips_err}")

    # Pillow fallback
    from PIL import Image, ImageOps

    img = Image.open(io.BytesIO(data))
    _orig_fmt = (img.format or "jpeg").lower()
    img = ImageOps.exif_transpose(img)

    w, h = img.width, img.height
    ratio = min(1.0, _IMGCOMPRESS_MAX_DIM / max(w, h))
    if ratio < 1.0:
        img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)

    has_alpha = img.mode in ("RGBA", "LA", "P")
    # Force webp for alpha images if jpeg requested (JPEG can't hold alpha)
    if has_alpha and output_format == "jpeg":
        output_format = "webp"

    # Adaptive quality based on source size
    if user_quality:
        quality = max(30, min(95, user_quality))
    elif orig_size > 3 * 1024 * 1024:
        quality = 65
    elif orig_size > 1024 * 1024:
        quality = 72
    else:
        quality = 82

    out = io.BytesIO()
    if output_format == "webp":
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGBA" if has_alpha else "RGB")
        img.save(out, format="WEBP", quality=quality, method=4)
        media_type, ext = "image/webp", "webp"
    elif output_format == "png":
        img.save(out, format="PNG", optimize=True, compress_level=9)
        media_type, ext = "image/png", "png"
    else:
        if img.mode != "RGB":
            img = img.convert("RGB")
        img.save(out, format="JPEG", quality=quality, optimize=True)
        media_type, ext = "image/jpeg", "jpg"

    out_bytes = out.getvalue()
    if len(out_bytes) >= orig_size:
        out_bytes = data
        quality   = 0
        # Revert media_type/ext to match the original file format
        _fmt_map = {"jpeg": ("image/jpeg", "jpg"), "png": ("image/png", "png"),
                    "webp": ("image/webp", "webp"), "gif": ("image/gif", "gif")}
        media_type, ext = _fmt_map.get(_orig_fmt, ("image/jpeg", "jpg"))

    saved = max(0, round((1 - len(out_bytes) / orig_size) * 100))
    info  = f"{w}×{h} → {img.width}×{img.height}, {output_format}, saved {saved}%"
    return out_bytes, info, saved, media_type, ext


@app.post("/api/imgcompress")
@limiter.limit("20/minute")
async def img_compress(
    request: Request,
    file: UploadFile = File(...),
    output_format: str = Form("jpeg"),
    quality: int = Form(0),
):
    try:
        data = await file.read()
        check_size(data, "/api/imgcompress")
        if len(data) < 4:
            raise HTTPException(status_code=422, detail="Rasm fayli ochib bo'lmadi. JPG, PNG yoki WebP yuklang.")
        # Quick magic-byte check: JPEG (FF D8), PNG (89 50 4E 47), WebP (52 49 46 46), GIF (47 49 46)
        sig = data[:4]
        is_image = (
            sig[:2] == b'\xff\xd8' or   # JPEG
            sig == b'\x89PNG' or        # PNG
            sig == b'RIFF' or           # WebP
            sig[:3] == b'GIF'           # GIF
        )
        if not is_image:
            raise HTTPException(status_code=422, detail="Rasm fayli emas. JPG, PNG yoki WebP yuklang.")
        output_format = output_format if output_format in ("jpeg", "png", "webp") else "jpeg"
        loop = asyncio.get_event_loop()
        try:
            out_bytes, info, saved, media_type, ext = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(_do_imgcompress, data, output_format, quality),
                ),
                timeout=30.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="Image compression timed out")
        return Response(
            content=out_bytes,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename=compressed.{ext}",
                "X-Info": safe_header(info),
                "X-Saved-Percent": str(saved),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"imgcompress xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"imgcompress: {type(e).__name__}: {str(e)[:160]}")

# ─── Background removal ───────────────────────────────────────────────────────

@app.post("/api/bgremove")
@limiter.limit("5/minute")
async def bgremove(request: Request):
    import base64
    from rembg import remove
    t0 = time.time()
    try:
        ct       = request.headers.get("content-type", "")
        bg_color = None
        if "multipart" in ct:
            form     = await request.form()
            f        = form.get("file")
            data     = await f.read()
            bg_color = (form.get("bg_color") or "").strip().lstrip("#") or None
        else:
            body     = await request.json()
            data     = base64.b64decode(body["data"])
            bg_color = (body.get("bg_color") or "").strip().lstrip("#") or None

        check_size(data, "/api/bgremove")

        session = get_rembg_session()
        loop    = asyncio.get_event_loop()
        try:
            result_png = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(remove, data, session=session),
                ),
                timeout=60.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Fon olib tashlash 60 soniyadan oshdi. Kichikroq rasm yuklang.")

        # Optional: composite over a solid background color
        _bg = bg_color.lstrip("#") if bg_color else ""
        if len(_bg) == 3:
            _bg = "".join(c * 2 for c in _bg)  # expand "fff" → "ffffff"
        if _bg and len(_bg) >= 6:
            from PIL import Image as _PILImg
            try:
                r = int(_bg[0:2], 16)
                g = int(_bg[2:4], 16)
                b = int(_bg[4:6], 16)
                fg  = _PILImg.open(io.BytesIO(result_png)).convert("RGBA")
                bg  = _PILImg.new("RGBA", fg.size, (r, g, b, 255))
                bg.paste(fg, mask=fg.split()[3])
                out_buf = io.BytesIO()
                bg.convert("RGB").save(out_buf, format="JPEG", quality=95, optimize=True)
                result_png = out_buf.getvalue()
                media_type, fname = "image/jpeg", "no-bg.jpg"
            except Exception:
                media_type, fname = "image/png", "no-bg.png"
        else:
            media_type, fname = "image/png", "no-bg.png"

        logger.info(f"bgremove: {len(data)//1024}KB → {len(result_png)//1024}KB, {time.time()-t0:.1f}s")
        return Response(
            content=result_png,
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={fname}"},
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"bgremove xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"bgremove: {type(e).__name__}: {str(e)[:160]}")

# ─── OCR helpers ──────────────────────────────────────────────────────────────

def _preprocess_for_ocr(img):
    """
    Prepare image for Tesseract:
      1. EXIF rotation fix
      2. Grayscale
      3. Scale: minimum 1800px long edge (Tesseract needs ~300 DPI)
      4. Cap at 4000px (prevent OOM + slow OCR)
      5. Auto-contrast (stretch histogram)
      6. Moderate contrast boost + sharpen
    Keeps preprocessing moderate — LSTM engine works well without
    aggressive binarization.
    """
    from PIL import Image, ImageFilter, ImageEnhance, ImageOps

    img = ImageOps.exif_transpose(img)
    gray = img.convert('L')
    w, h = gray.size

    # Scale up if image is too small for accurate OCR
    MIN_EDGE = 1800
    if max(w, h) < MIN_EDGE:
        scale = MIN_EDGE / max(w, h)
        gray = gray.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        w, h = gray.size

    # Cap resolution to prevent memory/speed issues
    MAX_EDGE = 4000
    if max(w, h) > MAX_EDGE:
        scale = MAX_EDGE / max(w, h)
        gray = gray.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

    # Stretch histogram (handles under/over-exposed photos)
    gray = ImageOps.autocontrast(gray, cutoff=2)
    # Moderate contrast boost
    gray = ImageEnhance.Contrast(gray).enhance(1.4)
    # Light sharpening (helps with slightly blurry scans)
    gray = gray.filter(ImageFilter.SHARPEN)

    return gray


def _clean_ocr_text(text: str) -> str:
    """Remove Tesseract noise: stray chars, excessive blank lines."""
    # Collapse 3+ consecutive blank lines → 2
    text = re.sub(r'\n{3,}', '\n\n', text)
    lines = text.split('\n')
    result = []
    for line in lines:
        s = line.strip()
        if s and len(s) >= 2:          # keep lines with ≥2 real chars
            result.append(s)
        elif not s and result and result[-1]:
            result.append('')          # keep at most one blank separator
    # Strip leading/trailing blanks
    while result and not result[0]:
        result.pop(0)
    while result and not result[-1]:
        result.pop()
    return '\n'.join(result)


# ✅ YANGILANDI: PaddleOCR ikki sessiya — en (Latin) va cyrillic (Rus/Uzbek-kirill)
_paddle_ocr_cache: dict = {}

def _get_paddle_ocr(lang: str = "en"):
    """lang: 'en' (Latin/Uzbek-latin) | 'cyrillic' (Rus/Uzbek-kirill)"""
    if lang in _paddle_ocr_cache:
        return _paddle_ocr_cache[lang]
    try:
        from paddleocr import PaddleOCR
        _paddle_ocr_cache[lang] = PaddleOCR(use_angle_cls=True, lang=lang, show_log=False)
    except (ImportError, Exception) as e:
        logger.debug(f"paddleocr {lang} sessiya yaratilmadi: {e}")
        _paddle_ocr_cache[lang] = None
    return _paddle_ocr_cache[lang]


def _paddle_run(ocr_obj, arr) -> str:
    if ocr_obj is None:
        return ""
    result = ocr_obj.ocr(arr, cls=True)
    lines = []
    for block in (result or []):
        for line in (block or []):
            if line and len(line) >= 2 and line[1]:
                lines.append(line[1][0])
    return "\n".join(lines)


def _do_ocr_paddle(img_bytes: bytes) -> str:
    """Run both Latin and Cyrillic models, return longer/better result."""
    import numpy as np
    from PIL import Image as _PIL
    en_ocr  = _get_paddle_ocr("en")
    cyr_ocr = _get_paddle_ocr("cyrillic")
    if en_ocr is None and cyr_ocr is None:
        raise ImportError("paddleocr not installed")
    img = _PIL.open(io.BytesIO(img_bytes)).convert("RGB")
    arr = np.array(img)
    en_text  = _paddle_run(en_ocr,  arr)
    cyr_text = _paddle_run(cyr_ocr, arr)
    # Cyrillic mavjudligini tekshir
    if re.search(r'[а-яёўқғҳА-ЯЁ]', cyr_text):
        return cyr_text if len(cyr_text) > len(en_text) * 0.5 else en_text
    return en_text if len(en_text) >= len(cyr_text) else cyr_text


def _do_ocr(data: bytes, is_pdf: bool) -> str:
    """
    Full OCR pipeline — runs in _converter_pool (CPU-bound).
    Primary: PaddleOCR. Fallback: Tesseract LSTM.

    PDF path:
      • Has extractable text → direct fitz.get_text() (fast, accurate)
      • Scanned PDF         → render 200 DPI + PaddleOCR → Tesseract

    Image path:
      • PaddleOCR → Tesseract fallback
    """
    import pytesseract
    from PIL import Image

    LANGS  = 'rus+eng+uzb'
    CONFIG = '--oem 3 --psm 6'

    if is_pdf:
        import fitz
        doc = fitz.open(stream=data, filetype='pdf')
        if doc.is_encrypted:
            doc.close()
            return 'PDF parol bilan himoyalangan. Avval parolini oching.'
        total_pages = doc.page_count
        if total_pages == 0:
            doc.close()
            return 'PDF bo\'sh (0 sahifa).'

        sample = ''.join(doc[i].get_text() for i in range(min(3, total_pages)))
        if len(sample.strip()) > 50:
            MAX_P = 15
            parts = []
            for i in range(min(total_pages, MAX_P)):
                t = doc[i].get_text().strip()
                if t:
                    if total_pages > 1:
                        parts.append(f'── Sahifa {i + 1} ──')
                    parts.append(t)
            doc.close()
            result = _clean_ocr_text('\n\n'.join(parts))
            if total_pages > MAX_P:
                result += f'\n\n⚠️ Faqat {MAX_P} sahifa o\'qildi (PDF jami {total_pages} sahifa).'
            return result or 'Matn topilmadi.'

        zoom  = 200 / 72
        mat   = fitz.Matrix(zoom, zoom)
        MAX_P = 8
        parts = []
        for i in range(min(total_pages, MAX_P)):
            pix      = doc[i].get_pixmap(matrix=mat, alpha=False)
            img      = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
            del pix
            img_bytes = io.BytesIO()
            img.save(img_bytes, format="PNG")
            img_bytes = img_bytes.getvalue()
            # PaddleOCR urinish
            try:
                text = _do_ocr_paddle(img_bytes)
            except Exception:
                img = _preprocess_for_ocr(img)
                text = pytesseract.image_to_string(img, lang=LANGS, config=CONFIG)
            cleaned = _clean_ocr_text(text)
            if cleaned:
                if total_pages > 1:
                    parts.append(f'── Sahifa {i + 1} ──')
                parts.append(cleaned)
        doc.close()
        result = '\n\n'.join(parts)
        if total_pages > MAX_P:
            result += (f'\n\n⚠️ Faqat {MAX_P} sahifa OCR qilindi '
                       f'(PDF jami {total_pages} sahifali).')
        return result or 'Matn topilmadi.'

    # Image: PaddleOCR → Tesseract fallback
    try:
        text = _do_ocr_paddle(data)
        if text.strip():
            return _clean_ocr_text(text) or 'Matn topilmadi.'
    except Exception:
        pass
    img  = Image.open(io.BytesIO(data))
    img  = _preprocess_for_ocr(img)
    text = pytesseract.image_to_string(img, lang=LANGS, config=CONFIG)
    return _clean_ocr_text(text) or 'Matn topilmadi.'


# ─── OCR ──────────────────────────────────────────────────────────────────────

@app.post("/api/ocr")
@limiter.limit("15/minute")
async def ocr(request: Request):
    t0 = time.time()
    try:
        import base64
        ct = request.headers.get("content-type", "")
        if "multipart" in ct:
            form  = await request.form()
            f     = form.get("file")
            fname = (f.filename or "").lower()
            data  = await f.read()
        else:
            body  = await request.json()
            data  = base64.b64decode(body["data"])
            fname = body.get("filename", "").lower()

        check_size(data, "/api/ocr")

        # Detect file type: filename OR magic bytes (%PDF)
        is_pdf = fname.endswith('.pdf') or data[:4] == b'%PDF'

        loop = asyncio.get_event_loop()
        try:
            text = await asyncio.wait_for(
                loop.run_in_executor(_converter_pool, _do_ocr, data, is_pdf),
                timeout=90.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="OCR 90 soniyadan oshdi. Rasmni kichiklashtiring yoki PDF'ni qisqartiring.")

        logger.info(f"ocr: {len(data)//1024}KB {'pdf' if is_pdf else 'img'} → {len(text)} belgi, {time.time()-t0:.1f}s")
        return JSONResponse({"text": text})

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"ocr xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"ocr: {type(e).__name__}: {str(e)[:160]}")

# ─── Translit ─────────────────────────────────────────────────────────────────

# Latin (Uzbek) → Cyrillic
_LTR_MAP = {
    # digraphs first (checked before single chars)
    "sh": "ш", "ch": "ч",
    "yo": "ё", "ye": "е", "yu": "ю", "ya": "я",
    "ts": "ц",
    # o' variants: apostrophe (U+0027), modifier-turned-comma (U+02BB),
    # modifier-apostrophe (U+02BC), right-single-quote (U+2019)
    "o’": "ў", "oʻ": "ў", "oʼ": "ў", "o'": "ў",
    # g' variants
    "g’": "ғ", "gʻ": "ғ", "gʼ": "ғ", "g'": "ғ",
    # single letters
    "a": "а", "b": "б", "v": "в", "d": "д", "e": "е",
    "f": "ф", "g": "г", "h": "ҳ", "i": "и", "j": "ж",
    "k": "к", "l": "л", "m": "м", "n": "н", "o": "о",
    "p": "п", "q": "қ", "r": "р", "s": "с", "t": "т",
    "u": "у", "x": "х", "y": "й", "z": "з",
    # standalone apostrophe variants → ъ
    "'": "ъ", "ʻ": "ъ", "ʼ": "ъ", "’": "ъ",
}
# ordered: longer keys first so digraphs are matched before single chars
_LTR_MULTI = [k for k in _LTR_MAP if len(k) == 2]

# Cyrillic → Latin (Uzbek)
_RTL_MAP = {
    "ш": "sh", "ч": "ch",
    "ё": "yo", "ю": "yu", "я": "ya",
    "ц": "ts", "ў": "o'", "ғ": "g'",
    "а": "a",  "б": "b",  "в": "v",  "г": "g",  "д": "d",
    "е": "e",  "э": "e",  "ж": "j",  "з": "z",
    "и": "i",  "й": "y",  "к": "k",  "л": "l",  "м": "m",
    "н": "n",  "о": "o",  "п": "p",  "р": "r",  "с": "s",
    "т": "t",  "у": "u",  "ф": "f",  "х": "x",
    "қ": "q",  "ҳ": "h",  "ъ": "'",  "ь": "",
}


def _translit_ltr(text: str) -> str:
    """Uzbek Latin → Cyrillic."""
    out = []
    i = 0
    n = len(text)
    while i < n:
        matched = False
        for key in _LTR_MULTI:
            kl = len(key)
            if text[i:i + kl].lower() == key:
                cyr = _LTR_MAP[key]
                out.append(cyr.upper() if text[i].isupper() else cyr)
                i += kl
                matched = True
                break
        if not matched:
            c = text[i]
            m = _LTR_MAP.get(c.lower(), c)
            out.append(m.upper() if c.isupper() else m)
            i += 1
    return "".join(out)


def _translit_rtl(text: str) -> str:
    """Uzbek Cyrillic → Latin.

    Handles three case patterns:
      lowercase  → lowercase  (salom → salom)
      Title-case → Title-case (Salom → Salom, Sh → Sh)
      ALL-CAPS   → ALL-CAPS   (SALOM → SALOM, SH → SH)
    """
    out = []
    chars = list(text)
    n = len(chars)
    for idx, c in enumerate(chars):
        lc = c.lower()
        m = _RTL_MAP.get(lc, c)
        if not m:           # soft sign → skip
            continue
        if c.isupper() and lc != c:
            # Look ahead: if next alphabetic char is also uppercase → ALL-CAPS
            next_alpha = next(
                (chars[j] for j in range(idx + 1, min(idx + 4, n))
                 if chars[j].isalpha()),
                None,
            )
            if next_alpha and next_alpha.isupper():
                out.append(m.upper())          # ALL-CAPS: SH, CH, YO …
            else:
                out.append(m[0].upper() + m[1:])  # Title: Sh, Ch, Yo …
        else:
            out.append(m)
    return "".join(out)


@app.post("/api/translit")
@limiter.limit("60/minute")
async def translit(request: Request):
    body = await request.json()
    text = str(body.get("text", "")).strip()
    if not text:
        return JSONResponse({"result": ""})
    if len(text) > 50_000:
        text = text[:50_000]
    # Detect by Uzbek Cyrillic letters (including ў қ ғ ҳ)
    is_cyr = bool(re.search(r'[а-яёА-ЯЁўқғҳ]', text))
    result = _translit_rtl(text) if is_cyr else _translit_ltr(text)
    return JSONResponse({"result": result})

# ─── Readtime ─────────────────────────────────────────────────────────────────

@app.post("/api/readtime")
@limiter.limit("60/minute")
async def readtime(request: Request):
    body = await request.json()
    text = (body.get("text") or "").strip()
    if not text:
        return JSONResponse({"result": "❌ Matn kiriting"})
    words = len(text.split())
    chars = len(text.replace(" ", ""))
    sentences = len(re.findall(r'[.!?]+', text)) or 1
    paragraphs = len([p for p in text.split("\n") if p.strip()])
    # O'zbek o'qish tezligi: ~160 so'z/daqiqa
    wpm = 160
    mins = max(1, round(words / wpm))
    return JSONResponse({"result": (
        f"📖 {words} so'z · {chars} belgi\n"
        f"📝 {sentences} gap · {paragraphs} paragraf\n"
        f"⏱ O'qish vaqti: ~{mins} daqiqa\n"
        f"({wpm} so'z/daqiqa hisobida)"
    )})

# ─── AI Summarize ─────────────────────────────────────────────────────────────
# Anthropic API bilan AI xulosalash. ANTHROPIC_API_KEY shart.
# Fallback: extractive (eng muhim jumlalarni tanlash).

ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")


def _extractive_summarize(text: str, max_sentences: int = 4) -> str:
    """Extractive xulosalash — eng tez-tez uchragan so'zlarga ega jumlalar."""
    from collections import Counter
    sents = re.split(r'(?<=[.!?])\s+', text.strip())
    if len(sents) <= max_sentences:
        return text.strip()
    # Stop-words minimal — barcha tillarda
    stop = {"va","bilan","uchun","uchun","ham","bu","shu","u","men","sen","biz","siz","ular",
            "in","of","the","a","an","to","is","and","or","for","on","at","by","with","as","it",
            "и","в","на","с","по","к","из","о","за","но","что","это","он","она","они","мы"}
    words = re.findall(r"\w+", text.lower())
    freq = Counter(w for w in words if w not in stop and len(w) > 2)
    if not freq:
        return ".".join(sents[:max_sentences])
    scored = []
    for i, s in enumerate(sents):
        sw = re.findall(r"\w+", s.lower())
        score = sum(freq.get(w, 0) for w in sw) / max(1, len(sw))
        scored.append((score, i, s))
    top = sorted(scored, reverse=True)[:max_sentences]
    return " ".join(s for _, _, s in sorted(top, key=lambda x: x[1]))


async def _claude_summarize(text: str, lang: str = "uz") -> str:
    """Claude API bilan AI xulosalash."""
    if not ANTHROPIC_API_KEY:
        raise RuntimeError("ANTHROPIC_API_KEY o'rnatilmagan")
    lang_label = {"uz": "o'zbek tilida", "ru": "по-русски", "en": "in English"}.get(lang, "o'zbek tilida")
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": ANTHROPIC_API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": "claude-haiku-4-5-20251001",
                "max_tokens": 600,
                "messages": [{
                    "role": "user",
                    "content": (
                        f"Quyidagi matnni {lang_label} 3-5 jumlada qisqacha xulosalab ber. "
                        f"Asosiy fikrlarni saqla, lekin keraksiz tafsilotlardan voz kech. "
                        f"Faqat xulosani qaytar, izoh berma.\n\n{text[:8000]}"
                    ),
                }],
            },
        )
    if r.status_code != 200:
        raise RuntimeError(f"Claude API xato: {r.status_code}")
    data = r.json()
    return (data.get("content") or [{}])[0].get("text", "").strip()


@app.post("/api/summarize")
@limiter.limit("10/minute")
async def summarize(request: Request):
    body = await request.json()
    text = (body.get("text") or "").strip()
    lang = (body.get("lang") or "uz").lower()[:2]
    if not text:
        return JSONResponse({"result": "❌ Matn kiriting"})
    if len(text) < 200:
        return JSONResponse({"result": "❌ Matn juda qisqa (kamida 200 belgi)"})
    text = text[:10000]

    # Try Claude first if API key set, otherwise extractive
    if ANTHROPIC_API_KEY:
        try:
            ai_result = await asyncio.wait_for(_claude_summarize(text, lang), timeout=25.0)
            if ai_result:
                return JSONResponse({"result": f"🤖 AI xulosa:\n\n{ai_result}"})
        except Exception as e:
            logger.warning(f"summarize claude xato: {e}")

    # Fallback: extractive
    try:
        loop = asyncio.get_event_loop()
        ext = await loop.run_in_executor(None, _extractive_summarize, text)
        return JSONResponse({"result": f"📝 Qisqacha:\n\n{ext}"})
    except Exception as e:
        logger.error(f"summarize xato: {e}")
        return JSONResponse({"result": f"❌ Xulosalash amalga oshmadi"})


# ─── Deadline ─────────────────────────────────────────────────────────────────

@app.post("/api/deadline")
@limiter.limit("60/minute")
async def deadline(request: Request):
    body = await request.json()
    text = (body.get("text") or "").strip()
    if not text:
        return JSONResponse({"result": "❌ Sana kiriting: 31.12.2025, kelasi juma, 3 kun keyin"})

    d = None
    # Try dateparser first — understands "kelasi juma", "next Monday", "3 days later"
    try:
        import dateparser
        d = dateparser.parse(
            text,
            languages=["uz", "ru", "en"],
            settings={"PREFER_DATES_FROM": "future", "RETURN_AS_TIMEZONE_AWARE": False},
        )
    except Exception:
        pass

    # Regex fallback for explicit numeric formats: 31.12.2025 / 2025-12-31
    if d is None:
        m = re.search(r'(\d{1,2})[.\-/](\d{1,2})[.\-/](\d{2,4})', text) or \
            re.search(r'(\d{4})[.\-/](\d{1,2})[.\-/](\d{1,2})', text)
        if m:
            try:
                g = m.groups()
                if len(g[0]) == 4:
                    d = datetime(int(g[0]), int(g[1]), int(g[2]))
                else:
                    yr = int(g[2])
                    if yr < 100: yr += 2000
                    d = datetime(yr, int(g[1]), int(g[0]))
            except ValueError:
                pass

    if d is None:
        return JSONResponse({"result": (
            "❌ Sana formatini aniqlab bo'lmadi.\n"
            "Misollar: 31.12.2025 · 2025-12-31 · kelasi juma · 3 kun keyin · next Monday"
        )})

    now = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)
    d0  = d.replace(hour=0, minute=0, second=0, microsecond=0)
    days = (d0 - now).days
    emoji = "🔴" if days < 0 else ("🚨" if days <= 3 else ("🟡" if days <= 7 else "🟢"))
    msg = f"{abs(days)} kun oldin o'tdi!" if days < 0 else ("BUGUN!" if days == 0 else f"{days} kun qoldi")
    return JSONResponse({"result": f"📅 {d.strftime('%d.%m.%Y')}\n{emoji} {msg}"})

# ─── Stats ────────────────────────────────────────────────────────────────────

@app.post("/api/stats")
@limiter.limit("60/minute")
async def stats(request: Request):
    body = await request.json()
    nums = [float(x) for x in re.findall(r'-?\d+(?:\.\d+)?', body.get("text", ""))]
    if not nums:
        return JSONResponse({"result": "❌ Raqamlar kiriting. Masalan: 4 7 2 9 1 5"})
    n = len(nums)
    total = sum(nums)
    mean  = total / n
    s     = sorted(nums)
    median = s[n // 2] if n % 2 else (s[n // 2 - 1] + s[n // 2]) / 2
    var   = sum((x - mean) ** 2 for x in nums) / n
    std   = var ** 0.5

    # Quartiles (linear interpolation)
    def _percentile(data, p):
        idx = (len(data) - 1) * p / 100
        lo, hi = int(idx), min(int(idx) + 1, len(data) - 1)
        return data[lo] + (data[hi] - data[lo]) * (idx - lo)

    q1  = _percentile(s, 25)
    q3  = _percentile(s, 75)
    iqr = q3 - q1

    # Mode (most frequent value; show only if appears >1 time)
    from collections import Counter
    freq = Counter(nums)
    top_val, top_cnt = freq.most_common(1)[0]
    mode_str = f"{top_val:g} ({top_cnt} marta)" if top_cnt > 1 else "—"

    # scipy: skewness, kurtosis, normality test
    extra = ""
    try:
        from scipy import stats as _sp
        skew = _sp.skew(nums)
        kurt = _sp.kurtosis(nums)
        extra += f"\n📐 Asimmetriya: {skew:.4f}"
        extra += f"\n⛰️  Kurtosis: {kurt:.4f}"
        if n >= 8:
            _, norm_p = _sp.normaltest(nums)
            norm_label = "Ha ✓" if norm_p > 0.05 else f"Yo'q (p={norm_p:.3f})"
            extra += f"\n🔔 Normal taqsimot: {norm_label}"
    except Exception:
        pass

    return JSONResponse({"result": (
        f"📊 n = {n}\n"
        f"∑  Yig'indi: {total:g}\n"
        f"x̄  O'rtacha: {mean:.4f}\n"
        f"M  Mediana: {median:g}\n"
        f"Mo Moda: {mode_str}\n"
        f"σ² Dispersiya: {var:.4f}\n"
        f"σ  Standart og'ish: {std:.4f}\n"
        f"Q1 / Q3: {q1:g} / {q3:g}\n"
        f"IQR: {iqr:g}\n"
        f"Min: {s[0]:g}  Max: {s[-1]:g}"
        + extra
    )})

# ─── QR Code ──────────────────────────────────────────────────────────────────

# EduBot accent palette
_QR_FILL  = (79, 70, 229)   # indigo-600
_QR_BACK  = (238, 242, 255) # indigo-50


def _do_qr(text: str, size: int = 400, ec_level: str = "M", fmt: str = "png") -> tuple:
    import segno
    from PIL import Image

    ec_map = {"L": "l", "M": "m", "Q": "q", "H": "h"}
    ec = ec_map.get(ec_level.upper(), "m")

    qr = segno.make(text, error=ec, micro=False)

    if fmt == "svg":
        buf = io.BytesIO()
        qr.save(buf, kind="svg", scale=4, dark="#0f172a", light="#f8f7ff", border=4)
        return buf.getvalue(), "image/svg+xml", "qrcode.svg"

    # Compute scale so output is close to requested size
    modules = qr.symbol_size()[0]
    scale = max(1, size // modules)
    buf = io.BytesIO()
    qr.save(buf, kind="png", scale=scale,
            dark="#4f46e5", light="#eef2ff", border=4)
    img = Image.open(buf)
    if img.width != size:
        img = img.resize((size, size), Image.NEAREST)
    out = io.BytesIO()
    img.save(out, format="PNG", optimize=True)
    return out.getvalue(), "image/png", "qrcode.png"


@app.post("/api/qr")
@limiter.limit("30/minute")
async def make_qr(request: Request):
    try:
        body     = await request.json()
        text     = (body.get("text") or "EduBot").strip() or "EduBot"
        size     = max(100, min(1000, int(body.get("size", 400))))
        ec_level = str(body.get("ec", "M")).upper()
        fmt      = str(body.get("format", "png")).lower()
        if fmt not in ("png", "svg"):
            fmt = "png"
        if len(text) > 2000:
            raise HTTPException(status_code=400, detail="Matn 2000 belgidan oshmasligi kerak")
        loop = asyncio.get_event_loop()
        try:
            result = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(_do_qr, text, size, ec_level, fmt),
                ),
                timeout=15.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="QR generation timed out")
        content, media_type, filename = result
        return Response(
            content=content,
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"},
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"qr xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"qr: {type(e).__name__}: {str(e)[:160]}")

# ─── Certificate ──────────────────────────────────────────────────────────────

_CERT_THEMES = {
    "classic": {"bg_top": (10, 8, 45),   "bg_bot": (35, 20, 75),  "accent": (251, 191, 36),  "text": (255, 255, 255), "sub": (167, 139, 250)},
    "modern":  {"bg_top": (15, 23, 42),  "bg_bot": (30, 50, 80),  "accent": (96, 165, 250),  "text": (255, 255, 255), "sub": (148, 210, 252)},
    "minimal": {"bg_top": (245, 245, 250),"bg_bot": (220, 220, 235),"accent": (79, 70, 229),  "text": (30, 30, 50),    "sub": (100, 80, 200)},
    "dark":    {"bg_top": (5, 5, 15),    "bg_bot": (20, 10, 30),  "accent": (167, 139, 250), "text": (240, 240, 255), "sub": (120, 100, 220)},
}


def _do_cert(name: str, course: str, issuer: str, theme: str = "classic") -> bytes:
    from PIL import Image, ImageDraw, ImageFont

    th     = _CERT_THEMES.get(theme, _CERT_THEMES["classic"])
    W, H   = 900, 620
    img    = Image.new("RGB", (W, H))
    draw   = ImageDraw.Draw(img)

    # Gradient background (theme top → theme bottom)
    r0, g0, b0 = th["bg_top"]
    r1, g1, b1 = th["bg_bot"]
    for i in range(H):
        t = i / H
        draw.line([(0, i), (W, i)], fill=(
            int(r0 + t * (r1 - r0)),
            int(g0 + t * (g1 - g0)),
            int(b0 + t * (b1 - b0)),
        ))

    GOLD   = th["accent"]
    GOLD2  = tuple(max(0, c - 70) for c in GOLD)
    WHITE  = th["text"]
    SILVER = (190, 190, 200)
    PURPLE = th["sub"]

    # Outer + inner border
    draw.rectangle([(16, 16), (W - 16, H - 16)], outline=GOLD,  width=3)
    draw.rectangle([(26, 26), (W - 26, H - 26)], outline=GOLD2, width=1)

    # Corner ornaments (simple cross-hatch)
    for cx_, cy_ in [(42, 42), (W - 42, 42), (42, H - 42), (W - 42, H - 42)]:
        draw.ellipse([(cx_ - 8, cy_ - 8), (cx_ + 8, cy_ + 8)], outline=GOLD, width=2)
        draw.line([(cx_ - 14, cy_), (cx_ + 14, cy_)], fill=GOLD2, width=1)
        draw.line([(cx_, cy_ - 14), (cx_, cy_ + 14)], fill=GOLD2, width=1)

    # Decorative horizontal rules
    draw.line([(60, 90), (W - 60, 90)],   fill=GOLD2, width=1)
    draw.line([(60, 92), (W - 60, 92)],   fill=GOLD,  width=1)
    draw.line([(60, H - 90), (W - 60, H - 90)], fill=GOLD,  width=1)
    draw.line([(60, H - 88), (W - 60, H - 88)], fill=GOLD2, width=1)

    def load_fonts():
        try:
            return (
                ImageFont.truetype(FONT_BOLD,    48),  # title
                ImageFont.truetype(FONT_BOLD,    44),  # name
                ImageFont.truetype(FONT_REGULAR, 22),  # course
                ImageFont.truetype(FONT_REGULAR, 15),  # sub
                ImageFont.truetype(FONT_REGULAR, 13),  # small
            )
        except Exception:
            d = ImageFont.load_default()
            return d, d, d, d, d

    fn_title, fn_name, fn_course, fn_sub, fn_sm = load_fonts()

    def center_x(text, font):
        try:
            bb = draw.textbbox((0, 0), text, font=font)
            return (W - (bb[2] - bb[0])) // 2
        except Exception:
            return W // 4

    # Title
    draw.text((center_x("SERTIFIKAT", fn_title), 108), "SERTIFIKAT",
              font=fn_title, fill=GOLD)

    # Subtitle
    sub = "ushbu kurs muvaffaqiyatli yakunlanganligi tasdiqlanganligi uchun beriladi"
    draw.text((center_x(sub, fn_sm), 168), sub, font=fn_sm, fill=SILVER)

    # Recipient name (Unicode-safe — DejaVu covers Latin + Cyrillic)
    draw.text((center_x(name, fn_name), 230), name, font=fn_name, fill=WHITE)

    # Divider below name
    draw.line([(W // 2 - 200, 295), (W // 2 + 200, 295)], fill=GOLD2, width=1)

    # Course label + value
    label = "kurs:"
    draw.text((center_x(label, fn_sm), 308), label, font=fn_sm, fill=SILVER)
    draw.text((center_x(course, fn_course), 328), course, font=fn_course, fill=PURPLE)

    # Seal circle (watermark-style)
    sc, sr = W // 2, H - 140
    draw.ellipse([(sc - 46, sr - 46), (sc + 46, sr + 46)], outline=GOLD2, width=2)
    draw.ellipse([(sc - 38, sr - 38), (sc + 38, sr + 38)], outline=GOLD,  width=1)
    seal_text = "✓"
    draw.text((center_x(seal_text, fn_sub), sr - 14), seal_text, font=fn_sub, fill=GOLD)

    # Issuer + date columns
    date_str = datetime.now().strftime("%d.%m.%Y")
    draw.text((70,  H - 75), issuer,   font=fn_sm, fill=SILVER)
    draw.text((70,  H - 58), "Mudir",  font=fn_sm, fill=SILVER)
    draw.line([(70, H - 50), (200, H - 50)], fill=GOLD2, width=1)

    draw.text((W - 200, H - 75), date_str,    font=fn_sm, fill=SILVER)
    draw.text((W - 200, H - 58), "Sana",      font=fn_sm, fill=SILVER)
    draw.line([(W - 200, H - 50), (W - 70, H - 50)], fill=GOLD2, width=1)

    draw.text((center_x("EduBot", fn_sm), H - 40), "EduBot",
              font=fn_sm, fill=(100, 100, 120))

    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


@app.post("/api/cert")
@limiter.limit("20/minute")
async def make_cert(request: Request):
    try:
        body   = await request.json()
        lines  = (body.get("text") or "Ism Familiya\nKurs nomi").strip().split("\n")
        name   = lines[0].strip() if lines else "Ism Familiya"
        course = lines[1].strip() if len(lines) > 1 else "Kurs nomi"
        issuer = lines[2].strip() if len(lines) > 2 else "EduBot"
        theme  = str(body.get("theme", "classic"))
        if theme not in _CERT_THEMES:
            theme = "classic"
        if not name:
            raise HTTPException(status_code=400, detail="Ism bo'sh bo'lmasligi kerak")
        loop = asyncio.get_event_loop()
        try:
            png = await asyncio.wait_for(
                loop.run_in_executor(
                    _converter_pool,
                    functools.partial(_do_cert, name, course, issuer, theme),
                ),
                timeout=20.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="Certificate generation timed out")
        return Response(
            content=png,
            media_type="image/png",
            headers={"Content-Disposition": "attachment; filename=certificate.png"},
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"cert xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"cert: {type(e).__name__}: {str(e)[:160]}")

# ─── Schedule ─────────────────────────────────────────────────────────────────

_TIME_RE = re.compile(r'\b(\d{1,2}:\d{2})\b')

def _parse_schedule(text: str) -> list:
    """Parse schedule input.
    Input format (one day per line):
        Dushanba: Matematika 8:00, Fizika 10:00
        Seshanba: Ingliz tili 9:00
    Returns: [(day, [(num, subject, time_str), ...]), ...]
    """
    result = []
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if ':' in line:
            # Split only on the first colon
            colon = line.index(':')
            day   = line[:colon].strip()
            rest  = line[colon + 1:].strip()
        else:
            result.append((line, []))
            continue

        lessons = []
        for num, part in enumerate(rest.split(','), 1):
            part = part.strip()
            if not part:
                continue
            m = _TIME_RE.search(part)
            if m:
                time_str = m.group(1)
                subject  = (part[:m.start()] + part[m.end():]).strip()
                subject  = subject or part
            else:
                subject  = part
                time_str = ""
            lessons.append((num, subject[:45], time_str))
        result.append((day[:20], lessons))
    return result

def _do_schedule(text: str) -> bytes:
    from PIL import Image, ImageDraw, ImageFont

    schedule = _parse_schedule(text)
    if not schedule:
        raise ValueError("Bo'sh jadval. Format:\nDushanba: Matematika 8:00, Fizika 10:00")

    # Flatten schedule → rows: (day_label, num_str, subject, time_str)
    flat = []
    for day, lessons in schedule:
        if not lessons:
            flat.append((day, "", "", ""))
        else:
            for idx, (num, subject, time_str) in enumerate(lessons):
                flat.append((day if idx == 0 else "", f"{num}.", subject, time_str))

    # ── Layout constants ────────────────────────────────────────────
    W         = 720
    PAD       = 14
    ROW_H     = 46
    TITLE_H   = 56
    HDR_H     = 32         # column header strip
    H         = TITLE_H + HDR_H + ROW_H * len(flat) + PAD

    # Column x-positions and widths: Day | # | Subject | Time
    COL_X = [PAD, PAD + 148, PAD + 148 + 32, PAD + 148 + 32 + 400]
    COL_W = [148, 32, 400, 100]

    # Colours
    C_BG      = (10,  10,  20)
    C_HDR_BG  = (22,  14,  48)
    C_ROW_A   = (18,  14,  38)
    C_ROW_B   = (12,  10,  26)
    C_BORDER  = (55,  42,  88)
    C_ACCENT  = (139, 92,  246)
    C_DAY     = (180, 150, 255)
    C_TEXT    = (225, 220, 255)
    C_TIME    = (251, 191, 36)
    C_NUM     = (100, 95,  155)
    C_HDR_TXT = (160, 140, 210)

    img  = Image.new("RGB", (W, H), C_BG)
    draw = ImageDraw.Draw(img)

    try:
        fn_title = ImageFont.truetype(FONT_BOLD,    22)
        fn_hdr   = ImageFont.truetype(FONT_BOLD,    13)
        fn_day   = ImageFont.truetype(FONT_BOLD,    15)
        fn_body  = ImageFont.truetype(FONT_REGULAR, 14)
        fn_time  = ImageFont.truetype(FONT_BOLD,    14)
        fn_num   = ImageFont.truetype(FONT_REGULAR, 13)
    except Exception:
        fn_title = fn_hdr = fn_day = fn_body = fn_time = fn_num = ImageFont.load_default()

    def cx(t, f):
        bb = draw.textbbox((0, 0), t, font=f)
        return (W - (bb[2] - bb[0])) // 2

    # ── Title ───────────────────────────────────────────────────────
    draw.text((cx("Dars Jadvali", fn_title), 17), "Dars Jadvali",
              font=fn_title, fill=C_ACCENT)

    # ── Column header strip ─────────────────────────────────────────
    y_hdr = TITLE_H
    draw.rectangle([(0, y_hdr), (W, y_hdr + HDR_H)], fill=C_HDR_BG)
    for col, label in enumerate(["Kun", "#", "Fan / Dars", "Vaqt"]):
        draw.text((COL_X[col] + 4, y_hdr + 9), label, font=fn_hdr, fill=C_HDR_TXT)
    draw.line([(0, y_hdr + HDR_H), (W, y_hdr + HDR_H)], fill=C_BORDER)

    # ── Data rows ───────────────────────────────────────────────────
    y0 = TITLE_H + HDR_H
    for i, (day_label, num_str, subject, time_str) in enumerate(flat):
        y  = y0 + i * ROW_H
        vy = y + (ROW_H - 16) // 2
        draw.rectangle([(0, y), (W, y + ROW_H)], fill=C_ROW_A if i % 2 == 0 else C_ROW_B)
        draw.line([(0, y + ROW_H - 1), (W, y + ROW_H - 1)], fill=C_BORDER)

        if day_label:
            draw.text((COL_X[0] + 4, vy), day_label, font=fn_day, fill=C_DAY)
        if num_str:
            draw.text((COL_X[1], vy), num_str, font=fn_num, fill=C_NUM)
        if subject:
            draw.text((COL_X[2], vy), subject, font=fn_body, fill=C_TEXT)
        if time_str:
            draw.text((COL_X[3], vy), time_str, font=fn_time, fill=C_TIME)

    # ── Vertical column dividers ────────────────────────────────────
    for col in range(1, 4):
        x_div = COL_X[col] - 2
        draw.line([(x_div, y_hdr), (x_div, H - PAD)], fill=C_BORDER)

    # ── Outer border ────────────────────────────────────────────────
    draw.rectangle([(0, 0), (W - 1, H - 1)], outline=C_ACCENT, width=2)

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

@app.post("/api/schedule")
@limiter.limit("20/minute")
async def make_schedule(request: Request):
    try:
        body = await request.json()
        text = body.get("text", "").strip()
        if not text:
            raise HTTPException(status_code=400,
                detail="Jadval kiriting.\nFormat: Dushanba: Matematika 8:00, Fizika 10:00")
        if len(text) > 3000:
            text = text[:3000]
        loop   = asyncio.get_event_loop()
        result = await loop.run_in_executor(_converter_pool, _do_schedule, text)
        return Response(
            content=result,
            media_type="image/png",
            headers={"Content-Disposition": "attachment; filename=schedule.png"},
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"schedule xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"schedule: {type(e).__name__}: {str(e)[:160]}")

# ─── Translate ────────────────────────────────────────────────────────────────

# Common 3-letter ISO → 2-letter codes Google Translate accepts
_LANG_ALIASES = {
    "uzb": "uz", "rus": "ru", "eng": "en", "tur": "tr",
    "deu": "de", "fra": "fr", "zho": "zh-CN", "jpn": "ja",
    "kor": "ko", "ara": "ar", "kaz": "kk", "kir": "ky",
    "spa": "es", "por": "pt", "ita": "it", "pol": "pl",
    "nld": "nl", "swe": "sv", "nor": "no", "dan": "da",
    "fin": "fi", "hin": "hi", "ben": "bn", "tha": "th",
    "vie": "vi", "ind": "id", "fas": "fa", "ukr": "uk",
}

# ✅ YANGILANDI: Google → argostranslate (offline) → MyMemory zanjiri
def _do_translate(query: str, lang: str) -> str:
    """Translate query → lang. Chain: Google → argostranslate (offline) → MyMemory."""
    import httpx as _httpx

    # 1. Google Translate (deep-translator)
    try:
        from deep_translator import GoogleTranslator
        result = GoogleTranslator(source="auto", target=lang).translate(query)
        if result and result.strip():
            return result.strip()
    except Exception as google_err:
        logger.warning(f"translate google xato: {google_err}")

    # 2. argostranslate — offline fallback (API key shart emas)
    try:
        import argostranslate.translate as _at
        installed = _at.get_installed_languages()
        lang_map = {l.code: l for l in installed}
        for src_code in ("en", "ru", "uz"):
            if src_code == lang:
                continue
            if src_code in lang_map and lang in lang_map:
                translation = lang_map[src_code].get_translation(lang_map[lang])
                if translation:
                    result = translation.translate(query)
                    if result and result.strip():
                        return result.strip()
    except Exception as argo_err:
        logger.debug(f"argostranslate xato: {argo_err}")

    # 3. MyMemory — oxirgi zaxira
    chunk = query[:500]
    suffix = "\n⚠️ Faqat 500 belgi tarjima qilindi (MyMemory chegarasi)" if len(query) > 500 else ""
    for src in ("uz", "ru", "en"):
        if src == lang:
            continue
        try:
            r = _httpx.get(
                "https://api.mymemory.translated.net/get",
                params={"q": chunk, "langpair": f"{src}|{lang}"},
                timeout=10,
            )
            d = r.json()
            txt = (d.get("responseData") or {}).get("translatedText", "")
            if d.get("responseStatus") == 200 and txt and txt.strip() != chunk.strip():
                return txt + suffix
        except Exception:
            continue
    raise RuntimeError("Tarjima amalga oshmadi. Keyinroq urinib ko'ring.")

@app.post("/api/translate")
@limiter.limit("20/minute")
async def translate(request: Request):
    body = await request.json()
    raw = body.get("text", "").strip()
    if not raw:
        return JSONResponse({"result": "❌ Matn kiriting.\n\nFormat:\nTarjima qilinadigan matn\nen  (til kodi — ixtiyoriy, default: en)"})

    lines = raw.split("\n")
    # If last line is a short token without spaces → treat as lang code
    last = lines[-1].strip()
    if len(lines) >= 2 and last and len(last) <= 7 and " " not in last:
        lang = _LANG_ALIASES.get(last.lower(), last.lower())
        query = "\n".join(lines[:-1]).strip()
    else:
        lang = "en"
        query = raw

    if not query:
        return JSONResponse({"result": "❌ Matn kiriting."})

    if len(query) > 4500:
        query = query[:4500]

    try:
        loop = asyncio.get_event_loop()
        result = await asyncio.wait_for(
            loop.run_in_executor(None, _do_translate, query, lang),
            timeout=20.0,
        )
        logger.info(f"translate: {len(query)} belgi → {lang}")
        return JSONResponse({"result": f"🌐 → {lang}\n\n{result}"})
    except asyncio.TimeoutError:
        return JSONResponse({"result": "❌ Tarjima vaqti tugadi (20s). Matnni qisqartiring."})
    except Exception as e:
        return JSONResponse({"result": f"❌ {e}"})

# ─── Wikipedia (proxy) ────────────────────────────────────────────────────────

# ✅ YANGILANDI: Wiki natijalari 1 soat cache (Redis/in-memory)
@_cache.cached(ttl=3600, key_prefix="wiki")
async def _wiki_lookup(q: str, pref_lang: str) -> str:
    langs = ([pref_lang] + [l for l in ["uz", "ru", "en"] if l != pref_lang]) if pref_lang else ["uz", "ru", "en"]
    async with httpx.AsyncClient(timeout=15) as client:
        for lang in langs:
            try:
                sr = await client.get(
                    f"https://{lang}.wikipedia.org/w/api.php",
                    params={"action": "opensearch", "search": q, "limit": 3,
                            "namespace": 0, "format": "json"},
                )
                titles = sr.json()[1] if sr.status_code == 200 else []
                search_title = titles[0] if titles else q
                encoded = url_quote(search_title, safe="")
                r = await client.get(
                    f"https://{lang}.wikipedia.org/api/rest_v1/page/summary/{encoded}"
                )
                if r.status_code == 200:
                    d = r.json()
                    if d.get("extract"):
                        url = d.get("content_urls", {}).get("desktop", {}).get("page", "")
                        result = f"📖 {d['title']} [{lang.upper()}]\n\n{d['extract']}"
                        if url:
                            result += f"\n\n🔗 {url}"
                        return result
            except Exception:
                continue
    return ""


@app.post("/api/wiki")
@limiter.limit("20/minute")
async def wiki(request: Request):
    body      = await request.json()
    q         = (body.get("text") or "").strip()[:200]
    pref_lang = (body.get("lang") or "").strip().lower()[:5]
    if not q:
        return JSONResponse({"result": "❌ Qidiruv so'zini kiriting"})
    result = await _wiki_lookup(q, pref_lang)
    if not result:
        return JSONResponse({"result": "❌ Wikipedia'da topilmadi. Boshqa so'z bilan qidiring."})
    return JSONResponse({"result": result})

# ─── Books (proxy) ────────────────────────────────────────────────────────────

# ✅ YANGILANDI: Books natijalari 6 soat cache
@_cache.cached(ttl=21600, key_prefix="books")
async def _books_lookup(q: str) -> str:
    try:
        async with httpx.AsyncClient(timeout=15) as client:
            gb_task = client.get(
                "https://www.googleapis.com/books/v1/volumes",
                params={"q": q, "maxResults": 6, "printType": "books"},
            )
            ol_task = client.get(
                "https://openlibrary.org/search.json",
                params={"q": q, "limit": 6,
                        "fields": "title,author_name,first_publish_year"},
            )
            gb_r, ol_r = await asyncio.gather(gb_task, ol_task, return_exceptions=True)

        lines = []
        seen = set()
        if not isinstance(gb_r, Exception) and gb_r.status_code == 200:
            for item in (gb_r.json().get("items") or [])[:6]:
                info = item.get("volumeInfo", {})
                title = info.get("title", "")
                if not title or title.lower() in seen:
                    continue
                seen.add(title.lower())
                line = f"📗 {title}"
                authors = info.get("authors", [])
                if authors:
                    line += f"\n   ✍️ {', '.join(authors[:2])}"
                year = (info.get("publishedDate") or "")[:4]
                if year:
                    line += f"  📅 {year}"
                preview = info.get("previewLink", "")
                if preview:
                    line += f"\n   🔗 {preview}"
                lines.append(line)
        if not isinstance(ol_r, Exception) and ol_r.status_code == 200:
            for b in (ol_r.json().get("docs") or [])[:6]:
                title = b.get("title", "")
                if not title or title.lower() in seen:
                    continue
                seen.add(title.lower())
                line = f"📘 {title}"
                if b.get("author_name"):
                    line += f"\n   ✍️ {b['author_name'][0]}"
                if b.get("first_publish_year"):
                    line += f"  📅 {b['first_publish_year']}"
                lines.append(line)
        if not lines:
            return ""
        return f"📚 Natijalar:\n\n" + "\n\n".join(lines[:8])
    except Exception:
        return ""


@app.post("/api/books")
@limiter.limit("20/minute")
async def books(request: Request):
    body = await request.json()
    q = (body.get("text") or "").strip()[:200]
    if not q:
        return JSONResponse({"result": "❌ Kitob nomi yoki muallifni kiriting"})
    result = await _books_lookup(q)
    if not result:
        return JSONResponse({"result": "📚 Kitob topilmadi"})
    return JSONResponse({"result": result})


# ─── ZIP ──────────────────────────────────────────────────────────────────────

@app.post("/api/zip")
@limiter.limit("20/minute")
async def make_zip(
    request: Request,
    files: List[UploadFile] = File(...),
    password: Optional[str] = Form(None),
):
    _ZIP_MAX_FILES = 100
    try:
        if password is not None and len(password) > 128:
            raise HTTPException(status_code=400, detail="Parol 128 belgidan oshmasligi kerak")
        if len(files) > _ZIP_MAX_FILES:
            raise HTTPException(status_code=400,
                detail=f"Maksimal {_ZIP_MAX_FILES} fayl ZIP ga qo'shilishi mumkin")

        buf = io.BytesIO()
        total = 0
        count = 0
        entries: list = []

        for f in files:
            data = await f.read()
            total += len(data)
            if total > MAX_FILE_BYTES * 3:
                raise HTTPException(status_code=413, detail="Fayllar jami hajmi juda katta")
            safe_name = os.path.basename(f.filename or "file") or "file"
            entries.append((safe_name, data))
            count += 1

        if password:
            import pyzipper
            with pyzipper.AESZipFile(buf, "w",
                                     compression=pyzipper.ZIP_DEFLATED,
                                     encryption=pyzipper.WZ_AES) as zf:
                zf.setpassword(password.encode("utf-8"))
                for name, content in entries:
                    zf.writestr(name, content)
        else:
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED, compresslevel=6) as zf:
                for name, content in entries:
                    zf.writestr(name, content)

        info = f"{count} fayl → archive.zip" + (" (parollangan)" if password else "")
        return Response(
            content=buf.getvalue(),
            media_type="application/zip",
            headers={
                "Content-Disposition": "attachment; filename=archive.zip",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"zip xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"zip: {type(e).__name__}: {str(e)[:160]}")

# ─── Unzip ────────────────────────────────────────────────────────────────────

_UNZIP_MAX_FILES       = 500
_UNZIP_MAX_UNCOMP_MB   = 200
_UNZIP_MAX_RATIO       = 50   # compressed:uncompressed ratio limit


def _build_response_from_entries(entries: list) -> Response:
    """entries = [(name, content_bytes), ...]"""
    if len(entries) == 1:
        name, content = entries[0]
        safe = os.path.basename(name) or "file"
        return Response(
            content=content,
            media_type="application/octet-stream",
            headers={"Content-Disposition": f"attachment; filename={url_quote(safe)}"},
        )
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, content in entries:
            zout.writestr(name, content)
    return Response(
        content=out.getvalue(),
        media_type="application/zip",
        headers={
            "Content-Disposition": "attachment; filename=extracted.zip",
            "X-Info": safe_header(f"{len(entries)} fayl"),
        },
    )


@app.post("/api/unzip")
@limiter.limit("20/minute")
async def unzip_file(request: Request, file: UploadFile = File(...)):
    try:
        data = await file.read()
        check_size(data, "/api/unzip")
        fname = (file.filename or "").lower()

        # ── 7z ───────────────────────────────────────────────────────────
        is_7z = fname.endswith(".7z") or data[:6] == b"7z\xbc\xaf'\x1c"
        if is_7z:
            try:
                import py7zr
                with py7zr.SevenZipFile(io.BytesIO(data), mode="r") as sz:
                    names = sz.getnames()
                    if not names:
                        raise HTTPException(status_code=400, detail="7z ichida fayl yo'q")
                    if len(names) > _UNZIP_MAX_FILES:
                        raise HTTPException(status_code=400,
                            detail=f"7z ichida {len(names)} fayl — maksimal {_UNZIP_MAX_FILES}")
                    all_data_dict = sz.readall() or {}
                entries = [
                    (os.path.basename(n) or "file", buf.read())
                    for n, buf in all_data_dict.items()
                    if buf is not None
                ]
                return _build_response_from_entries(entries)
            except HTTPException:
                raise
            except Exception as e:
                raise HTTPException(status_code=422, detail=f"7z ochib bo'lmadi: {e}")

        # ── TAR (tar, tar.gz, tgz, tar.bz2, tar.xz) ─────────────────────
        import tarfile as _tarfile
        is_tar = (
            fname.endswith((".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tar.xz", ".txz"))
            or data[:5] == b"ustar"
            or data[:2] == b"\x1f\x8b"  # gzip magic (may be .tar.gz)
        )
        if is_tar:
            try:
                tf = _tarfile.open(fileobj=io.BytesIO(data))
                members = [m for m in tf.getmembers() if m.isfile()]
                if not members:
                    raise HTTPException(status_code=400, detail="TAR ichida fayl yo'q")
                if len(members) > _UNZIP_MAX_FILES:
                    raise HTTPException(status_code=400,
                        detail=f"TAR ichida {len(members)} fayl — maksimal {_UNZIP_MAX_FILES}")
                total_uncomp = sum(m.size for m in members)
                if total_uncomp > _UNZIP_MAX_UNCOMP_MB * 1024 * 1024:
                    raise HTTPException(status_code=400,
                        detail=f"Siqilmagan hajm {total_uncomp // (1024*1024)} MB — maksimal {_UNZIP_MAX_UNCOMP_MB} MB")
                entries = []
                for m in members:
                    f_obj = tf.extractfile(m)
                    if f_obj:
                        safe = os.path.basename(m.name) or "file"
                        entries.append((safe, f_obj.read()))
                tf.close()
                return _build_response_from_entries(entries)
            except HTTPException:
                raise
            except Exception:
                pass  # fallback to ZIP path below

        # ── ZIP ───────────────────────────────────────────────────────────
        if data[:2] != b'PK':
            raise HTTPException(status_code=422,
                detail="Bu fayl ZIP arxivi emas. .zip fayl yuklang.")
        try:
            zf = zipfile.ZipFile(io.BytesIO(data))
        except zipfile.BadZipFile:
            raise HTTPException(status_code=422,
                detail="ZIP fayl buzilgan yoki noto'g'ri format. Boshqa .zip fayl yuklang.")
        all_infos = [i for i in zf.infolist() if not i.filename.endswith("/")]
        if not all_infos:
            raise HTTPException(status_code=400, detail="ZIP ichida fayl yo'q")

        if len(all_infos) > _UNZIP_MAX_FILES:
            raise HTTPException(status_code=400,
                detail=f"ZIP ichida {len(all_infos)} fayl — maksimal {_UNZIP_MAX_FILES}")
        total_uncomp = sum(i.file_size for i in all_infos)
        if total_uncomp > _UNZIP_MAX_UNCOMP_MB * 1024 * 1024:
            raise HTTPException(status_code=400,
                detail=f"Siqilmagan hajm {total_uncomp // (1024*1024)} MB — maksimal {_UNZIP_MAX_UNCOMP_MB} MB")
        if len(data) > 0 and total_uncomp / len(data) > _UNZIP_MAX_RATIO:
            raise HTTPException(status_code=400, detail="ZIP bomb aniqlandi — fayl xavfsiz emas")

        entries = [(i.filename, zf.read(i.filename)) for i in all_infos]
        return _build_response_from_entries(entries)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"unzip xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"unzip: {type(e).__name__}: {str(e)[:160]}")


# ─── Payment: Create ─────────────────────────────────────────────────────────

@app.post("/api/payment/create")
@limiter.limit("10/minute")
async def payment_create(request: Request):
    """
    Body: { user_id: int, plan: "monthly" | "yearly" }
    Returns: { payment_id, checkout_url, amount }
    """
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="JSON format noto'g'ri")

    user_id = body.get("user_id")
    plan    = body.get("plan", "")

    if not user_id or plan not in pay.PLAN_PRICES:
        raise HTTPException(status_code=400, detail="user_id va plan (monthly/yearly) kerak")

    amount     = pay.PLAN_PRICES[plan]
    payment_id = pay.new_payment_id()

    try:
        await db.create_payment(payment_id, int(user_id), plan, amount)
    except Exception as e:
        logger.error(f"DB create_payment xatosi: {e}")
        raise HTTPException(status_code=500, detail="To'lov yaratishda xatolik")

    checkout_url = pay.build_checkout_url(payment_id, plan)
    logger.info(f"To'lov yaratildi: {payment_id} user={user_id} plan={plan}")
    return {"payment_id": payment_id, "checkout_url": checkout_url, "amount": amount}


# ─── Payment: Payme RPC callback ──────────────────────────────────────────────

@app.post("/api/payment/payme")
async def payment_payme(request: Request):
    """Payme JSON-RPC 2.0 callback endpoint."""
    authorization = request.headers.get("Authorization", "")
    if not pay.verify_payme_auth(authorization):
        return JSONResponse(
            status_code=401,
            content={"error": {"code": -32504, "message": {"uz": "Ruxsat yo'q", "ru": "Доступ запрещён", "en": "Unauthorized"}}},
        )

    try:
        body = await request.json()
    except Exception:
        return JSONResponse(
            content={"error": {"code": pay.PaymeError.PARSE_ERROR, "message": {"uz": "JSON xatosi", "ru": "Ошибка JSON", "en": "Parse error"}}}
        )

    rpc_id = body.get("id", 1)
    method = body.get("method", "")
    params = body.get("params", {})

    db_fns = {
        "get_payment":          db.get_payment,
        "get_payment_by_payme": db.get_payment_by_payme,
        "confirm_payment":      db.confirm_payment,
        "cancel_payment":       db.cancel_payment,
    }

    result = await pay.handle_rpc(method, params, db_fns)
    return JSONResponse({"jsonrpc": "2.0", "id": rpc_id, **result})


# ─── Payment: Status ──────────────────────────────────────────────────────────

@app.get("/api/payment/status/{payment_id}")
@limiter.limit("30/minute")
async def payment_status(request: Request, payment_id: str):
    """
    Returns: { status: "pending" | "paid" | "cancelled", plan, amount }
    """
    payment = await db.get_payment(payment_id)
    if not payment:
        raise HTTPException(status_code=404, detail="To'lov topilmadi")
    return {
        "status":   payment["status"],
        "plan":     payment["plan"],
        "amount":   payment["amount"],
        "paid_at":  payment.get("completed_at"),
    }


# ─── User: Plan info ──────────────────────────────────────────────────────────

@app.get("/api/user/{user_id}/plan")
@limiter.limit("30/minute")
async def user_plan(request: Request, user_id: int):
    """
    Returns: { plan, plan_until, is_premium, usage_count }
    """
    user = await db.get_user(user_id)
    if not user:
        # Yangi foydalanuvchi — free plan
        return {"plan": "free", "plan_until": None, "is_premium": False, "usage_count": 0}
    premium = await db.is_premium(user_id)
    return {
        "plan":        user["plan"],
        "plan_until":  user["plan_until"],
        "is_premium":  premium,
        "usage_count": user["usage_count"],
    }


# ─── History endpoints ─────────────────────────────────────────────────────────

def _get_user_id(request: Request) -> int:
    """Extract user_id from Telegram WebApp init data via X-User-Id header."""
    uid = request.headers.get("X-User-Id") or request.headers.get("x-user-id")
    if not uid:
        raise HTTPException(status_code=401, detail="X-User-Id header kerak")
    try:
        return int(uid)
    except ValueError:
        raise HTTPException(status_code=400, detail="X-User-Id noto'g'ri")


@app.get("/api/history")
@limiter.limit("30/minute")
async def history_list(request: Request, limit: int = 20):
    user_id = _get_user_id(request)
    items = await db.get_history(user_id, limit=min(limit, 50))
    return {"items": items}


@app.delete("/api/history")
@limiter.limit("20/minute")
async def history_clear(request: Request):
    user_id = _get_user_id(request)
    await db.delete_history(user_id)
    return {"ok": True}


@app.delete("/api/history/{history_id}")
@limiter.limit("60/minute")
async def history_delete(request: Request, history_id: int):
    user_id = _get_user_id(request)
    await db.delete_history(user_id, history_id)
    return {"ok": True}


# ─── Admin helpers ─────────────────────────────────────────────────────────────

ADMIN_TOKEN = os.environ.get("ADMIN_TOKEN", "")

def _check_admin(request: Request):
    if not ADMIN_TOKEN:
        raise HTTPException(status_code=503, detail="ADMIN_TOKEN sozlanmagan — Railway Variables ga qo'shing")
    auth = request.headers.get("Authorization", "")
    if not auth.startswith("Bearer ") or not secrets.compare_digest(auth[7:], ADMIN_TOKEN):
        raise HTTPException(status_code=403, detail="Ruxsat yo'q")

# ─── Admin: Panel HTML ─────────────────────────────────────────────────────────

@app.get("/admin", include_in_schema=False)
async def admin_panel():
    import os as _os
    html_path = _os.path.join(_os.path.dirname(__file__), "admin.html")
    if not _os.path.exists(html_path):
        raise HTTPException(status_code=404, detail="admin.html topilmadi")
    from fastapi.responses import HTMLResponse
    with open(html_path, encoding="utf-8") as f:
        return HTMLResponse(f.read())

# ─── Admin: Stats ──────────────────────────────────────────────────────────────

# Note: GET /api/admin/stats — admin panel uchun. POST /api/stats — public math statistika.
@app.get("/api/admin/stats")
async def admin_stats(request: Request):
    _check_admin(request)
    stats = await db.get_stats()
    stats["uptime_seconds"] = int(time.time() - _start_time)
    return stats

# ─── Admin: Users ──────────────────────────────────────────────────────────────

@app.get("/api/admin/users")
async def admin_users(
    request: Request,
    offset: int = 0,
    limit: int = 50,
    search: str = "",
):
    _check_admin(request)
    return await db.get_users(offset=offset, limit=min(limit, 200), search=search)

# ─── Admin: Payments ───────────────────────────────────────────────────────────

@app.get("/api/admin/payments")
async def admin_payments(
    request: Request,
    offset: int = 0,
    limit: int = 50,
):
    _check_admin(request)
    return await db.get_payments(offset=offset, limit=min(limit, 200))

# ─── Admin: Set plan ───────────────────────────────────────────────────────────

@app.post("/api/admin/user/{user_id}/plan")
async def admin_set_plan(
    request: Request,
    user_id: int,
):
    _check_admin(request)
    body = await request.json()
    plan = body.get("plan", "free")
    days = int(body.get("days", 30))
    if plan not in ("free", "monthly", "yearly"):
        raise HTTPException(status_code=400, detail="plan: free | monthly | yearly")
    await db.admin_set_plan(user_id, plan, days)
    return {"ok": True, "user_id": user_id, "plan": plan, "days": days}
