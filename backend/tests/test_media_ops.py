"""Tests for media endpoints with full header + file-open verification.

Covers: imgcompress, cv, xlsx2pdf, compresspptx.
ML-heavy endpoints (bgremove, ocr, photo3x4) are @pytest.mark.slow
with pytest.importorskip guards so they are skipped when deps are absent.
"""
import io

import pytest

from conftest import (
    AUTH_HEADERS,
    assert_valid_png,
    assert_valid_jpeg,
    assert_content_type,
    assert_content_disposition,
    assert_x_header,
)


# ─── /api/imgcompress ─────────────────────────────────────────────────────────

class TestImgCompress:
    def test_compress_jpeg_returns_image_with_headers(self, client, minimal_jpeg):
        r = client.post(
            "/api/imgcompress",
            files=[("file", ("photo.jpg", minimal_jpeg, "image/jpeg"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 200
        assert_content_type(r.headers, "image")
        assert_content_disposition(r.headers, "compressed.")
        saved_hdr = assert_x_header(r.headers, "x-saved-percent")
        assert saved_hdr.lstrip("-").isdigit(), f"Non-numeric x-saved-percent: {saved_hdr!r}"
        assert_x_header(r.headers, "x-info")
        # Returned content must be a valid image
        ct = r.headers.get("content-type", "")
        if "jpeg" in ct or "jpg" in ct:
            assert_valid_jpeg(r.content)
        elif "png" in ct:
            assert_valid_png(r.content)

    def test_compress_png_returns_image(self, client, minimal_png):
        r = client.post(
            "/api/imgcompress",
            files=[("file", ("image.png", minimal_png, "image/png"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 200
        assert_content_type(r.headers, "image")

    def test_compress_not_image_returns_error(self, client):
        r = client.post(
            "/api/imgcompress",
            files=[("file", ("doc.pdf", b"%PDF-1.4 fake content", "application/pdf"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code in (400, 415, 422, 500)
        assert r.status_code != 200

    def test_compress_too_large_returns_413(self, client, too_large_bytes):
        r = client.post(
            "/api/imgcompress",
            files=[("file", ("huge.jpg", too_large_bytes, "image/jpeg"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 413


# ─── /api/cv ──────────────────────────────────────────────────────────────────

class TestCv:
    _CV_PAYLOAD = {
        "name": "Ali Valiyev",
        "phone": "+998901234567",
        "email": "ali@example.com",
        "summary": "Experienced Python backend developer",
        "experience": [
            {
                "company": "EduBot", "role": "Backend Dev",
                "period": "2023–2024", "desc": "Built FastAPI services",
            }
        ],
        "education": [
            {"institution": "TATU", "degree": "B.Sc. Computer Science", "period": "2019–2023"}
        ],
        "skills": ["Python", "FastAPI", "PostgreSQL", "Docker"],
        "languages": ["Uzbek (native)", "English (B2)"],
    }

    def test_cv_returns_pdf_with_x_info(self, client):
        r = client.post("/api/cv", json=self._CV_PAYLOAD, headers=AUTH_HEADERS)
        assert r.status_code == 200
        assert_content_type(r.headers, "pdf")
        assert_x_header(r.headers, "x-info")
        assert r.content[:4] == b"%PDF"

    def test_cv_content_disposition_has_filename(self, client):
        r = client.post("/api/cv", json=self._CV_PAYLOAD, headers=AUTH_HEADERS)
        assert r.status_code == 200
        cd = r.headers.get("content-disposition", "")
        assert cd, "Content-Disposition missing"
        # Filename should contain the person's name or cv
        assert "Ali" in cd or "cv" in cd.lower() or "CV" in cd

    def test_cv_missing_name_returns_error(self, client):
        payload = dict(self._CV_PAYLOAD)
        payload.pop("name")
        r = client.post("/api/cv", json=payload, headers=AUTH_HEADERS)
        assert r.status_code in (400, 422, 500)
        assert r.status_code != 200


# ─── /api/xlsx2pdf ────────────────────────────────────────────────────────────

class TestXlsx2Pdf:
    def test_xlsx2pdf_not_xlsx_returns_error(self, client):
        r = client.post(
            "/api/xlsx2pdf",
            files=[("file", ("doc.txt", b"hello world", "text/plain"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code in (400, 415, 422, 500)

    def test_xlsx2pdf_too_large_returns_413(self, client, too_large_bytes):
        r = client.post(
            "/api/xlsx2pdf",
            files=[(
                "file",
                ("big.xlsx", too_large_bytes,
                 "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
            )],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 413

    @pytest.mark.integration
    def test_xlsx2pdf_valid_xlsx_returns_pdf(self, client):
        openpyxl = pytest.importorskip("openpyxl")
        import io as _io
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Name"
        ws["B1"] = "Score"
        for i, (name, score) in enumerate([("Ali", 95), ("Zulfiya", 88), ("Bobur", 91)], 2):
            ws[f"A{i}"] = name
            ws[f"B{i}"] = score
        buf = _io.BytesIO()
        wb.save(buf)

        r = client.post(
            "/api/xlsx2pdf",
            files=[(
                "file",
                ("sheet.xlsx", buf.getvalue(),
                 "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
            )],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 200
        assert_content_type(r.headers, "pdf")
        assert_x_header(r.headers, "x-info")
        assert r.content[:4] == b"%PDF"


# ─── /api/compresspptx ────────────────────────────────────────────────────────

class TestCompressPptx:
    def test_compresspptx_not_pptx_returns_error(self, client):
        r = client.post(
            "/api/compresspptx",
            files=[("file", ("doc.txt", b"not a pptx", "text/plain"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code in (400, 415, 422, 500)

    def test_compresspptx_too_large_returns_413(self, client, too_large_bytes):
        r = client.post(
            "/api/compresspptx",
            files=[("file", ("big.pptx", too_large_bytes, "application/vnd.ms-powerpoint"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 413

    @pytest.mark.integration
    def test_compresspptx_valid_pptx_returns_pptx_with_headers(self, client):
        pptx = pytest.importorskip("pptx")
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        buf = io.BytesIO()
        prs.save(buf)

        r = client.post(
            "/api/compresspptx",
            files=[(
                "file",
                ("deck.pptx", buf.getvalue(),
                 "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
            )],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 200
        assert_content_type(r.headers, "pptx")
        assert_x_header(r.headers, "x-info")
        assert_x_header(r.headers, "x-saved-percent")
        assert_content_disposition(r.headers, "compressed.pptx")


# ─── ML endpoints — @pytest.mark.slow ────────────────────────────────────────

@pytest.mark.slow
class TestOcr:
    def test_ocr_image_returns_text_field(self, client, minimal_png):
        pytest.importorskip("paddleocr")
        r = client.post(
            "/api/ocr",
            files=[("file", ("image.png", minimal_png, "image/png"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 200
        body = r.json()
        assert "text" in body

    def test_ocr_not_image_not_pdf_returns_error(self, client):
        pytest.importorskip("paddleocr")
        r = client.post(
            "/api/ocr",
            files=[("file", ("doc.txt", b"plain text file", "text/plain"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code in (400, 415, 422, 500)

    def test_ocr_too_large_returns_413(self, client, too_large_bytes):
        pytest.importorskip("paddleocr")
        r = client.post(
            "/api/ocr",
            files=[("file", ("big.png", too_large_bytes, "image/png"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 413


@pytest.mark.slow
class TestBgRemove:
    def test_bgremove_returns_png(self, client, minimal_png):
        pytest.importorskip("rembg")
        r = client.post(
            "/api/bgremove",
            files=[("file", ("photo.png", minimal_png, "image/png"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code == 200
        assert_content_type(r.headers, "image/png")
        assert_valid_png(r.content)

    def test_bgremove_not_image_returns_error(self, client):
        pytest.importorskip("rembg")
        r = client.post(
            "/api/bgremove",
            files=[("file", ("doc.pdf", b"%PDF-1.4", "application/pdf"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code in (400, 415, 422, 500)


@pytest.mark.slow
class TestPhoto3x4:
    def test_photo3x4_returns_image_or_no_face(self, client, minimal_jpeg):
        pytest.importorskip("mediapipe")
        r = client.post(
            "/api/photo3x4",
            files=[("file", ("photo.jpg", minimal_jpeg, "image/jpeg"))],
            headers=AUTH_HEADERS,
        )
        # 200 → image returned; 422 → no face detected (valid for 10×10 pixel image)
        assert r.status_code in (200, 422)
        if r.status_code == 200:
            assert_content_type(r.headers, "image")

    def test_photo3x4_not_image_returns_error(self, client):
        pytest.importorskip("mediapipe")
        r = client.post(
            "/api/photo3x4",
            files=[("file", ("doc.pdf", b"%PDF-1.4", "application/pdf"))],
            headers=AUTH_HEADERS,
        )
        assert r.status_code in (400, 415, 422, 500)
