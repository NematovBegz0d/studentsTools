import io
import re
import time
import asyncio
import functools
import threading
from typing import Optional

from fastapi import APIRouter, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import JSONResponse, Response
from loguru import logger

from shared import (
    limiter, _io_pool, _ml_pool,
    check_size, safe_header, is_image_bytes,
    read_upload,
    safe_url_fetcher,
    get_rembg_session,
    acquire_user_ml_slot, release_user_ml_slot, _get_user_id,
    FONT_REGULAR, FONT_BOLD,
)

router = APIRouter()

# ─── ReportLab font helper (used by xlsx2pdf + cv fallback) ───────────────────

def _rl_fonts():
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    try:
        pdfmetrics.registerFont(TTFont('DV', FONT_REGULAR))
        pdfmetrics.registerFont(TTFont('DV-Bold', FONT_BOLD))
        return 'DV', 'DV-Bold'
    except Exception:
        return 'Helvetica', 'Helvetica-Bold'


# ─── 3×4 Passport foto ────────────────────────────────────────────────────────

_mp_face_detector = None
_mp_face_lock = threading.Lock()
_mp_process_lock = threading.Lock()


def _get_mp_face_detector():
    global _mp_face_detector
    if _mp_face_detector is False:
        return None
    if _mp_face_detector is not None:
        return _mp_face_detector
    with _mp_face_lock:
        if _mp_face_detector is False:
            return None
        if _mp_face_detector is not None:
            return _mp_face_detector
        try:
            import mediapipe as mp
            mp_face = mp.solutions.face_detection
            _mp_face_detector = mp_face.FaceDetection(
                model_selection=1,
                min_detection_confidence=0.5,
            )
            logger.info("MediaPipe FaceDetection init muvaffaqiyatli")
            return _mp_face_detector
        except Exception as e:
            logger.debug(f"MediaPipe yo'q yoki xato: {e}")
            _mp_face_detector = False
            return None


def _detect_face_mediapipe(img_pil):
    import numpy as np
    detector = _get_mp_face_detector()
    if detector is None:
        return None
    try:
        rgb_arr = np.array(img_pil)
        with _mp_process_lock:
            result = detector.process(rgb_arr)
        if not result or not result.detections:
            return None
        best = max(result.detections, key=lambda d: d.score[0])
        bb   = best.location_data.relative_bounding_box
        iw, ih = img_pil.size
        fx = max(0, int(bb.xmin * iw))
        fy = max(0, int(bb.ymin * ih))
        fw = max(1, int(bb.width * iw))
        fh = max(1, int(bb.height * ih))
        keypoints = {}
        names = ("right_eye", "left_eye", "nose_tip", "mouth_center", "right_ear", "left_ear")
        for name, kp in zip(names, best.location_data.relative_keypoints):
            keypoints[name] = (float(kp.x) * iw, float(kp.y) * ih)
        return {
            "bbox": (fx, fy, fw, fh),
            "keypoints": keypoints,
            "score": float(best.score[0]) if best.score else 0.0,
            "detector": "MediaPipe",
        }
    except Exception as e:
        logger.debug(f"MediaPipe detect xato: {e}")
        return None


def _detect_face_haar(img_pil):
    import cv2
    import numpy as np
    try:
        img_cv  = cv2.cvtColor(np.array(img_pil), cv2.COLOR_RGB2GRAY)
        cascade = cv2.CascadeClassifier(
            cv2.data.haarcascades + "haarcascade_frontalface_default.xml"
        )
        faces = cascade.detectMultiScale(
            img_cv, scaleFactor=1.1, minNeighbors=5, minSize=(60, 60)
        )
        if len(faces) == 0:
            return None
        fx, fy, fw, fh = tuple(sorted(faces, key=lambda f: f[2] * f[3], reverse=True)[0])
        return {
            "bbox": (int(fx), int(fy), int(fw), int(fh)),
            "keypoints": {},
            "score": 0.0,
            "detector": "Haar",
        }
    except Exception as e:
        logger.debug(f"Haar detect xato: {e}")
        return None


def _scale_face_profile(face: Optional[dict], scale: float) -> Optional[dict]:
    if not face or scale == 1:
        return face
    fx, fy, fw, fh = face["bbox"]
    inv = 1 / scale
    return {
        **face,
        "bbox": (
            int(round(fx * inv)),
            int(round(fy * inv)),
            int(round(fw * inv)),
            int(round(fh * inv)),
        ),
        "keypoints": {
            name: (x * inv, y * inv)
            for name, (x, y) in (face.get("keypoints") or {}).items()
        },
    }


def _detect_best_face(img_pil):
    from PIL import Image

    iw, ih = img_pil.size
    max_side = max(iw, ih)
    scale = min(1.0, 1600 / max_side) if max_side else 1.0
    detect_img = img_pil
    if scale < 1:
        resample = Image.Resampling.LANCZOS if hasattr(Image, "Resampling") else Image.LANCZOS
        detect_img = img_pil.resize(
            (max(1, int(iw * scale)), max(1, int(ih * scale))),
            resample,
        )

    try:
        face = _detect_face_mediapipe(detect_img)
    except Exception as e:
        logger.debug(f"MediaPipe detect import/run xato: {e}")
        face = None
    if face is None:
        try:
            face = _detect_face_haar(detect_img)
        except Exception as e:
            logger.debug(f"Haar detect import/run xato: {e}")
            face = None
    return _scale_face_profile(face, scale)


def _photo3x4_crop_rect(iw: int, ih: int, face: Optional[dict], framing: str = "formal") -> tuple:
    target_ratio = 3 / 4
    if not face:
        if iw / max(1, ih) > target_ratio:
            ch = ih
            cw = int(ch * target_ratio)
        else:
            cw = iw
            ch = int(cw / target_ratio)
        x0 = (iw - cw) // 2
        y0 = max(0, int((ih - ch) * 0.38))
        return (x0, y0, x0 + cw, y0 + ch)

    fx, fy, fw, fh = face["bbox"]
    kps = face.get("keypoints") or {}
    eye_pts = [kps[name] for name in ("left_eye", "right_eye") if name in kps]
    eye_cx = sum(p[0] for p in eye_pts) / len(eye_pts) if eye_pts else fx + fw / 2
    eye_cy = sum(p[1] for p in eye_pts) / len(eye_pts) if eye_pts else fy + fh * 0.38

    if framing == "tight":
        face_target_h = 0.50
        eye_target_y = 0.35
        face_top_target_y = 0.10
    else:
        face_target_h = 0.42
        eye_target_y = 0.32
        face_top_target_y = 0.11

    crop_h = max(fh / face_target_h, fw / 0.46)
    crop_w = crop_h * target_ratio
    if crop_w < fw * 2.05:
        crop_w = fw * 2.05
        crop_h = crop_w / target_ratio

    cx = eye_cx
    y_from_eyes = eye_cy - crop_h * eye_target_y
    y_from_face = fy - crop_h * face_top_target_y
    y0 = int(round(y_from_eyes * 0.55 + y_from_face * 0.45))
    x0 = int(round(cx - crop_w / 2))
    x1 = int(round(x0 + crop_w))
    y1 = int(round(y0 + crop_h))
    return (x0, y0, x1, y1)


def _crop_with_padding(img, rect: tuple, fill: tuple):
    from PIL import Image

    x0, y0, x1, y1 = rect
    out_w = max(1, x1 - x0)
    out_h = max(1, y1 - y0)
    canvas = Image.new("RGB", (out_w, out_h), fill)
    ix0, iy0 = max(0, x0), max(0, y0)
    ix1, iy1 = min(img.width, x1), min(img.height, y1)
    if ix1 > ix0 and iy1 > iy0:
        canvas.paste(img.crop((ix0, iy0, ix1, iy1)), (ix0 - x0, iy0 - y0))
    return canvas


def _map_face_to_crop(face: Optional[dict], rect: tuple, out_size: tuple) -> Optional[dict]:
    if not face:
        return None
    x0, y0, x1, y1 = rect
    sx = out_size[0] / max(1, x1 - x0)
    sy = out_size[1] / max(1, y1 - y0)
    fx, fy, fw, fh = face["bbox"]
    return {
        **face,
        "bbox": (
            int(round((fx - x0) * sx)),
            int(round((fy - y0) * sy)),
            int(round(fw * sx)),
            int(round(fh * sy)),
        ),
    }


def _replace_photo_background(crop, bg_rgb: tuple) -> tuple:
    try:
        from PIL import Image, ImageFilter
        from rembg import remove as rembg_remove
        try:
            crop_no_bg = rembg_remove(
                crop,
                alpha_matting=True,
                alpha_matting_foreground_threshold=240,
                alpha_matting_background_threshold=10,
                alpha_matting_erode_size=10,
            )
        except TypeError:
            crop_no_bg = rembg_remove(crop)
        if crop_no_bg.mode != "RGBA":
            crop_no_bg = crop_no_bg.convert("RGBA")
        alpha = crop_no_bg.getchannel("A").filter(ImageFilter.GaussianBlur(0.45))
        bg_layer = Image.new("RGB", crop.size, bg_rgb)
        bg_layer.paste(crop_no_bg.convert("RGB"), mask=alpha)
        return bg_layer, True
    except Exception as e:
        logger.debug(f"rembg xato photo3x4: {e}")
        return crop, False


def _draw_formal_suit(photo, face_out: Optional[dict], attire: str = "suit"):
    if attire != "suit":
        return photo
    from PIL import ImageDraw

    img = photo.copy()
    draw = ImageDraw.Draw(img, "RGBA")
    W, H = img.size
    if face_out:
        fx, fy, fw, fh = face_out["bbox"]
        cx = fx + fw / 2
        chin_y = fy + fh
        shoulder_w = max(W * 0.74, fw * 2.45)
        neck_y = int(min(max(chin_y - H * 0.015, H * 0.50), H * 0.66))
    else:
        cx = W / 2
        shoulder_w = W * 0.78
        neck_y = int(H * 0.57)

    left_shoulder = int(max(0, cx - shoulder_w / 2))
    right_shoulder = int(min(W, cx + shoulder_w / 2))
    center = int(round(cx))
    bottom = H + 10

    jacket    = (25, 32, 45, 238)
    jacket_hi = (38, 50, 68, 232)
    shirt     = (250, 250, 248, 245)
    collar    = (242, 242, 240, 246)
    tie       = (38, 55, 104, 242)

    shoulder_y = min(H - 40, neck_y + int(H * 0.18))
    draw.polygon(
        [(left_shoulder, shoulder_y), (center - 22, neck_y + 18), (center - 58, bottom), (0, bottom), (0, H - 55)],
        fill=jacket,
    )
    draw.polygon(
        [(right_shoulder, shoulder_y), (center + 22, neck_y + 18), (center + 58, bottom), (W, bottom), (W, H - 55)],
        fill=jacket,
    )
    draw.polygon(
        [(center - 34, neck_y + 8), (center + 34, neck_y + 8), (center + 68, bottom), (center - 68, bottom)],
        fill=shirt,
    )
    draw.polygon([(center - 40, neck_y + 8), (center - 6, neck_y + 54), (center - 58, neck_y + 36)], fill=collar)
    draw.polygon([(center + 40, neck_y + 8), (center + 6, neck_y + 54), (center + 58, neck_y + 36)], fill=collar)
    draw.polygon([(center - 13, neck_y + 44), (center + 13, neck_y + 44), (center + 20, bottom), (center - 20, bottom)], fill=tie)
    draw.polygon([(center - 16, neck_y + 32), (center, neck_y + 48), (center + 16, neck_y + 32), (center, neck_y + 20)], fill=tie)
    draw.line([(center - 74, neck_y + 54), (center - 28, bottom)], fill=jacket_hi, width=2)
    draw.line([(center + 74, neck_y + 54), (center + 28, bottom)], fill=jacket_hi, width=2)
    return img


def _enhance_passport_photo(photo):
    from PIL import ImageEnhance, ImageOps

    photo = ImageOps.autocontrast(photo, cutoff=0.4)
    photo = ImageEnhance.Color(photo).enhance(1.03)
    photo = ImageEnhance.Contrast(photo).enhance(1.04)
    photo = ImageEnhance.Sharpness(photo).enhance(1.12)
    return photo


def _do_photo3x4(
    data: bytes,
    bg: str = "white",
    sheet: bool = True,
    attire: str = "suit",
    framing: str = "formal",
) -> tuple:
    from PIL import Image, ImageOps, ImageDraw

    W, H = 354, 472
    BG_COLORS = {
        "white":     (255, 255, 255),
        "blue":      (100, 149, 237),
        "lightblue": (222, 238, 255),
        "gray":      (220, 224, 230),
    }
    bg_rgb  = BG_COLORS.get(bg, BG_COLORS["white"])
    attire  = attire  if attire  in ("suit", "natural") else "suit"
    framing = framing if framing in ("formal", "tight") else "formal"

    img_pil = Image.open(io.BytesIO(data))
    img_pil = ImageOps.exif_transpose(img_pil)
    if img_pil.mode != "RGB":
        img_pil = img_pil.convert("RGB")

    iw, ih = img_pil.size
    face = _detect_best_face(img_pil)
    detector_used = face.get("detector", "face") if face else "markaz crop"
    rect  = _photo3x4_crop_rect(iw, ih, face, framing)
    crop  = _crop_with_padding(img_pil, rect, bg_rgb)
    crop, bg_replaced = _replace_photo_background(crop, bg_rgb)

    resample = Image.Resampling.LANCZOS if hasattr(Image, "Resampling") else Image.LANCZOS
    photo    = crop.resize((W, H), resample)
    face_out = _map_face_to_crop(face, rect, (W, H))
    photo    = _draw_formal_suit(photo, face_out, attire)
    photo    = _enhance_passport_photo(photo)

    bg_note   = " - fon almashtirildi" if bg_replaced else ""
    suit_note = " - kostyum" if attire == "suit" else ""

    if not sheet:
        buf = io.BytesIO()
        photo.save(buf, format="JPEG", quality=96, subsampling=0, optimize=True, dpi=(300, 300))
        return buf.getvalue(), f"3x4 sm - 300 DPI - {detector_used}{suit_note}{bg_note}"

    A4W, A4H = 2480, 3508
    COLS, ROWS = 3, 2
    GAP_X, GAP_Y = 70, 70
    grid_w  = COLS * W + (COLS - 1) * GAP_X
    start_x = (A4W - grid_w) // 2
    start_y = 230
    canvas  = Image.new("RGB", (A4W, A4H), (255, 255, 255))
    draw    = ImageDraw.Draw(canvas)
    for row in range(ROWS):
        for col in range(COLS):
            px = start_x + col * (W + GAP_X)
            py = start_y + row * (H + GAP_Y)
            canvas.paste(photo, (px, py))
            draw.rectangle([px - 1, py - 1, px + W, py + H], outline=(170, 170, 170), width=1)
            cut = 16
            draw.line([(px - cut, py), (px - 3, py)], fill=(210, 210, 210), width=1)
            draw.line([(px, py - cut), (px, py - 3)], fill=(210, 210, 210), width=1)
            draw.line([(px + W + 3, py), (px + W + cut, py)], fill=(210, 210, 210), width=1)
            draw.line([(px + W, py - cut), (px + W, py - 3)], fill=(210, 210, 210), width=1)
            draw.line([(px - cut, py + H), (px - 3, py + H)], fill=(210, 210, 210), width=1)
            draw.line([(px, py + H + 3), (px, py + H + cut)], fill=(210, 210, 210), width=1)
            draw.line([(px + W + 3, py + H), (px + W + cut, py + H)], fill=(210, 210, 210), width=1)
            draw.line([(px + W, py + H + 3), (px + W, py + H + cut)], fill=(210, 210, 210), width=1)

    buf = io.BytesIO()
    canvas.save(buf, format="JPEG", quality=96, subsampling=0, optimize=True, dpi=(300, 300))
    return buf.getvalue(), f"A4 - 6 ta 3x4 - 300 DPI - {detector_used}{suit_note}{bg_note}"


@router.post("/api/photo3x4")
@limiter.limit("15/minute")
async def photo3x4(
    request: Request,
    file: UploadFile = File(...),
    bg: str      = Form("white"),
    attire: str  = Form("suit"),
    framing: str = Form("formal"),
    sheet: str   = Form("true"),
):
    t0 = time.time()
    user_id = _get_user_id(request)
    if not await acquire_user_ml_slot(user_id):
        raise HTTPException(status_code=429, detail="Parallel ML so'rov rad etildi — oldingi so'rov tugashini kuting")
    try:
        data = await read_upload(file, "/api/photo3x4")

        if not is_image_bytes(data):
            raise HTTPException(status_code=422,
                detail="Rasm fayli emas. JPG, PNG, WebP yoki HEIC yuklang.")

        bg       = bg if bg in ("white", "blue", "lightblue", "gray") else "white"
        do_sheet = sheet.lower() not in ("false", "0", "no")
        attire   = attire  if attire  in ("suit", "natural") else "suit"
        framing  = framing if framing in ("formal", "tight") else "formal"

        loop = asyncio.get_running_loop()
        try:
            jpg_bytes, info = await asyncio.wait_for(
                loop.run_in_executor(
                    _ml_pool,
                    functools.partial(_do_photo3x4, data, bg, do_sheet, attire, framing),
                ),
                timeout=60.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408, detail="Konversiya 60 soniyadan oshdi.")

        fname = "photo3x4_sheet.jpg" if do_sheet else "photo3x4.jpg"
        logger.info(f"photo3x4: {len(data)//1024}KB → {len(jpg_bytes)//1024}KB [{bg} sheet={do_sheet}] {time.time()-t0:.1f}s")
        return Response(
            content=jpg_bytes,
            media_type="image/jpeg",
            headers={
                "Content-Disposition": f"attachment; filename={fname}",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"photo3x4 xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"photo3x4: {type(e).__name__}: {str(e)[:160]}")
    finally:
        await release_user_ml_slot(user_id)


# ─── CV (Resume) Generator ────────────────────────────────────────────────────

_CV_TEMPLATES = {
    "modern":  {"accent": (79, 70, 229),  "bg": (248, 247, 255), "text": (15, 15, 30)},
    "classic": {"accent": (30, 64, 175),  "bg": (255, 255, 255), "text": (17, 24, 39)},
    "minimal": {"accent": (5, 150, 105),  "bg": (255, 255, 255), "text": (17, 24, 39)},
    "dark":    {"accent": (167, 139, 250), "bg": (15, 15, 30),   "text": (240, 240, 255)},
}

_CV_HTML_THEMES = {
    "modern":  {"accent": "#4F46E5", "bg": "#F8F7FF", "text": "#0F0F1E", "header_bg": "#4F46E5", "header_fg": "#FFFFFF"},
    "classic": {"accent": "#1E40AF", "bg": "#FFFFFF", "text": "#111827", "header_bg": "#1E40AF", "header_fg": "#FFFFFF"},
    "minimal": {"accent": "#059669", "bg": "#FFFFFF", "text": "#111827", "header_bg": "#059669", "header_fg": "#FFFFFF"},
    "dark":    {"accent": "#A78BFA", "bg": "#0F0F1E", "text": "#F0F0FF", "header_bg": "#1A1A2E", "header_fg": "#F0F0FF"},
}

_CV_HTML_TMPL = """<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
  @page { size: A4; margin: 0; }
  * { box-sizing: border-box; margin: 0; padding: 0; }
  body { font-family: 'DejaVu Sans', sans-serif; font-size: 10pt;
         background: {{ bg }}; color: {{ text }}; }
  .header { background: {{ header_bg }}; color: {{ header_fg }}; padding: 22px 28px; }
  .header h1 { font-size: 22pt; font-weight: 700; line-height: 1.1; }
  .header .title { font-size: 12pt; opacity: .85; margin-top: 4px; }
  .contacts { display: flex; flex-wrap: wrap; gap: 14px; margin-top: 12px;
              font-size: 9pt; opacity: .92; }
  .body { padding: 18px 28px 24px; }
  .section { margin-bottom: 14px; page-break-inside: avoid; }
  .section h2 { font-size: 11pt; font-weight: 700; color: {{ accent }};
                border-bottom: 1.5px solid {{ accent }}; padding-bottom: 3px;
                margin-bottom: 8px; text-transform: uppercase; letter-spacing: .04em; }
  .skills { display: flex; flex-wrap: wrap; gap: 5px; }
  .skill { background: rgba(0,0,0,0.04); color: {{ accent }};
           border: 1px solid {{ accent }}40;
           border-radius: 4px; padding: 2px 8px; font-size: 9pt; }
  .entry { margin-bottom: 8px; page-break-inside: avoid; }
  .entry-head { display: flex; justify-content: space-between; font-weight: 600; }
  .entry-sub { font-size: 9pt; opacity: .75; margin: 1px 0; }
  .entry-desc { font-size: 9pt; margin-top: 2px; opacity: .85; }
  .summary { font-size: 10pt; line-height: 1.55; opacity: .92; }
</style></head><body>
<div class="header">
  <h1>{{ name }}</h1>
  {% if title %}<div class="title">{{ title }}</div>{% endif %}
  <div class="contacts">
    {% if email %}<span>Email: {{ email }}</span>{% endif %}
    {% if phone %}<span>Tel: {{ phone }}</span>{% endif %}
    {% if location %}<span>Manzil: {{ location }}</span>{% endif %}
  </div>
</div>
<div class="body">
  {% if summary %}
  <div class="section">
    <h2>Haqimda</h2>
    <div class="summary">{{ summary }}</div>
  </div>
  {% endif %}
  {% if skills %}
  <div class="section">
    <h2>Ko'nikmalar</h2>
    <div class="skills">
      {% for s in skills %}<span class="skill">{{ s }}</span>{% endfor %}
    </div>
  </div>
  {% endif %}
  {% if experience %}
  <div class="section">
    <h2>Ish tajribasi</h2>
    {% for e in experience %}
    <div class="entry">
      <div class="entry-head"><span>{{ e.position or "" }}</span><span>{{ e.period or "" }}</span></div>
      {% if e.company %}<div class="entry-sub">{{ e.company }}</div>{% endif %}
      {% if e.desc %}<div class="entry-desc">{{ e.desc }}</div>{% endif %}
    </div>
    {% endfor %}
  </div>
  {% endif %}
  {% if education %}
  <div class="section">
    <h2>Ta'lim</h2>
    {% for e in education %}
    <div class="entry">
      <div class="entry-head"><span>{{ e.degree or "" }}</span><span>{{ e.year or "" }}</span></div>
      {% if e.school %}<div class="entry-sub">{{ e.school }}</div>{% endif %}
    </div>
    {% endfor %}
  </div>
  {% endif %}
  {% if languages %}
  <div class="section">
    <h2>Tillar</h2>
    <div class="skills">
      {% for l in languages %}<span class="skill">{{ l }}</span>{% endfor %}
    </div>
  </div>
  {% endif %}
</div>
</body></html>"""


def _do_cv(
    name: str, title: str, email: str, phone: str, location: str,
    summary: str, skills: list, education: list, experience: list,
    languages: list, template: str = "modern",
) -> bytes:
    try:
        from jinja2 import Environment, BaseLoader
        from weasyprint import HTML

        theme  = _CV_HTML_THEMES.get(template, _CV_HTML_THEMES["modern"])
        env    = Environment(loader=BaseLoader(), autoescape=True)
        jtmpl  = env.from_string(_CV_HTML_TMPL)
        html_str = jtmpl.render(
            name=name, title=title, email=email, phone=phone,
            location=location, summary=summary, skills=skills,
            education=education, experience=experience, languages=languages,
            **theme,
        )
        # SSRF guard: block any external/file/http resources WeasyPrint might try
        # to fetch (CV template only needs system fonts — data:/no URLs at all).
        return HTML(string=html_str, url_fetcher=safe_url_fetcher).write_pdf()
    except ImportError as e:
        logger.warning(f"CV WeasyPrint/Jinja2 yo'q ({e}) — ReportLab fallback")
    except Exception as e:
        logger.warning(f"CV WeasyPrint xato: {type(e).__name__}: {e} — ReportLab fallback")

    return _do_cv_reportlab(
        name, title, email, phone, location, summary,
        skills, education, experience, languages, template,
    )


def _do_cv_reportlab(
    name: str, title: str, email: str, phone: str, location: str,
    summary: str, skills: list, education: list, experience: list,
    languages: list, template: str = "modern",
) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib.colors import Color
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle,
    )
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import os as _os

    tmpl    = _CV_TEMPLATES.get(template, _CV_TEMPLATES["modern"])
    accent  = Color(tmpl["accent"][0]/255, tmpl["accent"][1]/255, tmpl["accent"][2]/255)
    txt_col = Color(tmpl["text"][0]/255,   tmpl["text"][1]/255,   tmpl["text"][2]/255)

    _font_reg, _font_bold = "Helvetica", "Helvetica-Bold"
    try:
        from reportlab.pdfbase.pdfmetrics import registerFontFamily

        candidates = [
            (FONT_REGULAR, FONT_BOLD),
            ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
             "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        ]

        for reg_path, bold_path in candidates:
            if reg_path and bold_path and _os.path.exists(reg_path) and _os.path.exists(bold_path):
                pdfmetrics.registerFont(TTFont("DejaVuSans", reg_path))
                pdfmetrics.registerFont(TTFont("DejaVuSans-Bold", bold_path))
                registerFontFamily(
                    "DejaVuSans",
                    normal="DejaVuSans",
                    bold="DejaVuSans-Bold",
                    italic="DejaVuSans",
                    boldItalic="DejaVuSans-Bold",
                )
                _font_reg, _font_bold = "DejaVuSans", "DejaVuSans-Bold"
                break
    except Exception as e:
        logger.debug(f"CV ReportLab font fallback: {e}")
        _font_reg, _font_bold = "Helvetica", "Helvetica-Bold"

    W, H = A4
    ML, MR, MT, MB = 18*mm, 18*mm, 18*mm, 18*mm

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=ML, rightMargin=MR, topMargin=MT, bottomMargin=MB,
    )

    def _style(name_s, **kw):
        kw.setdefault("fontName", _font_reg)
        kw.setdefault("textColor", txt_col)
        return ParagraphStyle(name_s, **kw)

    s_name    = _style("cvname",  fontName=_font_bold, fontSize=26, leading=30, textColor=accent)
    s_title   = _style("cvtitle", fontName=_font_reg,  fontSize=13, leading=16, textColor=txt_col, spaceAfter=2)
    s_contact = _style("cvcont",  fontName=_font_reg,  fontSize=9,  leading=13, textColor=txt_col)
    s_section = _style("cvsec",   fontName=_font_bold, fontSize=11, leading=14, textColor=accent,
                        spaceBefore=10, spaceAfter=3)
    s_body    = _style("cvbody",  fontName=_font_reg,  fontSize=9.5, leading=14, textColor=txt_col)
    s_bold    = _style("cvbold",  fontName=_font_bold, fontSize=9.5, leading=14, textColor=txt_col)
    s_small   = _style("cvsmall", fontName=_font_reg,  fontSize=8.5, leading=12, textColor=txt_col)
    s_skill   = _style("cvskill", fontName=_font_reg,  fontSize=9,   leading=13, textColor=txt_col)

    def hr():
        return HRFlowable(width="100%", thickness=0.8, color=accent, spaceAfter=4, spaceBefore=2)

    story = []
    story.append(Paragraph(name or "Ism Familiya", s_name))
    if title:
        story.append(Paragraph(title, s_title))
    contacts = [x for x in [email, phone, location] if x]
    if contacts:
        story.append(Paragraph("  •  ".join(contacts), s_contact))
    story.append(Spacer(1, 4*mm))

    if summary:
        story.append(hr())
        story.append(Paragraph("Haqida / О себе", s_section))
        story.append(Paragraph(summary, s_body))
        story.append(Spacer(1, 2*mm))

    if experience:
        story.append(hr())
        story.append(Paragraph("Ish tajribasi / Опыт работы", s_section))
        for exp in experience:
            pos    = exp.get("position", "")
            company = exp.get("company", "")
            period  = exp.get("period", "")
            desc    = exp.get("desc", "")
            header  = f"<b>{pos}</b>" + (f" — {company}" if company else "")
            story.append(Paragraph(header, s_bold))
            if period:
                story.append(Paragraph(period, s_small))
            if desc:
                story.append(Paragraph(desc, s_body))
            story.append(Spacer(1, 2*mm))

    if education:
        story.append(hr())
        story.append(Paragraph("Ta'lim / Образование", s_section))
        for edu in education:
            degree = edu.get("degree", "")
            school = edu.get("school", "")
            year   = edu.get("year", "")
            header = f"<b>{degree}</b>" + (f" — {school}" if school else "")
            story.append(Paragraph(header, s_bold))
            if year:
                story.append(Paragraph(year, s_small))
            story.append(Spacer(1, 2*mm))

    if skills:
        story.append(hr())
        story.append(Paragraph("Ko'nikmalar / Навыки", s_section))
        rows = [skills[i:i+3] for i in range(0, len(skills), 3)]
        tdata = []
        for row in rows:
            while len(row) < 3:
                row.append("")
            tdata.append([Paragraph(f"• {s}", s_skill) for s in row])
        if tdata:
            col_w = (W - ML - MR) / 3
            tbl = Table(tdata, colWidths=[col_w]*3)
            tbl.setStyle(TableStyle([
                ("VALIGN",        (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING",   (0, 0), (-1, -1), 0),
                ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]))
            story.append(tbl)
        story.append(Spacer(1, 2*mm))

    if languages:
        story.append(hr())
        story.append(Paragraph("Tillar / Языки", s_section))
        story.append(Paragraph("  •  ".join(languages), s_body))

    doc.build(story)
    return buf.getvalue()


@router.post("/api/cv")
@limiter.limit("10/minute")
async def make_cv(request: Request):
    t0 = time.time()
    user_id = _get_user_id(request)
    try:
        body = await request.json()

        name     = (body.get("name") or "").strip()[:80]
        title    = (body.get("title") or "").strip()[:100]
        email    = (body.get("email") or "").strip()[:80]
        phone    = (body.get("phone") or "").strip()[:30]
        location = (body.get("location") or "").strip()[:80]
        summary  = (body.get("summary") or "").strip()[:800]
        template = (body.get("template") or "modern").strip().lower()
        if template not in _CV_TEMPLATES:
            template = "modern"

        raw_skills = body.get("skills") or []
        if isinstance(raw_skills, str):
            raw_skills = [s.strip() for s in raw_skills.split(",") if s.strip()]
        skills = [str(s).strip()[:40] for s in raw_skills[:20]]

        raw_langs = body.get("languages") or []
        if isinstance(raw_langs, str):
            raw_langs = [s.strip() for s in raw_langs.split(",") if s.strip()]
        languages = [str(s).strip()[:40] for s in raw_langs[:10]]

        education = []
        for edu in (body.get("education") or [])[:8]:
            if isinstance(edu, dict):
                education.append({
                    "degree": str(edu.get("degree") or "")[:80],
                    "school": str(edu.get("school") or "")[:80],
                    "year":   str(edu.get("year")   or "")[:20],
                })

        experience = []
        for exp in (body.get("experience") or [])[:8]:
            if isinstance(exp, dict):
                experience.append({
                    "position": str(exp.get("position") or "")[:80],
                    "company":  str(exp.get("company")  or "")[:80],
                    "period":   str(exp.get("period")   or "")[:40],
                    "desc":     str(exp.get("desc")     or "")[:400],
                })

        if not name:
            raise HTTPException(status_code=400, detail="Ism (name) majburiy")

        loop = asyncio.get_running_loop()
        pdf_bytes = await asyncio.wait_for(
            loop.run_in_executor(
                _io_pool,
                functools.partial(
                    _do_cv, name, title, email, phone, location,
                    summary, skills, education, experience, languages, template,
                ),
            ),
            timeout=30.0,
        )
        fname = f"cv_{name.replace(' ', '_')[:30]}.pdf"
        logger.info(f"cv: {name!r} {template} {len(pdf_bytes)//1024}KB {time.time()-t0:.1f}s")
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{fname}"',
                "X-Info": safe_header(f"{template} · {len(experience)} tajriba · {len(education)} ta'lim"),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"cv xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"cv: {type(e).__name__}: {str(e)[:160]}")


# ─── Excel → PDF ──────────────────────────────────────────────────────────────

def _do_xlsx2pdf(data: bytes) -> tuple:
    import subprocess
    import tempfile
    import os
    import shutil

    lo_path = shutil.which("libreoffice") or shutil.which("soffice")
    if lo_path:
        with tempfile.TemporaryDirectory() as tmpdir:
            xlsx_path = os.path.join(tmpdir, "input.xlsx")
            pdf_path  = os.path.join(tmpdir, "input.pdf")
            try:
                with open(xlsx_path, "wb") as f:
                    f.write(data)
                result = subprocess.run(
                    [lo_path, "--headless", "--nologo", "--nofirststartwizard",
                     "--convert-to", "pdf", "--outdir", tmpdir, xlsx_path],
                    capture_output=True, timeout=60,
                )
                if result.returncode == 0 and os.path.exists(pdf_path):
                    with open(pdf_path, "rb") as f:
                        pdf_bytes = f.read()
                    return pdf_bytes, "✅ LibreOffice · formatlar saqlandi"
                logger.warning(
                    f"xlsx2pdf LibreOffice rc={result.returncode}: "
                    f"{result.stderr[:200].decode('utf-8', errors='ignore')}"
                )
            except subprocess.TimeoutExpired:
                logger.warning("xlsx2pdf LibreOffice timeout > 60s — ReportLab fallback")
            except Exception as e:
                logger.warning(f"xlsx2pdf LibreOffice xato: {e} — ReportLab fallback")

    return _do_xlsx2pdf_reportlab(data)


def _do_xlsx2pdf_reportlab(data: bytes) -> tuple:
    from openpyxl import load_workbook
    from reportlab.platypus import (
        SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak,
    )
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from datetime import datetime as _dt, date as _date, time as _time

    fn, fn_bold = _rl_fonts()
    MAX_ROWS    = 500
    MAX_SHEETS  = 10
    MAX_COLS    = 20
    MARGIN      = 1.2 * cm

    def fmt_val(v) -> str:
        if v is None:
            return ''
        if isinstance(v, _dt):
            return v.strftime('%d.%m.%Y %H:%M') if v.hour or v.minute else v.strftime('%d.%m.%Y')
        if isinstance(v, _date):
            return v.strftime('%d.%m.%Y')
        if isinstance(v, _time):
            return v.strftime('%H:%M')
        if isinstance(v, bool):
            return 'Ha' if v else "Yo'q"
        if isinstance(v, float):
            if abs(v) < 1e14 and v == int(v):
                return str(int(v))
            return f'{v:g}'
        return str(v)

    def esc(t: str) -> str:
        return t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

    def calc_col_widths(all_rows, col_n: int, avail_w: float) -> list:
        max_lens = [1] * col_n
        for row in all_rows[:100]:
            for j in range(col_n):
                v = row[j] if j < len(row) else None
                max_lens[j] = max(max_lens[j], min(len(fmt_val(v)), 40))
        total = sum(max_lens) or col_n
        MIN_W, MAX_W = 1.0 * cm, 9.0 * cm
        widths = [max(MIN_W, min(MAX_W, avail_w * l / total)) for l in max_lens]
        scale = avail_w / sum(widths)
        return [w * scale for w in widths]

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    max_col_n = 0
    for name in wb.sheetnames[:MAX_SHEETS]:
        ws = wb[name]
        for row in ws.iter_rows(values_only=True, max_row=5, max_col=MAX_COLS):
            non_empty = sum(1 for v in row if v is not None)
            max_col_n = max(max_col_n, non_empty)
    wb.close()

    if max_col_n > 7:
        page_w, page_h = A4[1], A4[0]
        orientation = "Landscape"
    else:
        page_w, page_h = A4
        orientation = "Portrait"
    avail_w = page_w - 2 * MARGIN

    wb = load_workbook(io.BytesIO(data), read_only=False, data_only=True)

    def _cell_bg(cell):
        try:
            fill = cell.fill
            if fill and fill.fill_type == "solid":
                rgb = fill.fgColor.rgb
                if rgb and rgb not in ("00000000", "FF000000", "FFFFFFFF"):
                    return f"#{rgb[2:]}"
        except Exception:
            pass
        return None

    title_st = ParagraphStyle('xt', fontName=fn_bold, fontSize=12,
                               spaceAfter=6, spaceBefore=4,
                               textColor=colors.HexColor('#1a3a5c'))
    note_st  = ParagraphStyle('xn', fontName=fn, fontSize=8,
                               textColor=colors.HexColor('#888888'), spaceAfter=4)

    story        = []
    total_rows   = 0
    total_sheets = 0
    any_truncated = False

    for name in wb.sheetnames[:MAX_SHEETS]:
        ws = wb[name]
        raw_cells = list(ws.iter_rows(values_only=False,
                                      max_row=MAX_ROWS + 1, max_col=MAX_COLS))
        if not raw_cells:
            continue

        truncated = len(raw_cells) > MAX_ROWS
        raw_cells = raw_cells[:MAX_ROWS]
        if truncated:
            any_truncated = True

        raw = [[cell.value for cell in row] for row in raw_cells]

        col_n = MAX_COLS
        while col_n > 1:
            if all((r[col_n - 1] if col_n - 1 < len(r) else None) is None
                   for r in raw[:30]):
                col_n -= 1
            else:
                break
        col_n = max(1, col_n)

        while raw and all(v is None for v in raw[-1][:col_n]):
            raw.pop()
            raw_cells.pop()
        if not raw:
            continue

        if   col_n <= 4:  fs = 10
        elif col_n <= 7:  fs = 9
        elif col_n <= 11: fs = 8
        else:             fs = 7

        tc_st  = ParagraphStyle(f'tc_{name}',  fontName=fn,      fontSize=fs, leading=fs + 3)
        tch_st = ParagraphStyle(f'tch_{name}', fontName=fn_bold, fontSize=fs, leading=fs + 3)

        col_widths = calc_col_widths(raw, col_n, avail_w)

        table_data = []
        cell_bgs   = []
        for r_idx, (row_cells, row_vals) in enumerate(zip(raw_cells, raw)):
            is_hdr = r_idx == 0
            cells  = []
            for j in range(col_n):
                cell_obj = row_cells[j] if j < len(row_cells) else None
                v        = row_vals[j]  if j < len(row_vals)  else None
                txt      = esc(fmt_val(v))
                is_bold  = is_hdr or bool(
                    cell_obj and cell_obj.font and cell_obj.font.bold
                )
                cells.append(Paragraph(txt, tch_st if is_bold else tc_st))
                if cell_obj and not is_hdr:
                    bg = _cell_bg(cell_obj)
                    if bg:
                        cell_bgs.append((r_idx, j, bg))
            table_data.append(cells)

        if story:
            story.append(PageBreak())
        story.append(Paragraph(name, title_st))
        if truncated:
            story.append(Paragraph(
                f"⚠️ Faqat birinchi {MAX_ROWS} qator ko'rsatildi", note_st))

        tbl = Table(table_data, colWidths=col_widths, repeatRows=1)
        style_cmds = [
            ('BACKGROUND',    (0, 0), (-1,  0), colors.HexColor('#cfe2f3')),
            ('LINEBELOW',     (0, 0), (-1,  0), 1.2, colors.HexColor('#2e7cbf')),
            ('ROWBACKGROUNDS',(0, 1), (-1, -1),
             [colors.white, colors.HexColor('#f0f6fb')]),
            ('GRID',          (0, 0), (-1, -1), 0.35, colors.HexColor('#b0c8e0')),
            ('TOPPADDING',    (0, 0), (-1, -1), 3),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
            ('LEFTPADDING',   (0, 0), (-1, -1), 4),
            ('RIGHTPADDING',  (0, 0), (-1, -1), 4),
            ('VALIGN',        (0, 0), (-1, -1), 'TOP'),
            ('WORDWRAP',      (0, 0), (-1, -1), 'WORD'),
        ]
        for r_idx, c_idx, hex_color in cell_bgs:
            style_cmds.append(
                ('BACKGROUND', (c_idx, r_idx), (c_idx, r_idx), colors.HexColor(hex_color))
            )
        tbl.setStyle(TableStyle(style_cmds))
        story.append(tbl)
        total_rows   += len(raw)
        total_sheets += 1

    wb.close()

    if not story:
        raise ValueError("Excel fayl bo'sh yoki o'qib bo'lmadi.")

    out = io.BytesIO()
    doc = SimpleDocTemplate(
        out, pagesize=(page_w, page_h),
        leftMargin=MARGIN, rightMargin=MARGIN,
        topMargin=MARGIN + 0.3 * cm, bottomMargin=MARGIN,
        title="EduBot — Excel to PDF",
    )
    doc.build(story)

    sheets_word = f"{total_sheets} varaq" if total_sheets > 1 else "1 varaq"
    info = f"✅ {sheets_word} · {total_rows} qator · {orientation}"
    if any_truncated:
        info += f" · (max {MAX_ROWS} qator/varaq)"
    return out.getvalue(), info


@router.post("/api/xlsx2pdf")
@limiter.limit("10/minute")
async def xlsx_to_pdf(request: Request, file: UploadFile = File(...)):
    t0 = time.time()
    user_id = _get_user_id(request)
    try:
        data = await read_upload(file, "/api/xlsx2pdf")
        if len(data) < 4 or data[:2] != b'PK':
            raise HTTPException(status_code=422,
                detail="Bu fayl Excel (.xlsx) emas. .xlsx yoki .xls fayl yuklang.")

        loop = asyncio.get_running_loop()
        try:
            pdf_bytes, info = await asyncio.wait_for(
                loop.run_in_executor(_io_pool, _do_xlsx2pdf, data),
                timeout=60.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Konversiya 60 soniyadan oshdi. Excel fayl juda katta.")

        logger.info(f"xlsx2pdf: {len(data)//1024}KB → {len(pdf_bytes)//1024}KB, {time.time()-t0:.1f}s")
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=spreadsheet.pdf",
                "X-Info": safe_header(info),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"xlsx2pdf xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"xlsx2pdf: {type(e).__name__}: {str(e)[:160]}")


# ─── PPTX Compress ────────────────────────────────────────────────────────────

_PPTX_MAX_DIM = 1920


def _do_compresspptx(data: bytes) -> tuple:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from PIL import Image

    orig_size  = len(data)
    compressed = 0

    prs = Presentation(io.BytesIO(data))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                img_blob = shape.image.blob
                img_ext  = shape.image.ext.lower()
                with Image.open(io.BytesIO(img_blob)) as img:
                    if max(img.width, img.height) > _PPTX_MAX_DIM:
                        ratio   = _PPTX_MAX_DIM / max(img.width, img.height)
                        resized = img.resize(
                            (int(img.width * ratio), int(img.height * ratio)),
                            Image.LANCZOS,
                        )
                    else:
                        resized = img
                    if len(img_blob) > 2 * 1024 * 1024:
                        quality = 65
                    elif len(img_blob) > 512 * 1024:
                        quality = 72
                    else:
                        quality = 82
                    buf = io.BytesIO()
                    if resized.mode == "RGBA" or img_ext in ("png", "gif"):
                        resized.save(buf, format="PNG", optimize=True)
                    else:
                        resized.convert("RGB").save(buf, format="JPEG", quality=quality, optimize=True)
                    candidate = buf.getvalue()
                    if resized is not img:
                        resized.close()
                if len(candidate) < len(img_blob):
                    blip  = shape.element.blipFill.blip
                    rId   = blip.get(qn("r:embed"))
                    image_part = shape.part.related_parts[rId]
                    image_part._blob = candidate
                    compressed += 1
            except Exception as e:
                logger.debug(f"pptx image compress skip: {e}")
                continue

    out_buf   = io.BytesIO()
    prs.save(out_buf)
    out_bytes = out_buf.getvalue()
    if len(out_bytes) >= orig_size:
        out_bytes = data
    saved = max(0, round((1 - len(out_bytes) / orig_size) * 100))
    info  = f"{compressed} ta rasm siqildi, {saved}% kichiklashdi"
    return out_bytes, info, saved


@router.post("/api/compresspptx")
@limiter.limit("10/minute")
async def compress_pptx(request: Request, file: UploadFile = File(...)):
    user_id = _get_user_id(request)
    try:
        data = await read_upload(file, "/api/compresspptx")
        loop = asyncio.get_running_loop()
        try:
            out_bytes, info, saved = await asyncio.wait_for(
                loop.run_in_executor(_io_pool, _do_compresspptx, data),
                timeout=45.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="PPTX compression timed out")
        logger.info(f"compresspptx: {info}")
        return Response(
            content=out_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": "attachment; filename=compressed.pptx",
                "X-Info": safe_header(info),
                "X-Saved-Percent": str(saved),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"compresspptx xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"compresspptx: {type(e).__name__}: {str(e)[:160]}")


# ─── Image compress ───────────────────────────────────────────────────────────

_IMGCOMPRESS_MAX_DIM = 4000


def _do_imgcompress(data: bytes, output_format: str = "jpeg", user_quality: int = 0) -> tuple:
    orig_size = len(data)

    try:
        import pyvips
        vimg = pyvips.Image.new_from_buffer(data, "")
        w, h = vimg.width, vimg.height
        if max(w, h) > _IMGCOMPRESS_MAX_DIM:
            scale = _IMGCOMPRESS_MAX_DIM / max(w, h)
            vimg  = vimg.resize(scale)
        if user_quality:
            quality = max(30, min(95, user_quality))
        elif orig_size > 3 * 1024 * 1024:
            quality = 65
        elif orig_size > 1024 * 1024:
            quality = 72
        else:
            quality = 82
        has_alpha = vimg.hasalpha()
        fmt = output_format
        if has_alpha and fmt == "jpeg":
            fmt = "webp"
        if fmt == "webp":
            out_bytes  = vimg.webpsave_buffer(Q=quality)
            media_type, ext = "image/webp", "webp"
        elif fmt == "png":
            out_bytes  = vimg.pngsave_buffer()
            media_type, ext = "image/png", "png"
        else:
            if has_alpha:
                vimg = vimg.flatten()
            out_bytes  = vimg.jpegsave_buffer(Q=quality, optimize_coding=True)
            media_type, ext = "image/jpeg", "jpg"
        if len(out_bytes) >= orig_size:
            out_bytes = data
            quality   = 0
            if data[:3] == b"\xff\xd8\xff":
                media_type, ext = "image/jpeg", "jpg"
            elif data[:8] == b"\x89PNG\r\n\x1a\n":
                media_type, ext = "image/png", "png"
            elif data[:4] == b"RIFF" and data[8:12] == b"WEBP":
                media_type, ext = "image/webp", "webp"
            else:
                media_type, ext = "image/jpeg", "jpg"
        saved = max(0, round((1 - len(out_bytes) / orig_size) * 100))
        info  = f"{w}x{h} -> {vimg.width}x{vimg.height}, {fmt}, saved {saved}%"
        return out_bytes, info, saved, media_type, ext
    except ImportError:
        pass
    except Exception as _vips_err:
        logger.debug(f"pyvips xato, Pillow fallback: {_vips_err}")

    from PIL import Image, ImageOps

    img      = Image.open(io.BytesIO(data))
    _orig_fmt = (img.format or "jpeg").lower()
    img      = ImageOps.exif_transpose(img)

    w, h  = img.width, img.height
    ratio = min(1.0, _IMGCOMPRESS_MAX_DIM / max(w, h))
    if ratio < 1.0:
        img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)

    has_alpha = img.mode in ("RGBA", "LA", "P")
    if has_alpha and output_format == "jpeg":
        output_format = "webp"

    if user_quality:
        quality = max(30, min(95, user_quality))
    elif orig_size > 3 * 1024 * 1024:
        quality = 65
    elif orig_size > 1024 * 1024:
        quality = 72
    else:
        quality = 82

    out = io.BytesIO()
    if output_format == "webp":
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGBA" if has_alpha else "RGB")
        img.save(out, format="WEBP", quality=quality, method=4)
        media_type, ext = "image/webp", "webp"
    elif output_format == "png":
        img.save(out, format="PNG", optimize=True, compress_level=9)
        media_type, ext = "image/png", "png"
    else:
        if img.mode != "RGB":
            img = img.convert("RGB")
        img.save(out, format="JPEG", quality=quality, optimize=True)
        media_type, ext = "image/jpeg", "jpg"

    out_bytes = out.getvalue()
    if len(out_bytes) >= orig_size:
        out_bytes = data
        quality   = 0
        _fmt_map  = {"jpeg": ("image/jpeg", "jpg"), "png": ("image/png", "png"),
                     "webp": ("image/webp", "webp"), "gif": ("image/gif", "gif")}
        media_type, ext = _fmt_map.get(_orig_fmt, ("image/jpeg", "jpg"))

    saved = max(0, round((1 - len(out_bytes) / orig_size) * 100))
    info  = f"{w}×{h} → {img.width}×{img.height}, {output_format}, saved {saved}%"
    return out_bytes, info, saved, media_type, ext


@router.post("/api/imgcompress")
@limiter.limit("20/minute")
async def img_compress(
    request: Request,
    file: UploadFile = File(...),
    output_format: str = Form("jpeg"),
    quality: int = Form(0),
):
    user_id = _get_user_id(request)
    try:
        data = await read_upload(file, "/api/imgcompress")
        if len(data) < 4:
            raise HTTPException(status_code=422,
                detail="Rasm fayli ochib bo'lmadi. JPG, PNG yoki WebP yuklang.")
        sig = data[:4]
        is_image = (
            sig[:2] == b'\xff\xd8' or
            sig      == b'\x89PNG' or
            sig      == b'RIFF'    or
            sig[:3]  == b'GIF'
        )
        if not is_image:
            raise HTTPException(status_code=422,
                detail="Rasm fayli emas. JPG, PNG yoki WebP yuklang.")
        output_format = output_format if output_format in ("jpeg", "png", "webp") else "jpeg"
        loop = asyncio.get_running_loop()
        try:
            out_bytes, info, saved, media_type, ext = await asyncio.wait_for(
                loop.run_in_executor(
                    _io_pool,
                    functools.partial(_do_imgcompress, data, output_format, quality),
                ),
                timeout=30.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=504, detail="Image compression timed out")
        return Response(
            content=out_bytes,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename=compressed.{ext}",
                "X-Info": safe_header(info),
                "X-Saved-Percent": str(saved),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"imgcompress xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"imgcompress: {type(e).__name__}: {str(e)[:160]}")


# ─── Background removal ───────────────────────────────────────────────────────

@router.post("/api/bgremove")
@limiter.limit("5/minute")
async def bgremove(request: Request):
    # Auth FIRST — heavy `from rembg import remove` below would raise
    # ModuleNotFoundError on dev machines without rembg, bypassing auth.
    user_id = _get_user_id(request)
    import base64
    from rembg import remove
    t0 = time.time()
    if not await acquire_user_ml_slot(user_id):
        raise HTTPException(status_code=429, detail="Parallel ML so'rov rad etildi — oldingi so'rov tugashini kuting")
    try:
        ct       = request.headers.get("content-type", "")
        bg_color = None
        if "multipart" in ct:
            form     = await request.form()
            data     = await read_upload(form.get("file"), "/api/bgremove")
            bg_color = (form.get("bg_color") or "").strip().lstrip("#") or None
        else:
            body     = await request.json()
            data     = base64.b64decode(body["data"])
            if not data:
                raise HTTPException(status_code=400, detail="Fayl bo'sh. Boshqa fayl tanlang.")
            check_size(data, "/api/bgremove")
            bg_color = (body.get("bg_color") or "").strip().lstrip("#") or None

        session = get_rembg_session()
        loop    = asyncio.get_running_loop()
        try:
            result_png = await asyncio.wait_for(
                loop.run_in_executor(
                    _ml_pool,
                    functools.partial(remove, data, session=session),
                ),
                timeout=60.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="Fon olib tashlash 60 soniyadan oshdi. Kichikroq rasm yuklang.")

        _bg = bg_color.lstrip("#") if bg_color else ""
        if len(_bg) == 3:
            _bg = "".join(c * 2 for c in _bg)
        if _bg and len(_bg) >= 6:
            from PIL import Image as _PILImg
            try:
                r = int(_bg[0:2], 16)
                g = int(_bg[2:4], 16)
                b = int(_bg[4:6], 16)
                fg  = _PILImg.open(io.BytesIO(result_png)).convert("RGBA")
                bg  = _PILImg.new("RGBA", fg.size, (r, g, b, 255))
                bg.paste(fg, mask=fg.split()[3])
                out_buf = io.BytesIO()
                bg.convert("RGB").save(out_buf, format="JPEG", quality=95, optimize=True)
                result_png = out_buf.getvalue()
                media_type, fname = "image/jpeg", "no-bg.jpg"
            except Exception:
                media_type, fname = "image/png", "no-bg.png"
        else:
            media_type, fname = "image/png", "no-bg.png"

        logger.info(f"bgremove: {len(data)//1024}KB → {len(result_png)//1024}KB, {time.time()-t0:.1f}s")
        return Response(
            content=result_png,
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={fname}"},
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"bgremove xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"bgremove: {type(e).__name__}: {str(e)[:160]}")
    finally:
        await release_user_ml_slot(user_id)


# ─── OCR ──────────────────────────────────────────────────────────────────────

def _preprocess_for_ocr(img):
    from PIL import Image, ImageFilter, ImageEnhance, ImageOps

    img  = ImageOps.exif_transpose(img)
    gray = img.convert('L')
    w, h = gray.size

    MIN_EDGE = 1800
    if max(w, h) < MIN_EDGE:
        scale = MIN_EDGE / max(w, h)
        gray  = gray.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        w, h  = gray.size

    MAX_EDGE = 4000
    if max(w, h) > MAX_EDGE:
        scale = MAX_EDGE / max(w, h)
        gray  = gray.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

    gray = ImageOps.autocontrast(gray, cutoff=2)
    gray = ImageEnhance.Contrast(gray).enhance(1.4)
    gray = gray.filter(ImageFilter.SHARPEN)
    return gray


def _clean_ocr_text(text: str) -> str:
    text   = re.sub(r'\n{3,}', '\n\n', text)
    lines  = text.split('\n')
    result = []
    for line in lines:
        s = line.strip()
        if s and len(s) >= 2:
            result.append(s)
        elif not s and result and result[-1]:
            result.append('')
    while result and not result[0]:
        result.pop(0)
    while result and not result[-1]:
        result.pop()
    return '\n'.join(result)


_paddle_ocr_cache: dict = {}
_paddle_ocr_lock = threading.Lock()


def _get_paddle_ocr(lang: str = "en"):
    if lang in _paddle_ocr_cache:
        return _paddle_ocr_cache[lang]
    with _paddle_ocr_lock:
        if lang in _paddle_ocr_cache:
            return _paddle_ocr_cache[lang]
        try:
            from paddleocr import PaddleOCR
            _paddle_ocr_cache[lang] = PaddleOCR(use_angle_cls=True, lang=lang, show_log=False)
            logger.info(f"PaddleOCR init muvaffaqiyatli: lang={lang}")
        except Exception as e:
            logger.debug(f"paddleocr {lang} sessiya yaratilmadi: {e}")
            _paddle_ocr_cache[lang] = None
        return _paddle_ocr_cache[lang]


def _paddle_run(ocr_obj, arr) -> str:
    if ocr_obj is None:
        return ""
    result = ocr_obj.ocr(arr, cls=True)
    lines  = []
    for block in (result or []):
        for line in (block or []):
            if line and len(line) >= 2 and line[1]:
                lines.append(line[1][0])
    return "\n".join(lines)


def _do_ocr_paddle(img_bytes: bytes) -> str:
    import numpy as np
    from PIL import Image as _PIL
    en_ocr  = _get_paddle_ocr("en")
    cyr_ocr = _get_paddle_ocr("cyrillic")
    if en_ocr is None and cyr_ocr is None:
        raise ImportError("paddleocr not installed")
    img     = _PIL.open(io.BytesIO(img_bytes)).convert("RGB")
    arr     = np.array(img)
    en_text  = _paddle_run(en_ocr,  arr)
    cyr_text = _paddle_run(cyr_ocr, arr)
    if re.search(r'[а-яёўқғҳА-ЯЁ]', cyr_text):
        return cyr_text if len(cyr_text) > len(en_text) * 0.5 else en_text
    return en_text if len(en_text) >= len(cyr_text) else cyr_text


def _do_ocr(data: bytes, is_pdf: bool) -> str:
    import pytesseract
    from PIL import Image

    LANGS  = 'rus+eng+uzb'
    CONFIG = '--oem 3 --psm 6'

    if is_pdf:
        import fitz
        doc = fitz.open(stream=data, filetype='pdf')
        try:
            if doc.is_encrypted:
                return 'PDF parol bilan himoyalangan. Avval parolini oching.'
            total_pages = doc.page_count
            if total_pages == 0:
                return "PDF bo'sh (0 sahifa)."

            sample = ''.join(doc[i].get_text() for i in range(min(3, total_pages)))
            if len(sample.strip()) > 50:
                MAX_P = 15
                parts = []
                for i in range(min(total_pages, MAX_P)):
                    t = doc[i].get_text().strip()
                    if t:
                        if total_pages > 1:
                            parts.append(f'── Sahifa {i + 1} ──')
                        parts.append(t)
                result = _clean_ocr_text('\n\n'.join(parts))
                if total_pages > MAX_P:
                    result += f"\n\n⚠️ Faqat {MAX_P} sahifa o'qildi (PDF jami {total_pages} sahifa)."
                return result or 'Matn topilmadi.'

            zoom  = 200 / 72
            mat   = fitz.Matrix(zoom, zoom)
            MAX_P = 8
            parts = []
            for i in range(min(total_pages, MAX_P)):
                pix       = doc[i].get_pixmap(matrix=mat, alpha=False)
                img       = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
                del pix
                img_bytes = io.BytesIO()
                img.save(img_bytes, format="PNG")
                img_bytes = img_bytes.getvalue()
                try:
                    text = _do_ocr_paddle(img_bytes)
                except Exception:
                    img  = _preprocess_for_ocr(img)
                    text = pytesseract.image_to_string(img, lang=LANGS, config=CONFIG)
                cleaned = _clean_ocr_text(text)
                if cleaned:
                    if total_pages > 1:
                        parts.append(f'── Sahifa {i + 1} ──')
                    parts.append(cleaned)
            result = '\n\n'.join(parts)
            if total_pages > MAX_P:
                result += (f'\n\n⚠️ Faqat {MAX_P} sahifa OCR qilindi '
                           f'(PDF jami {total_pages} sahifali).')
            return result or 'Matn topilmadi.'
        finally:
            doc.close()

    try:
        text = _do_ocr_paddle(data)
        if text.strip():
            return _clean_ocr_text(text) or 'Matn topilmadi.'
    except Exception:
        pass
    img  = Image.open(io.BytesIO(data))
    img  = _preprocess_for_ocr(img)
    text = pytesseract.image_to_string(img, lang=LANGS, config=CONFIG)
    return _clean_ocr_text(text) or 'Matn topilmadi.'


@router.post("/api/ocr")
@limiter.limit("15/minute")
async def ocr(request: Request):
    t0 = time.time()
    user_id = _get_user_id(request)
    if not await acquire_user_ml_slot(user_id):
        raise HTTPException(status_code=429, detail="Parallel ML so'rov rad etildi — oldingi so'rov tugashini kuting")
    try:
        import base64
        ct = request.headers.get("content-type", "")
        if "multipart" in ct:
            form  = await request.form()
            f     = form.get("file")
            fname = (f.filename or "").lower() if f is not None else ""
            data  = await read_upload(f, "/api/ocr")
        else:
            body  = await request.json()
            data  = base64.b64decode(body["data"])
            if not data:
                raise HTTPException(status_code=400, detail="Fayl bo'sh. Boshqa fayl tanlang.")
            check_size(data, "/api/ocr")
            fname = body.get("filename", "").lower()

        is_pdf = fname.endswith('.pdf') or data[:4] == b'%PDF'

        loop = asyncio.get_running_loop()
        try:
            text = await asyncio.wait_for(
                loop.run_in_executor(_ml_pool, _do_ocr, data, is_pdf),
                timeout=90.0,
            )
        except asyncio.TimeoutError:
            raise HTTPException(status_code=408,
                detail="OCR 90 soniyadan oshdi. Rasmni kichiklashtiring yoki PDF'ni qisqartiring.")

        logger.info(f"ocr: {len(data)//1024}KB {'pdf' if is_pdf else 'img'} → {len(text)} belgi, {time.time()-t0:.1f}s")
        return JSONResponse({"text": text})

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"ocr xato: {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=f"ocr: {type(e).__name__}: {str(e)[:160]}")
    finally:
        await release_user_ml_slot(user_id)
